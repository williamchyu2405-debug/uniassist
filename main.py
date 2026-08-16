import os
import re
import json
import random
import sqlite3
import hashlib
import secrets
from collections import defaultdict
from datetime import datetime, date, timedelta
from pathlib import Path
from typing import Optional
import anthropic
from fastapi import FastAPI, UploadFile, File, HTTPException, Form, Request, Header, Depends
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, HTMLResponse, StreamingResponse, Response
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import pdfplumber
from pptx import Presentation
import base64

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass

# Bulletproof .env loader — handles broken/old dotenv packages
def _load_env_file():
    import pathlib
    env_path = pathlib.Path(__file__).parent / ".env"
    if not env_path.exists():
        return
    try:
        with open(env_path) as f:
            for line in f:
                line = line.strip()
                if not line or line.startswith('#') or '=' not in line:
                    continue
                if line.startswith('.') or line.startswith('/'):
                    continue  # Skip lines like './start.sh'
                k, _, v = line.partition('=')
                k = k.strip()
                v = v.strip().strip('"').strip("'")
                if k and k not in os.environ:
                    os.environ[k] = v
    except Exception:
        pass

_load_env_file()


def _normalize_topic(topic: str) -> str:
    """Shorten verbose AI-generated topics to 2-3 word broad categories.
    E.g. 'Acidity and pKa of Carboxylic Acids' → 'Carboxylic Acids'"""
    if not topic or len(topic) <= 25:
        return topic  # already short
    # Common patterns: "X of Y" → keep Y; "X and Y of Z" → keep Z
    # Strategy: take the last noun phrase (after last 'of'/'in'/'for') or first 2-3 words
    for sep in [' of ', ' in ', ' for ', ' during ']:
        if sep in topic:
            tail = topic.split(sep)[-1].strip()
            if len(tail) <= 30:
                return tail
    # Fallback: first 3 words
    words = topic.split()
    # Drop leading filler words
    fillers = {'the', 'a', 'an', 'and', 'or'}
    cleaned = [w for w in words if w.lower() not in fillers]
    return ' '.join(cleaned[:3])


# Optional passcode that must be entered before the app loads.
# Set ACCESS_CODE=yourcode in .env — anyone opening the URL must enter it first.
ACCESS_CODE = os.getenv("ACCESS_CODE", "").strip()

app = FastAPI(title="MedVault")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])


from fastapi.responses import JSONResponse

@app.middleware("http")
async def access_guard(request: Request, call_next):
    """Block all API routes unless the correct access code is provided.
    Accepts the code as either the X-Access-Code header (app) or the
    ?ac= query param (bookmarklet — keeps the request 'simple' so no CORS preflight is needed)."""
    if not ACCESS_CODE:
        return await call_next(request)
    path  = request.url.path
    method = request.method
    # Always pass through: HTML root, static assets, CORS preflight, and the info endpoint itself
    if (path == "/" or path == "/guides" or path == "/bookmarklet" or path.startswith("/static") or path.startswith("/images")
            or method == "OPTIONS" or path == "/api/access-check" or path == "/api/beacon"):
        return await call_next(request)
    provided = (request.headers.get("X-Access-Code", "")
                or request.query_params.get("ac", ""))
    if provided != ACCESS_CODE:
        return JSONResponse({"detail": "access_code_required"}, status_code=403)
    return await call_next(request)


@app.get("/api/access-check")
def access_check_get():
    """Tells the client whether an access code is required (no auth needed to call this)."""
    return {"required": bool(ACCESS_CODE)}


@app.post("/api/access-check")
async def access_check_post(request: Request):
    """Validate an access code. This endpoint is exempt from the middleware,
    so we check the code manually here."""
    if not ACCESS_CODE:
        return {"ok": True}
    provided = (request.headers.get("X-Access-Code", "")
                or request.query_params.get("ac", ""))
    if provided != ACCESS_CODE:
        return JSONResponse({"detail": "Incorrect access code"}, status_code=403)
    return {"ok": True}

MODEL = "claude-sonnet-4-6"
HAIKU = "claude-haiku-4-5-20251001"
WRITE_MODEL = "claude-opus-4-8"   # GAMSAT writing module — grading / stimulus / drills (see writing_* helpers)

# #13 Structured outputs on the flashcard generator (opt-in, default OFF).
# When "1", generate_flashcards uses output_config.format json_schema so Haiku
# returns schema-valid JSON directly; any failure falls back to the legacy
# generate_json + parse_json_response path. Deployed behavior is unchanged
# unless the env var is set.
STRUCTURED_FLASHCARDS = os.getenv("STRUCTURED_FLASHCARDS", "0") == "1"
GEN_CONTENT_CHARS = 50000   # ~12k tokens; fits a full multi-lesson SCORM module clip (cached, so cheap on repeat generators)

# Adaptive quiz bank — built up over time so studying stays fresh without re-spending API.
QUIZ_MIN_UNSEEN = 8    # while ≥ this many UNSEEN questions remain, serve from the bank with NO AI call
QUIZ_BANK_CAP   = 60   # max saved questions per material; oldest *seen* ones pruned beyond this
QUIZ_SESSION    = 60   # questions served per sitting (unseen first); = bank cap so one sitting can serve a guide's whole bank

# Performance-driven difficulty: a 4-rung ladder the app climbs/descends per topic
# based purely on the student's answer history (no AI call needed).
DIFF_RANK = {"easy": 0, "medium": 1, "hard": 2, "daredevil": 3}
RANK_NAME = {0: "easy", 1: "medium", 2: "hard", 3: "daredevil"}

# Fast re-adaptation tuning: short recency-weighted window so the ladder
# responds within ~2 attempts, with a streak fast-track up and a firm drop.
ADAPT_WINDOW      = 5      # attempts per topic considered (was 8)
ADAPT_MIN_HISTORY = 2      # attempts before we leave the medium default (was 3)
ADAPT_WEIGHTS     = [1.0, 0.8, 0.65, 0.5, 0.4]   # most-recent attempt first

def adaptive_targets(db, user_id: int, material_id: Optional[int] = None) -> dict:
    """Per-topic target difficulty rank derived from recent answers.
    Recency-weighted accuracy over the last ADAPT_WINDOW attempts so new form
    registers fast: acing → one rung harder (+2 on a hot streak of 3 correct at
    or above the current level); struggling → firm drop (−1, or −2 when the two
    most-recent answers are both wrong). Topics with fewer than
    ADAPT_MIN_HISTORY attempts default to medium (rank 1). Clamped to [0, 3].
    Pure logic over quiz_attempts — costs nothing."""
    where, args = "WHERE a.user_id = ?", [user_id]
    if material_id:
        where += " AND a.material_id = ?"
        args.append(material_id)
    rows = db.execute(
        f"""SELECT a.topic, a.is_correct, q.difficulty
            FROM quiz_attempts a JOIN quiz_questions q ON a.question_id = q.id
            {where}
            ORDER BY a.attempted_at DESC""",
        tuple(args)
    ).fetchall()
    by_topic: dict = {}
    for r in rows:
        recent = by_topic.setdefault(r["topic"] or "General", [])
        if len(recent) < ADAPT_WINDOW:  # most-recent first (rows are DESC)
            recent.append((1 if r["is_correct"] else 0,
                           DIFF_RANK.get(r["difficulty"] or "medium", 1)))
    targets = {}
    for topic, recent in by_topic.items():
        if len(recent) < ADAPT_MIN_HISTORY:
            targets[topic] = 1
            continue
        wts   = ADAPT_WEIGHTS[:len(recent)]
        wsum  = sum(wts)
        acc      = sum(w * c  for (c, _), w in zip(recent, wts)) / wsum   # recency-weighted
        cur_rank = round(sum(w * rk for (_, rk), w in zip(recent, wts)) / wsum)

        # Hot streak: last 3 answers all correct at/above the current level → jump 2.
        last3 = recent[:3]
        hot_streak = len(last3) == 3 and all(c and rk >= cur_rank for c, rk in last3)
        # Cold streak: two most-recent both wrong → drop firmly.
        cold_streak = len(recent) >= 2 and not recent[0][0] and not recent[1][0]

        if hot_streak:
            target = cur_rank + 2                   # dominating → fast-track harder
        elif acc >= 0.8:
            target = cur_rank + 1                   # mastering → harder
        elif cold_streak or acc < 0.3:
            target = cur_rank - 2                   # clearly struggling → firm drop
        elif acc < 0.5:
            target = cur_rank - 1                   # struggling → easier
        else:
            target = cur_rank                       # hold steady
        targets[topic] = max(0, min(3, target))
    return targets


# ── Concept graph — turn the (previously unused) AI `related_topics` into links
#    between concepts, and let weakness spread along those links. All zero-AI. ──

def concept_links(db, user_id: int) -> dict:
    """Undirected concept adjacency built from the `related_topics` the AI already
    writes on every question/card but that nothing ever read back.
    Returns {topic: {neighbour: weight}} with both endpoints normalised; weight is
    how many questions/cards asserted that link (a stronger, repeated connection)."""
    adj: dict = {}
    def link(a, b):
        a, b = _normalize_topic(a or ""), _normalize_topic(b or "")
        if not a or not b or a == b:
            return
        adj.setdefault(a, {})[b] = adj.setdefault(a, {}).get(b, 0) + 1
        adj.setdefault(b, {})[a] = adj.setdefault(b, {}).get(a, 0) + 1
    for tbl in ("quiz_questions", "flashcards"):
        for r in db.execute(
            f"SELECT topic, related_topics FROM {tbl} WHERE user_id = ? AND related_topics IS NOT NULL",
            (user_id,)
        ).fetchall():
            try:
                rel = json.loads(r["related_topics"] or "[]")
            except Exception:
                rel = []
            for t in rel:
                if isinstance(t, str):
                    link(r["topic"], t)
    return adj

def topic_accuracy(db, user_id: int) -> dict:
    """{normalised_topic: accuracy 0..1} over attempted topics.
    Normalisation collisions (two raw topics mapping to the same key) are merged
    by aggregating raw attempt/correct counts and computing accuracy once —
    order-independent and weighted by sample size (mirrors the combined_topics
    aggregation in /api/progress)."""
    agg: dict = {}
    for r in db.execute(
        "SELECT topic, COUNT(*) AS n, SUM(is_correct) AS c FROM quiz_attempts WHERE user_id = ? GROUP BY topic",
        (user_id,)
    ).fetchall():
        t = _normalize_topic(r["topic"] or "")
        if not t or not r["n"]:
            continue
        d = agg.setdefault(t, {"attempts": 0, "correct": 0})
        d["attempts"] += r["n"]
        d["correct"]  += int(r["c"] or 0)
    return {t: d["correct"] / d["attempts"] for t, d in agg.items() if d["attempts"]}

def propagated_weakness(db, user_id: int, lam: float = 0.5) -> dict:
    """Per-topic weakness that SPREADS along concept links: a topic adjacent to ones
    you're failing becomes 'at risk' before you've even slipped on it.
    score = own_weakness + lam * (link-weighted mean of neighbours' weakness).
    Higher = needs work more. Pure logic over attempts + related_topics."""
    acc = topic_accuracy(db, user_id)
    adj = concept_links(db, user_id)
    topics = set(acc) | set(adj)
    own = {t: (1 - acc[t]) if t in acc else 0.4 for t in topics}  # unattempted ⇒ mild unknown
    score = {}
    for t in topics:
        nbrs = adj.get(t, {})
        if nbrs:
            wsum = sum(nbrs.values())
            neigh = sum(own.get(n, 0.4) * w for n, w in nbrs.items()) / wsum
            score[t] = own[t] + lam * neigh
        else:
            score[t] = own[t]
    return score


WEAK_MIN_ATTEMPTS = 3

def compute_weak_topics(topic_agg: dict, min_attempts: int = WEAK_MIN_ATTEMPTS, limit: int = 10) -> list:
    """Reliability-aware weak-topic ranking (pure logic, unit-testable).
    Laplace-smoothed accuracy (correct+1)/(attempts+2) pulls tiny samples toward
    50% so one bad answer can't dominate; a topic needs `min_attempts` attempts
    to be ranked at all. Sorted weakest first, ties broken by larger sample.
    `topic_agg` is {topic: {"attempts": int, "correct": int}}."""
    out = []
    for t, d in topic_agg.items():
        att, corr = d["attempts"], d["correct"]
        if att < min_attempts:
            continue
        out.append({
            "topic": t,
            "attempts": att,
            "correct": corr,
            "accuracy": corr / att,
            "smoothed_accuracy": round((corr + 1) / (att + 2), 4),
        })
    out.sort(key=lambda x: (x["smoothed_accuracy"], -x["attempts"]))
    return out[:limit]


def sm2_schedule(ease, interval, count, correct):
    """One SM-2 step with a binary grade (correct=4, wrong=1). Returns
    (interval_days, ease_factor, review_count, next_review_iso). Shared by flashcards
    and quiz questions so spaced-repetition scheduling stays identical across both."""
    ease     = float(ease or 2.5)
    interval = int(interval or 1)
    count    = int(count or 0)
    quality  = 4 if correct else 1
    if quality >= 3:  # correct → grow the interval
        new_interval = 1 if count == 0 else 6 if count == 1 else max(1, round(interval * ease))
        new_ease     = max(1.3, ease + 0.1 - (5 - quality) * (0.08 + (5 - quality) * 0.02))
        new_count    = count + 1
    else:             # wrong → reset to tomorrow, small ease penalty
        new_interval, new_ease, new_count = 1, max(1.3, ease - 0.2), 0
    next_review = (date.today() + timedelta(days=new_interval)).isoformat()
    return new_interval, round(new_ease, 4), new_count, next_review


# ── Storage paths — override with env vars for cloud deployment ───────────────
# Locally (no DATA_DIR):  data/study.db, uploads/, static/material_images/
# Railway (DATA_DIR=/data): /data/data/study.db, /data/uploads/, /data/material_images/
_DATA_DIR_ENV = os.getenv("DATA_DIR", "").strip()
if _DATA_DIR_ENV:
    _DATA_ROOT = Path(_DATA_DIR_ENV)
    IMAGES_DIR = _DATA_ROOT / "material_images"   # served via /images/ route
else:
    _DATA_ROOT = Path(".")
    IMAGES_DIR = Path("static/material_images")   # served by StaticFiles + /images/ route

DB_PATH    = str(_DATA_ROOT / "data" / "study.db")
UPLOAD_DIR = _DATA_ROOT / "uploads"

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
IMAGES_DIR.mkdir(parents=True, exist_ok=True)
(_DATA_ROOT / "data").mkdir(parents=True, exist_ok=True)
Path("static").mkdir(exist_ok=True)


def get_client():
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not set in .env file")
    return anthropic.Anthropic(api_key=api_key)


def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


# ── Password hashing (stdlib PBKDF2 — no extra dependencies) ────────────────

def hash_password(password: str) -> str:
    """Hash a password with a random salt. Returns 'salt$hash' hex string."""
    salt = secrets.token_hex(16)
    h = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 260000)
    return salt + "$" + h.hex()

def verify_password(password: str, stored: str) -> bool:
    """Verify a password against a 'salt$hash' string."""
    if not stored or "$" not in stored:
        return False
    salt, _, expected = stored.partition("$")
    h = hashlib.pbkdf2_hmac("sha256", password.encode(), salt.encode(), 260000)
    return secrets.compare_digest(h.hex(), expected)


def init_db():
    conn = get_db()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS user_materials (
            user_id INTEGER NOT NULL,
            material_id INTEGER NOT NULL,
            added_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            PRIMARY KEY (user_id, material_id)
        );
        CREATE TABLE IF NOT EXISTS materials (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            visibility TEXT DEFAULT 'private',
            filename TEXT NOT NULL,
            original_name TEXT NOT NULL,
            subject TEXT DEFAULT 'Medicine',
            content TEXT,
            file_type TEXT,
            images TEXT DEFAULT '[]',
            uploaded_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS flashcards (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            material_id INTEGER,
            user_id INTEGER,
            topic TEXT,
            question TEXT NOT NULL,
            answer TEXT NOT NULL,
            times_seen INTEGER DEFAULT 0,
            times_correct INTEGER DEFAULT 0,
            last_seen DATETIME,
            next_review TEXT DEFAULT NULL,
            srs_interval INTEGER DEFAULT 1,
            ease_factor REAL DEFAULT 2.5,
            review_count INTEGER DEFAULT 0,
            FOREIGN KEY (material_id) REFERENCES materials(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS flashcard_log (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            flashcard_id INTEGER,
            material_id INTEGER,
            user_id INTEGER,
            topic TEXT,
            correct INTEGER DEFAULT 0,
            reviewed_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (flashcard_id) REFERENCES flashcards(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS quiz_questions (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            material_id INTEGER,
            user_id INTEGER,
            topic TEXT,
            question TEXT NOT NULL,
            options TEXT,
            correct_answer TEXT NOT NULL,
            explanation TEXT,
            difficulty TEXT DEFAULT 'medium',
            FOREIGN KEY (material_id) REFERENCES materials(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS quiz_attempts (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            question_id INTEGER,
            material_id INTEGER,
            user_id INTEGER,
            topic TEXT,
            user_answer TEXT,
            is_correct INTEGER,
            attempted_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS revision_slides (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            material_id INTEGER,
            user_id INTEGER,
            title TEXT,
            content TEXT,
            slide_order INTEGER DEFAULT 0,
            FOREIGN KEY (material_id) REFERENCES materials(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS exam_dates (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            subject TEXT,
            exam_date DATE,
            notes TEXT,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS mind_maps (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            material_id INTEGER,
            user_id INTEGER,
            title TEXT,
            data TEXT,
            FOREIGN KEY (material_id) REFERENCES materials(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS chat_messages (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            session_id TEXT,
            role TEXT,
            content TEXT,
            timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
        );
    """)
    conn.commit()
    # ── Backward-compat migrations for existing databases ─────────────────
    # (These are no-ops on new databases where CREATE TABLE already includes these columns)
    for stmt in [
        "ALTER TABLE materials ADD COLUMN images TEXT DEFAULT '[]'",
        "ALTER TABLE quiz_questions ADD COLUMN difficulty TEXT DEFAULT 'medium'",
        "ALTER TABLE flashcards ADD COLUMN next_review TEXT DEFAULT NULL",
        "ALTER TABLE flashcards ADD COLUMN srs_interval INTEGER DEFAULT 1",
        "ALTER TABLE flashcards ADD COLUMN ease_factor REAL DEFAULT 2.5",
        "ALTER TABLE flashcards ADD COLUMN review_count INTEGER DEFAULT 0",
        "ALTER TABLE quiz_attempts ADD COLUMN material_id INTEGER",
        # Multi-user columns
        "ALTER TABLE materials ADD COLUMN user_id INTEGER",
        "ALTER TABLE materials ADD COLUMN visibility TEXT DEFAULT 'private'",
        "ALTER TABLE flashcards ADD COLUMN user_id INTEGER",
        "ALTER TABLE flashcard_log ADD COLUMN user_id INTEGER",
        "ALTER TABLE quiz_questions ADD COLUMN user_id INTEGER",
        "ALTER TABLE quiz_attempts ADD COLUMN user_id INTEGER",
        "ALTER TABLE revision_slides ADD COLUMN user_id INTEGER",
        "ALTER TABLE mind_maps ADD COLUMN user_id INTEGER",
        "ALTER TABLE exam_dates ADD COLUMN user_id INTEGER",
        "ALTER TABLE chat_messages ADD COLUMN user_id INTEGER",
        # Organisation
        "ALTER TABLE materials ADD COLUMN sort_order INTEGER DEFAULT 0",
        # Auth columns
        "ALTER TABLE users ADD COLUMN password_hash TEXT",
        # Auto-linking: related topics on flashcards and quiz questions
        "ALTER TABLE flashcards ADD COLUMN related_topics TEXT DEFAULT '[]'",
        "ALTER TABLE quiz_questions ADD COLUMN related_topics TEXT DEFAULT '[]'",
        # Chemistry: SMILES notation for structure rendering
        "ALTER TABLE flashcards ADD COLUMN smiles TEXT DEFAULT NULL",
        "ALTER TABLE quiz_questions ADD COLUMN smiles TEXT DEFAULT NULL",
        # Spaced repetition on quiz questions (same SM-2 engine as flashcards)
        "ALTER TABLE quiz_questions ADD COLUMN next_review TEXT DEFAULT NULL",
        "ALTER TABLE quiz_questions ADD COLUMN srs_interval INTEGER DEFAULT 0",
        "ALTER TABLE quiz_questions ADD COLUMN ease_factor REAL DEFAULT 2.5",
        "ALTER TABLE quiz_questions ADD COLUMN review_count INTEGER DEFAULT 0",
        "ALTER TABLE quiz_questions ADD COLUMN last_seen TEXT DEFAULT NULL",
    ]:
        try:
            conn.execute(stmt)
            conn.commit()
        except Exception:
            pass  # Column already exists

    # Sessions table for token-based auth
    conn.execute("""
        CREATE TABLE IF NOT EXISTS sessions (
            token TEXT PRIMARY KEY,
            user_id INTEGER NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            expires_at DATETIME NOT NULL,
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE
        )
    """)
    # Quiz battles
    conn.execute("""
        CREATE TABLE IF NOT EXISTS quiz_battles (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            creator_id INTEGER NOT NULL,
            topic TEXT NOT NULL,
            question_ids TEXT NOT NULL,
            status TEXT DEFAULT 'open',
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (creator_id) REFERENCES users(id) ON DELETE CASCADE
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS battle_participants (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            battle_id INTEGER NOT NULL,
            user_id INTEGER NOT NULL,
            score INTEGER DEFAULT 0,
            total INTEGER DEFAULT 0,
            completed INTEGER DEFAULT 0,
            finished_at DATETIME,
            FOREIGN KEY (battle_id) REFERENCES quiz_battles(id) ON DELETE CASCADE,
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE,
            UNIQUE(battle_id, user_id)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS graph_hidden_nodes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            node_id TEXT NOT NULL,
            hidden_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(user_id, node_id)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS graph_hidden_edges (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            source TEXT NOT NULL,
            target TEXT NOT NULL,
            hidden_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(user_id, source, target)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS graph_custom_edges (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            source TEXT NOT NULL,
            target TEXT NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(user_id, source, target)
        )
    """)
    # ── GAMSAT writing module ─────────────────────────────────────────────
    conn.execute("""
        CREATE TABLE IF NOT EXISTS writing_essays (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            task TEXT,
            theme TEXT,
            essay_text TEXT NOT NULL,
            overall_band INTEGER,
            assessment_json TEXT
        )
    """)
    # Word Bank — user-added vocabulary (merged with the frontend seed list).
    # UNIQUE(user_id, word) lets re-imports use INSERT OR IGNORE to dedupe.
    conn.execute("""
        CREATE TABLE IF NOT EXISTS word_bank (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            category TEXT DEFAULT 'vocab',
            word TEXT NOT NULL,
            pos TEXT DEFAULT '',
            definition TEXT DEFAULT '',
            example TEXT DEFAULT '',
            source TEXT DEFAULT 'import',
            UNIQUE(user_id, word)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS writing_errors (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            essay_id INTEGER,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            original TEXT,
            corrected TEXT,
            category TEXT,
            explanation TEXT,
            FOREIGN KEY (essay_id) REFERENCES writing_essays(id) ON DELETE CASCADE
        )
    """)
    # One SM-2 card per (user, grammar category). Columns mirror the flashcards
    # SM-2 fields so sm2_schedule() is reused unchanged.
    conn.execute("""
        CREATE TABLE IF NOT EXISTS writing_cards (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            category TEXT NOT NULL,
            error_count INTEGER DEFAULT 0,
            last_seen TEXT DEFAULT NULL,
            next_review TEXT DEFAULT NULL,
            srs_interval INTEGER DEFAULT 0,
            ease_factor REAL DEFAULT 2.5,
            review_count INTEGER DEFAULT 0,
            UNIQUE(user_id, category)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS writing_reviews (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            card_id INTEGER,
            created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            quality INTEGER,
            was_correct INTEGER,
            FOREIGN KEY (card_id) REFERENCES writing_cards(id) ON DELETE CASCADE
        )
    """)
    # ── Russian learning module ───────────────────────────────────────────
    # One drillable card per (user, cyrillic, english). SM-2 columns mirror
    # flashcards so sm2_schedule() is reused unchanged. phase 0=alphabet, 1..4.
    conn.execute("""
        CREATE TABLE IF NOT EXISTS russian_vocab (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            phase INTEGER DEFAULT 1,
            category TEXT DEFAULT 'vocab',
            cyrillic TEXT NOT NULL,
            translit TEXT DEFAULT '',
            english TEXT NOT NULL,
            example TEXT DEFAULT '',
            note TEXT DEFAULT '',
            source TEXT DEFAULT 'seed',
            times_seen INTEGER DEFAULT 0,
            times_correct INTEGER DEFAULT 0,
            last_seen TEXT DEFAULT NULL,
            next_review TEXT DEFAULT NULL,
            srs_interval INTEGER DEFAULT 1,
            ease_factor REAL DEFAULT 2.5,
            review_count INTEGER DEFAULT 0,
            UNIQUE(user_id, cyrillic, english)
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS russian_review_log (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            vocab_id INTEGER,
            user_id INTEGER NOT NULL,
            phase INTEGER,
            correct INTEGER DEFAULT 0,
            reviewed_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (vocab_id) REFERENCES russian_vocab(id) ON DELETE CASCADE
        )
    """)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS russian_progress (
            user_id INTEGER NOT NULL,
            phase INTEGER NOT NULL,
            status TEXT DEFAULT 'not_started',
            updated_at DATETIME DEFAULT CURRENT_TIMESTAMP,
            PRIMARY KEY (user_id, phase)
        )
    """)
    # Access log for the public study-guide gallery (who opened the shared link, when & roughly where)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS access_log (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            ts DATETIME DEFAULT CURRENT_TIMESTAMP,
            ip TEXT,
            country TEXT,
            region TEXT,
            city TEXT,
            org TEXT,
            guide TEXT,
            referer TEXT,
            user_agent TEXT,
            source TEXT DEFAULT 'gallery'
        )
    """)
    conn.commit()

    # ── Performance indexes (CREATE INDEX IF NOT EXISTS = safe to re-run) ──
    for idx in [
        "CREATE INDEX IF NOT EXISTS idx_flashcards_user    ON flashcards(user_id)",
        "CREATE INDEX IF NOT EXISTS idx_flashcards_review  ON flashcards(user_id, next_review)",
        "CREATE INDEX IF NOT EXISTS idx_fc_log_user        ON flashcard_log(user_id, reviewed_at)",
        "CREATE INDEX IF NOT EXISTS idx_quiz_attempts_user ON quiz_attempts(user_id, attempted_at)",
        "CREATE INDEX IF NOT EXISTS idx_quiz_review        ON quiz_questions(user_id, next_review)",
        "CREATE INDEX IF NOT EXISTS idx_materials_user     ON materials(user_id)",
        "CREATE INDEX IF NOT EXISTS idx_user_materials     ON user_materials(user_id)",
        "CREATE INDEX IF NOT EXISTS idx_slides_material    ON revision_slides(material_id, user_id)",
        "CREATE INDEX IF NOT EXISTS idx_mm_material        ON mind_maps(material_id, user_id)",
        "CREATE INDEX IF NOT EXISTS idx_writing_essays     ON writing_essays(user_id, created_at)",
        "CREATE INDEX IF NOT EXISTS idx_writing_errors     ON writing_errors(user_id, category)",
        "CREATE INDEX IF NOT EXISTS idx_writing_cards      ON writing_cards(user_id, next_review)",
        "CREATE INDEX IF NOT EXISTS idx_russian_vocab_user ON russian_vocab(user_id)",
        "CREATE INDEX IF NOT EXISTS idx_russian_review     ON russian_vocab(user_id, next_review)",
        "CREATE INDEX IF NOT EXISTS idx_russian_log_user   ON russian_review_log(user_id, reviewed_at)",
    ]:
        conn.execute(idx)
    conn.commit()

    # ── One-time topic cleanup: normalize verbose AI-generated topic names ─────
    long_topics = conn.execute("SELECT DISTINCT topic FROM quiz_attempts WHERE LENGTH(topic) > 25").fetchall()
    for row in long_topics:
        old = row[0]
        new = _normalize_topic(old)
        if new != old:
            conn.execute("UPDATE quiz_attempts SET topic = ? WHERE topic = ?", (new, old))
            conn.execute("UPDATE quiz_questions SET topic = ? WHERE topic = ?", (new, old))
    if long_topics:
        conn.commit()

    # Also fix flashcard_log topics
    long_fc = conn.execute("SELECT DISTINCT topic FROM flashcard_log WHERE topic IS NOT NULL AND LENGTH(topic) > 25").fetchall()
    for row in long_fc:
        old = row[0]
        new = _normalize_topic(old)
        if new != old:
            conn.execute("UPDATE flashcard_log SET topic = ? WHERE topic = ?", (new, old))
    if long_fc:
        conn.commit()

    # ── One-time backfill: claim pre-multi-user data for a default profile ─────
    # If data exists but no profiles do yet, create a "Me" profile and assign all
    # existing materials/cards/quizzes/etc. to it so nothing is lost on upgrade.
    user_count = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
    material_count = conn.execute("SELECT COUNT(*) FROM materials").fetchone()[0]
    if user_count == 0 and material_count > 0:
        conn.execute("INSERT INTO users (username) VALUES ('Me')")
        uid = conn.execute("SELECT id FROM users WHERE username = 'Me'").fetchone()[0]
        for tbl in ["materials", "flashcards", "flashcard_log", "quiz_questions",
                    "quiz_attempts", "revision_slides", "mind_maps", "exam_dates", "chat_messages"]:
            conn.execute(f"UPDATE {tbl} SET user_id = ? WHERE user_id IS NULL", (uid,))
        for r in conn.execute("SELECT id FROM materials").fetchall():
            conn.execute("INSERT OR IGNORE INTO user_materials (user_id, material_id) VALUES (?,?)", (uid, r[0]))
        conn.commit()

    conn.close()


init_db()


# ── Text extraction ──────────────────────────────────────────────────────────

def extract_pdf(path: str) -> str:
    text = ""
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text += t + "\n\n"
    except Exception as e:
        text = f"[PDF extraction error: {e}]"
    return text.strip()


def extract_pptx(path: str) -> str:
    text = ""
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, 1):
            text += f"--- Slide {i} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text.strip() + "\n"
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text += f"[Notes: {notes}]\n"
            text += "\n"
    except Exception as e:
        text = f"[PPTX extraction error: {e}]"
    return text.strip()


def extract_image(path: str) -> str:
    try:
        client = get_client()
        with open(path, "rb") as f:
            data = base64.standard_b64encode(f.read()).decode("utf-8")
        ext = Path(path).suffix.lower()
        mime = {".jpg": "image/jpeg", ".jpeg": "image/jpeg", ".png": "image/png", ".webp": "image/webp"}.get(ext, "image/jpeg")
        msg = client.messages.create(
            model=HAIKU, max_tokens=2000,  # Haiku has vision — far cheaper than Sonnet for transcription
            messages=[{"role": "user", "content": [
                {"type": "image", "source": {"type": "base64", "media_type": mime, "data": data}},
                {"type": "text", "text": "Transcribe and describe all text, diagrams, tables, and key information from this medical/health study image in full detail."}
            ]}]
        )
        return msg.content[0].text
    except Exception as e:
        return f"[Image extraction error: {e}]"


def extract_pdf_images(path: str, mat_id: int) -> list:
    """Extract embedded images from a PDF and save to IMAGES_DIR."""
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    urls = []
    try:
        with pdfplumber.open(path) as pdf:
            idx = 0
            for page in pdf.pages[:20]:
                for img_meta in (page.images or [])[:2]:
                    try:
                        x0 = float(img_meta.get('x0', 0))
                        top = float(img_meta.get('top', 0))
                        x1 = float(img_meta.get('x1', page.width))
                        bottom = float(img_meta.get('bottom', page.height))
                        if (x1 - x0) < 80 or (bottom - top) < 80:
                            continue
                        cropped = page.crop((x0, top, x1, bottom))
                        fname = f"mat{mat_id}_{idx}.png"
                        cropped.to_image(resolution=120).save(str(IMAGES_DIR / fname))
                        urls.append(f"/images/{fname}")
                        idx += 1
                        if idx >= 8:
                            return urls
                    except Exception:
                        continue
    except Exception:
        pass
    return urls


def extract_pptx_images(path: str, mat_id: int) -> list:
    """Extract embedded pictures from a PPTX and save to IMAGES_DIR."""
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    urls = []
    try:
        prs = Presentation(path)
        idx = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        image = shape.image
                        ext = (image.ext or 'png').lower()
                        if ext not in ('png', 'jpg', 'jpeg'):
                            ext = 'png'
                        fname = f"mat{mat_id}_{idx}.{ext}"
                        (IMAGES_DIR / fname).write_bytes(image.blob)
                        urls.append(f"/images/{fname}")
                        idx += 1
                        if idx >= 10:
                            return urls
                except Exception:
                    continue
    except Exception:
        pass
    return urls


def _salvage_json_array(text: str):
    """Recover the complete objects from a JSON array that was cut off mid-way
    (e.g. the model hit max_tokens, leaving a truncated trailing object).
    Returns a list of the objects that DID parse, or None if nothing usable."""
    start = text.find("[")
    if start == -1:
        return None
    objs, depth, obj_start = [], 0, None
    in_str = esc = False
    for i in range(start + 1, len(text)):
        ch = text[i]
        if in_str:
            if esc:            esc = False
            elif ch == "\\":   esc = True
            elif ch == '"':    in_str = False
            continue
        if ch == '"':
            in_str = True
        elif ch == "{":
            if depth == 0:
                obj_start = i
            depth += 1
        elif ch == "}":
            depth -= 1
            if depth == 0 and obj_start is not None:
                try:
                    objs.append(json.loads(text[obj_start:i + 1]))
                except json.JSONDecodeError:
                    pass
                obj_start = None
    return objs or None


def parse_json_response(text: str):
    text = text.strip()
    if "```json" in text:
        text = text.split("```json")[1].split("```")[0].strip()
    elif "```" in text:
        text = text.split("```")[1].split("```")[0].strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        # The model may have been truncated at max_tokens, leaving invalid JSON.
        # Salvage every complete object so we keep the questions that DID arrive
        # instead of discarding the whole batch.
        salvaged = _salvage_json_array(text)
        if salvaged:
            return salvaged
        raise


# ── Efficient generation (prompt caching) ──────────────────────────────────────

def gen_source_block(mat, content: Optional[str] = None) -> str:
    """Identical source-material text for every generator on a given material.
    Sent as the first content block and cache-tagged, so generating slides then
    flashcards/quiz/mindmap on the same material within ~5 min reuses the cache.
    `content` overrides the material body — used by Option B sectioned generation
    to send one section at a time; when None, the full body is clipped at
    GEN_CONTENT_CHARS (unchanged single-pass behaviour)."""
    body = content if content is not None else (mat['content'] or '')[:GEN_CONTENT_CHARS]
    return (
        "SOURCE STUDY MATERIAL\n"
        f"Title: {mat['original_name']}\n"
        f"Subject: {mat['subject']}\n\n"
        f"{body}"
    )


def generate_json(mat, instructions: str, model: str = HAIKU, max_tokens: int = 4000,
                  temperature: Optional[float] = None, content: Optional[str] = None) -> str:
    """One-shot generation. The (large, reusable) source material is the first block and
    is marked for prompt caching; the (small, varying) task instructions follow it.
    Repeat calls on the same material hit the cache and cost ~10% on the cached portion.
    Pass `temperature` higher (e.g. 0.9) when you want more variety between runs.
    Pass `content` to send a specific section of the material (Option B sectioned
    generation) instead of the GEN_CONTENT_CHARS-clipped full body."""
    kwargs = dict(
        model=model, max_tokens=max_tokens,
        messages=[{"role": "user", "content": [
            {"type": "text", "text": gen_source_block(mat, content), "cache_control": {"type": "ephemeral"}},
            {"type": "text", "text": instructions},
        ]}],
    )
    if temperature is not None:
        kwargs["temperature"] = temperature
    resp = get_client().messages.create(**kwargs)
    return resp.content[0].text


def _split_material_sections(text: str, target: int = 18000, overlap: int = 600) -> list:
    """Option B — order-preserving split of long study material into sections of
    roughly <= `target` chars, so a girthy module can be generated section-by-section
    at full fidelity instead of being truncated at GEN_CONTENT_CHARS.

    Prefers the module's OWN structure (SCORM/Rise lesson titles, numbered/ALL-CAPS
    headings); falls back to ~target-char windows (broken on paragraph boundaries,
    with a small overlap) when there are no clean headings. Pure/deterministic —
    unit-testable without any API call. Returns [text] unchanged when it already fits."""
    text = text or ""
    if len(text) <= target:
        return [text]
    lines = text.split("\n")

    def is_heading(ln: str) -> bool:
        s = ln.strip()
        if not (3 <= len(s) <= 80):
            return False
        if re.match(r'^(lesson|part|section|module|chapter|topic|unit)\b', s, re.I):
            return True
        if re.match(r'^\d+(\.\d+)*[\.\)]?\s+\S', s):          # "1. Foo", "2.3) Bar"
            return True
        if s == s.upper() and re.search(r'[A-Za-z]', s) and len(s.split()) <= 8:  # ALL-CAPS heading
            return True
        return False

    heads = [i for i, ln in enumerate(lines) if is_heading(ln)]
    sections: list = []
    if len(heads) >= 2:
        bounds = heads + [len(lines)]
        blocks = []
        for b in range(len(heads)):
            s_i = 0 if (b == 0 and heads[0] > 0) else heads[b]   # keep any pre-heading preamble
            blocks.append("\n".join(lines[s_i:bounds[b + 1]]))
        cur = ""
        for blk in blocks:                                       # greedily merge blocks up to target
            if cur and len(cur) + len(blk) > target:
                sections.append(cur); cur = blk
            else:
                cur = (cur + "\n" + blk) if cur else blk
        if cur:
            sections.append(cur)

    if not sections:                                             # fallback: overlapping char windows
        i, n = 0, len(text)
        while i < n:
            end = min(i + target, n)
            brk = text.rfind("\n\n", i, end)
            if brk == -1 or brk <= i + target // 2:
                brk = end
            sections.append(text[i:brk])
            if brk >= n:
                break
            i = max(brk - overlap, i + 1)

    return [s for s in sections if s.strip()] or [text[:target]]


# ── #13 Structured outputs (flashcards only, behind STRUCTURED_FLASHCARDS) ────

# Schema for output_config.format json_schema. All properties are required and
# additionalProperties is false (structured-outputs constraint); optionality is
# expressed as an empty array (related_topics) / null (smiles).
FLASHCARD_SCHEMA = {
    "type": "object",
    "properties": {
        "cards": {
            "type": "array",
            "items": {
                "type": "object",
                "properties": {
                    "topic": {"type": "string", "description": "Broad subject category, 2-3 words"},
                    "question": {"type": "string"},
                    "answer": {"type": "string"},
                    "related_topics": {"type": "array", "items": {"type": "string"},
                                       "description": "0-2 genuinely related topic names; empty array if none"},
                    "smiles": {"type": ["string", "null"],
                               "description": "SMILES string only if a chemical structure is directly relevant, else null"},
                },
                "required": ["topic", "question", "answer", "related_topics", "smiles"],
                "additionalProperties": False,
            },
        },
    },
    "required": ["cards"],
    "additionalProperties": False,
}


def generate_flashcards_structured(mat, instructions: str, max_tokens: int = 6000) -> list:
    """Flashcard generation via structured outputs (Haiku 4.5 supports
    output_config.format json_schema). Same cache-tagged source-block layout as
    generate_json, so the prompt cache is shared with the other generators.
    Raises on any problem — the caller falls back to the legacy path."""
    resp = get_client().messages.create(
        model=HAIKU, max_tokens=max_tokens,
        messages=[{"role": "user", "content": [
            {"type": "text", "text": gen_source_block(mat), "cache_control": {"type": "ephemeral"}},
            {"type": "text", "text": instructions},
        ]}],
        output_config={"format": {"type": "json_schema", "schema": FLASHCARD_SCHEMA}},
    )
    text = next(b.text for b in resp.content if b.type == "text")
    cards = json.loads(text)["cards"]
    if not isinstance(cards, list) or not cards:
        raise ValueError("structured output returned no cards")
    return cards


# ══════════════════════════════════════════════════════════════════════════════
# GAMSAT WRITING COACH — AI core (model: WRITE_MODEL / claude-opus-4-8)
#
# Opus 4.8 differs from MODEL/HAIKU in ways this module must honor:
#   • temperature/top_p/top_k are removed (they 400) — standardisation comes
#     from the frozen anchored rubric + output_config effort, not sampling knobs.
#   • Adaptive thinking may emit a thinking block first — always extract the
#     first `text`-type block (_write_text), never resp.content[0].text.
#   • Guard stop_reason == "refusal" before reading content.
# ══════════════════════════════════════════════════════════════════════════════

def _write_text(resp) -> str:
    """First text block from an Opus 4.8 response (adaptive thinking may put a
    thinking block before it, so never assume resp.content[0] is text)."""
    for block in resp.content:
        if getattr(block, "type", None) == "text":
            return block.text
    raise HTTPException(502, "Writing AI returned no text content")


def _write_call(system_blocks, user_text: str, effort: str = "high", max_tokens: int = 8000) -> str:
    """One call to the writing model. `system_blocks` is a list of system content
    blocks (the cached rubric prefix) or None. No temperature args — removed on
    Opus 4.8. Depth/consistency via output_config effort."""
    kwargs = dict(
        model=WRITE_MODEL,
        max_tokens=max_tokens,
        output_config={"effort": effort},
        messages=[{"role": "user", "content": user_text}],
    )
    if system_blocks:
        kwargs["system"] = system_blocks
    resp = get_client().messages.create(**kwargs)
    if resp.stop_reason == "refusal":
        raise HTTPException(502, "The writing AI declined this text — please revise and resubmit")
    return _write_text(resp)


# Fixed grammar taxonomy — one SM-2 drill card per category per user.
WRITING_CATEGORIES = [
    "article", "preposition", "verb-tense", "agreement", "word-choice",
    "word-order", "plurals", "spelling", "punctuation", "register",
]

# ── Frozen anchored rubric (v1.0). DO NOT edit casually: it is the calibration
#    artifact — every byte change resets the prompt cache AND shifts the scale.
#    After ANY edit, re-run the consistency self-test (/api/writing/calibrate).
WRITING_RUBRIC = """GAMSAT SECTION II — STANDARDISED MARKING RUBRIC (FROZEN v1.0)

ROLE
You are a standardised GAMSAT Written Communication examiner. You mark essays
against this rubric and NOTHING else. Your judgements must be reproducible: the
same essay must always receive the same bands. Anchor every band decision to
the descriptors and the benchmark exemplars below — never to personal taste,
never to leniency, never to the essay's length.

SCALE
Each of the four criteria is banded 1–6 (6 highest). The overall band is a
holistic 1–6 judgement consistent with the criterion bands (normally within 1
of their mean). Band meanings at a glance:
6 = exceptional · 5 = strong · 4 = competent · 3 = developing · 2 = weak · 1 = rudimentary

────────────────────────────────────────────────────────────────────
CRITERION 1 — QUALITY OF THOUGHT & ARGUMENT
Assesses: depth, originality and control of ideas; strength, nuance and
consistency of the position (Task A) or insight (Task B) developed.

Band 6: A precise, arguable thesis or genuine personal insight is sustained
throughout. Ideas are complex and nuanced — tensions, counter-positions or
qualifications are actively engaged, not ignored. Reasoning is layered: claims
develop through consequence, example and implication rather than restatement.
Nothing important is merely asserted.
Band 5: A clear thesis or insight developed with real reasoning and apt,
specific examples. Some complexity or acknowledgement of the counter-view is
present, though one or two moves remain underdeveloped or safe.
Band 4: A discernible position developed logically but predictably. Ideas are
sensible yet general; examples support rather than extend the thinking; the
counter-view is absent or token. Competent, not probing.
Band 3: A position exists but is thin, drifting, or partly inconsistent.
Assertion outweighs reasoning; examples are vague, clichéd, or only loosely
tied to the claim being made.
Band 2: Ideas are fragmentary — largely restatement of the quotes or of one
obvious commonplace. Little visible reasoning; contradictions go unnoticed.
Band 1: No discernible position or development of thought; disconnected
remarks on or near the topic.

────────────────────────────────────────────────────────────────────
CRITERION 2 — STRUCTURE & DEVELOPMENT
Assesses: the shape of the whole piece — orientation, paragraphing, logical
progression, cohesion between paragraphs, and a close that completes (rather
than repeats) the thinking.

Band 6: The essay has an architecture: each paragraph advances a single move
in a visible line of development, transitions carry logic (not just sequence
words), and the conclusion reframes or resolves rather than summarises.
Nothing could be reordered without loss.
Band 5: Clear, purposeful organisation with well-controlled paragraphs and
mostly logical transitions; the opening or close may be slightly mechanical,
or one paragraph may overreach or repeat.
Band 4: Recognisable intro–body–conclusion shape; paragraphs are unified but
the progression is additive ("another point is…") rather than cumulative; the
conclusion restates the introduction.
Band 3: Paragraphing is inconsistent or arbitrary; ideas appear in an order
that could be shuffled; links between paragraphs are missing or purely verbal;
the piece stops rather than concludes.
Band 2: Little visible organisation — one long block or fragments; ideas
repeat or interrupt each other; no functional opening or close.
Band 1: No structural control at all.

────────────────────────────────────────────────────────────────────
CRITERION 3 — LANGUAGE & EXPRESSION
Assesses: control of grammar, syntax, spelling and punctuation; precision of
vocabulary; sentence variety; fluency; consistency of register.

IMPORTANT — OBJECTIVE ANCHOR: every essay arrives with a separately computed
language-error audit (count, density per 100 words, breakdown by category) and
a HARD CAP derived from that density. This criterion's band MUST NOT exceed
the cap. Within the cap, judge precision, variety, rhythm and register.

Band 6: Error-free or near error-free prose with genuine stylistic control —
varied sentence architecture used for emphasis, exact word choice, confident
punctuation, consistent register. The language does work, not just service.
Band 5: Accurate, fluent prose with occasional minor slips; vocabulary is
precise; some sentence variety; register consistent.
Band 4: Generally accurate; errors present but rarely obstruct meaning.
Vocabulary adequate but sometimes approximate; sentences correct yet
monotonous in shape; register mostly consistent.
Band 3: Frequent errors of the taxonomy (articles, prepositions, tense,
agreement…) that a reader must read past; word choice often approximate;
limited sentence control; register wobbles.
Band 2: Dense error patterns that impede reading; very restricted syntax and
vocabulary; meaning sometimes recoverable only by guesswork.
Band 1: Errors so pervasive that meaning is frequently lost.

────────────────────────────────────────────────────────────────────
CRITERION 4 — ENGAGEMENT WITH THE STIMULUS
Assesses: whether the piece genuinely responds to the THEME of the quotes and
honours the task instruction. Task A: argumentative/analytical treatment of a
socio-political theme. Task B: reflective/personal treatment of an
interpersonal theme. Quoting the stimulus is NOT required; thinking with its
theme is.

Band 6: The theme is the essay's centre of gravity. The piece engages the
tension BETWEEN the quotes' positions (or takes one quote's idea somewhere
genuinely its own), and the response fits the task mode exactly.
Band 5: Clear, sustained engagement with the theme; may lean on one reading of
it rather than exploring the tension; task mode respected.
Band 4: On-theme throughout but treats it at the most general level; the essay
could have been half-written before seeing these particular quotes.
Band 3: Partial engagement — starts on-theme then drifts to an adjacent,
easier topic; or answers the wrong mode (e.g. narrates when asked to argue).
Band 2: Tangential — the theme appears only in the first lines or as decoration
around a prepared piece on something else.
Band 1: Essentially off-topic.

────────────────────────────────────────────────────────────────────
BENCHMARK EXEMPLARS — FIXED ANCHORS (theme for all: "comfort")
These marked excerpts define the scale. Band a new essay by asking which
anchor it sits closest to (6, 4 or 2 level), then refine ±1 by descriptors.
Bands given here are FINAL and non-negotiable.

EXEMPLAR 1 — Task A, overall band 6
"The modern state no longer disciplines its citizens; it upholsters them. We
mistake this for freedom because nothing visibly forbids us — yet a society
optimised for comfort quietly prices out its dissenters, for whom discomfort
is the entry fee of speech. Consider the office worker who will not raise a
safety concern: no law silences him; the prospect of an awkward meeting does.
Comfort, at scale, becomes a politics. To object that comfort is what
civilisation is FOR mistakes the means for the end: we build heated homes so
that we may do difficult things, not so that difficulty itself becomes
obscene. A decent society keeps its citizens warm; a serious one keeps them
capable of choosing cold."
Bands: Thought 6 · Structure 6 · Language 6 · Engagement 6.
Why: arguable thesis sustained and complicated (engages its own strongest
counter-argument); every claim developed through consequence or example;
sentence architecture does rhetorical work; the quotes' tension (comfort as
achievement vs comfort as sedative) is the essay's engine.

EXEMPLAR 2 — Task B, overall band 6
"My grandmother's kitchen was never comfortable. The chairs were hard, the
radio argued, and you could not sit long before being handed a task. It took
me twenty years and one very quiet apartment of my own to understand that I
had confused comfort with welcome. The sofa I bought asks nothing of me, and
that is precisely its poverty: nothing asked, nothing belonged to. When I
visit friends now I notice I relax most in houses that put me to work —
shelling peas, minding a pot. Perhaps ease is what we offer guests, and
demand is what we offer family; the deepest rest I have known came disguised
as a chore."
Bands: Thought 6 · Structure 6 · Language 6 · Engagement 6.
Why: a genuinely personal insight (comfort vs welcome) discovered, not
announced; concrete memory does the reasoning; the reflective mode is exact;
the close reframes rather than repeats.

EXEMPLAR 3 — Task A, overall band 4
"Comfort is one of the most important things in modern society, but it has
both advantages and disadvantages. On one hand, comfort improves people's
lives. Modern medicine, heating and transport mean people no longer suffer as
they did in the past, and this is clearly a good thing. On the other hand,
too much comfort can make people lazy. For example, many people order food
instead of cooking and drive instead of walking, which causes health problems.
Furthermore, students who are too comfortable may not push themselves to
study hard. In conclusion, comfort is beneficial but people should not let it
control their lives, and a balance should be found between comfort and
challenge."
Bands: Thought 4 · Structure 4 · Language 5 · Engagement 4.
Why: clear position, logical but entirely predictable "balance" argument;
additive structure with a restating conclusion; accurate but shape-poor prose;
on-theme at the most general level.

EXEMPLAR 4 — Task B, overall band 4
"The most comfortable place I know is my bedroom at my parents' house. When I
moved away for university I missed it a lot. My new room was smaller and the
bed was different, and for the first months I could not sleep well. Slowly I
added things: a lamp, photos of my friends, a blanket from home. One day I
realised the new room felt like mine. This experience taught me that comfort
is not about the place itself but about the memories and effort we put into
it. Now when I feel uncomfortable in a new situation, I remember that room
and know that with time I can make anywhere feel like home."
Bands: Thought 4 · Structure 4 · Language 5 · Engagement 4.
Why: sincere and unified but the insight is announced as a moral ("taught me
that…") rather than explored; competent chronology; clean but plain language.

EXEMPLAR 5 — Task A, overall band 2
"Comfort is very important in the life. Like the quote say, people want to be
comfort in their houses and jobs. I agree with this because everybody like
comfortable. In old times people was not comfortable and now they are more.
Also technology make many comforts for example phone and car and internet.
Some people say comfort is bad but I think is good because nobody want to
suffer. In conclusion comfort is very important for the people and the
society and we must to have more comfort in the future."
Bands: Thought 2 · Structure 2 · Language 2 · Engagement 3.
Why: restates the stimulus and one commonplace with no reasoning; no
functional paragraphs; dense agreement/article/verb-form error patterns
(audit cap applies); stays near the theme but only at its surface.

EXEMPLAR 6 — Task B, overall band 2
"I remember one time I was very comfortable. It was holiday with my family in
the beach. The weather it was hot and we swim every day. My mother cook fish.
It was very nice time and I was feeling comfort. Comfort is when you don't
have problems and you can relax with the persons you love. Everyone have a
place where they feel the comfort. This is important for the mental health.
So people should find their comfortable place and go there when they have
stress."
Bands: Thought 2 · Structure 2 · Language 2 · Engagement 3.
Why: a listed memory with a bolted-on moral; fragments and repeated
agreement/article errors; no development between sentences.

────────────────────────────────────────────────────────────────────
MARKING RULES (apply in order, every time)
1. Read the whole essay once without judging. Ignore its length.
2. Band each criterion by locating the essay against the exemplar anchors
   (closest to 6-, 4- or 2-level?), then refine ±1 using the descriptors.
3. Language & Expression: apply the HARD CAP from the error audit. Never
   exceed it; you may band below it on precision/variety grounds.
4. Every criterion band MUST cite `evidence`: a short VERBATIM quote from the
   candidate's essay (not from the stimulus, not paraphrased) that most
   influenced that band, plus a one-to-two sentence justification.
5. Overall band: holistic 1–6, normally within 1 of the criterion mean. Do
   not average mechanically, but never contradict the criterion bands.
6. Strengths and priority_improvements must be specific and actionable —
   name the habit, show the fix. No generic advice ("add more detail").
7. Output STRICTLY the requested JSON. No preamble, no markdown fences.
"""


def _language_band_cap(word_count: int, error_count: int) -> int:
    """Objective ceiling on the Language & Expression band, from error density
    (errors per 100 words). Anchors the grammar judgement to countable evidence."""
    if word_count <= 0:
        return 1
    density = error_count * 100.0 / word_count
    if density <= 0.8:  return 6
    if density <= 1.6:  return 5
    if density <= 3.0:  return 4
    if density <= 5.0:  return 3
    if density <= 8.0:  return 2
    return 1


def writing_analyze_language(text: str) -> list:
    """Deterministic-ish grammar pass over a fixed taxonomy. Returns a list of
    {original, corrected, category, explanation}. This is the objective anchor
    that caps the grader's Language & Expression band."""
    prompt = (
        "You are a precise, conservative English grammar auditor. Find every GENUINE "
        "language error in the text below — do NOT flag stylistic preferences, "
        "acceptable informal usage, or debatable comma choices. Be exhaustive on real "
        "errors, silent on style.\n\n"
        "Each error must be classified into EXACTLY one of these categories:\n"
        f"{', '.join(WRITING_CATEGORIES)}\n\n"
        "(register = wrong formality level for an essay, e.g. slang or texting language)\n\n"
        "Return ONLY a JSON array (no prose, no fences). Each element:\n"
        '{"original": "<smallest exact substring containing the error>", '
        '"corrected": "<the fixed version of that substring>", '
        '"category": "<one category>", '
        '"explanation": "<one short sentence naming the rule broken>"}\n'
        "Report each distinct error once. If the text has no genuine errors, return [].\n\n"
        "TEXT TO AUDIT:\n" + text
    )
    raw = _write_call(None, prompt, effort="medium", max_tokens=4000)
    data = parse_json_response(raw)
    if isinstance(data, dict):
        data = data.get("errors") or []
    out = []
    for e in data or []:
        if not isinstance(e, dict):
            continue
        cat = (e.get("category") or "").strip().lower().replace(" ", "-").replace("_", "-")
        if cat not in WRITING_CATEGORIES:
            cat = "word-choice"
        out.append({
            "original":    (e.get("original") or "").strip(),
            "corrected":   (e.get("corrected") or "").strip(),
            "category":    cat,
            "explanation": (e.get("explanation") or "").strip(),
        })
    return out


def writing_assess(stimulus: dict, essay: str, error_stats: dict) -> dict:
    """The standardised grader. The frozen rubric is sent as a cache-tagged system
    block (byte-identical every call → cached, cheap, and anchored). The grader
    returns grounded JSON: every criterion band cites verbatim essay evidence.
    The Language & Expression band is hard-capped by the objective error audit."""
    system_blocks = [{
        "type": "text",
        "text": WRITING_RUBRIC,
        "cache_control": {"type": "ephemeral"},
    }]
    cap = error_stats.get("language_band_cap", 6)
    quotes = "\n".join(f"- {q}" for q in (stimulus.get("quotes") or []))
    by_cat = json.dumps(error_stats.get("by_category") or {}, sort_keys=True)
    user = (
        f"STIMULUS (Task {stimulus.get('task','A')})\n"
        f"Theme: {stimulus.get('theme','')}\n"
        f"Quotes:\n{quotes}\n"
        f"Instruction given to the candidate: {stimulus.get('instruction','')}\n\n"
        f"CANDIDATE ESSAY (verbatim):\n\"\"\"\n{essay}\n\"\"\"\n\n"
        "OBJECTIVE LANGUAGE-ERROR AUDIT (computed separately — treat as ground truth):\n"
        f"- word count: {error_stats.get('word_count')}\n"
        f"- genuine errors found: {error_stats.get('error_count')}\n"
        f"- errors per 100 words: {error_stats.get('per_100_words')}\n"
        f"- by category: {by_cat}\n"
        f"- HARD CAP on the Language & Expression band: {cap}. Do not exceed it.\n\n"
        "Mark this essay against the frozen rubric. Be concrete and specific — cite the\n"
        "candidate's OWN words as evidence, never generic praise. For each criterion give\n"
        "`evidence` as an ARRAY of 1-3 short verbatim quotes from the essay that justify the\n"
        "band (use more quotes for longer essays — a long essay must not receive a one-line\n"
        "verdict). `justification` should explain, in 2-3 sentences, exactly what those quotes\n"
        "show and what would lift the band. Make `strengths` and `priority_improvements`\n"
        "specific and actionable, and scale their number to the essay's length (aim ~1 item\n"
        "per 120 words, min 2, max 6). Return ONLY this JSON object:\n"
        "{\n"
        '  "criteria": [\n'
        '    {"name": "Quality of Thought & Argument", "band": <1-6>, "evidence": ["<verbatim quote>", "..."], "justification": "<2-3 sentences>"},\n'
        '    {"name": "Structure & Development", "band": <1-6>, "evidence": ["..."], "justification": "..."},\n'
        '    {"name": "Language & Expression", "band": <1-6>, "evidence": ["..."], "justification": "..."},\n'
        '    {"name": "Engagement with the Stimulus", "band": <1-6>, "evidence": ["..."], "justification": "..."}\n'
        "  ],\n"
        '  "overall_band": <1-6>,\n'
        '  "strengths": ["<specific strength>", "..."],\n'
        '  "priority_improvements": ["<specific, actionable fix>", "..."]\n'
        "}"
    )
    raw = _write_call(system_blocks, user, effort="high", max_tokens=8000)
    a = parse_json_response(raw)
    # The grader sometimes returns the four criteria as a BARE ARRAY instead of the
    # wrapped {"criteria": [...]} object. Detect that shape and re-wrap it, rather
    # than keeping only a[0] (which silently discarded 3 criteria and defaulted the
    # whole essay to band 1).
    if isinstance(a, list):
        crit_like = [x for x in a if isinstance(x, dict) and ("band" in x or "name" in x)]
        if len(crit_like) >= 2:
            a = {"criteria": crit_like}
        else:
            a = a[0] if a and isinstance(a[0], dict) else {}
    if not isinstance(a, dict):
        a = {}
    # ── Enforce shape + the objective cap in code (never trust the model alone) ──
    crits = [c for c in (a.get("criteria") or []) if isinstance(c, dict) and c.get("band") is not None]
    # If parsing still yielded no scored criteria, the grade is meaningless — fail
    # loudly so the caller can retry, instead of persisting a bogus band 1.
    if not crits:
        raise HTTPException(502, "The grader returned an unreadable response — please submit again.")
    capped = False
    for c in crits:
        try:
            c["band"] = max(1, min(6, int(c.get("band", 1))))
        except (TypeError, ValueError):
            c["band"] = 1
        if "language" in (c.get("name") or "").lower() and c["band"] > cap:
            c["band"] = cap
            c["capped_by_error_density"] = True
            capped = True
    a["criteria"] = crits
    try:
        overall = max(1, min(6, int(a.get("overall_band", 0))))
    except (TypeError, ValueError):
        overall = 0
    if crits and (capped or not overall):
        overall = max(1, min(6, round(sum(c["band"] for c in crits) / len(crits))))
    a["overall_band"] = overall or 1
    a.setdefault("strengths", [])
    a.setdefault("priority_improvements", [])
    return a


# ── Stimulus generation + hand-written seed bank ──────────────────────────────

WRITING_TASK_INSTRUCTIONS = {
    "A": ("Consider the following statements on a common theme. Write an argumentative or "
          "analytical response to the ideas they raise. Develop and defend a position of your "
          "own — you may draw on any of the quotes, but you are not required to reference them."),
    "B": ("Consider the following statements on a common theme. Write a reflective or personal "
          "response to the ideas they raise. Explore what the theme means in your own experience "
          "and understanding — you may draw on any of the quotes, but you are not required to "
          "reference them."),
}

WRITING_SEED_STIMULI = [
    {"task": "A", "theme": "Freedom and security",
     "quotes": [
        "Those who would give up essential liberty to purchase a little temporary safety deserve neither. — Benjamin Franklin",
        "Freedom is not worth having if it does not include the freedom to make mistakes. — Mahatma Gandhi",
        "The cage went in search of a bird. — Franz Kafka",
        "Most people do not really want freedom, because freedom involves responsibility. — Sigmund Freud",
     ]},
    {"task": "A", "theme": "Technology and progress",
     "quotes": [
        "It has become appallingly obvious that our technology has exceeded our humanity. — attributed to Albert Einstein",
        "The real problem is not whether machines think but whether men do. — B. F. Skinner",
        "We shape our tools, and thereafter our tools shape us. — attributed to Marshall McLuhan",
        "Progress is a comfortable disease. — E. E. Cummings",
        "Any sufficiently advanced technology is indistinguishable from magic. — Arthur C. Clarke",
     ]},
    {"task": "A", "theme": "Wealth and inequality",
     "quotes": [
        "The law, in its majestic equality, forbids rich and poor alike to sleep under bridges. — Anatole France",
        "Poverty is the parent of revolution and crime. — Aristotle",
        "No one has ever become poor by giving. — Anne Frank",
        "The rich would have to eat money if the poor did not provide food. — Assyrian proverb",
     ]},
    {"task": "B", "theme": "Failure",
     "quotes": [
        "Ever tried. Ever failed. No matter. Try again. Fail again. Fail better. — Samuel Beckett",
        "Success is stumbling from failure to failure with no loss of enthusiasm. — attributed to Winston Churchill",
        "There is no failure except in no longer trying. — Elbert Hubbard",
        "I have not failed. I've just found ten thousand ways that won't work. — Thomas Edison",
     ]},
    {"task": "B", "theme": "Belonging",
     "quotes": [
        "The ache for home lives in all of us, the safe place where we can go as we are. — Maya Angelou",
        "I am not an Athenian or a Greek, but a citizen of the world. — attributed to Socrates",
        "You only are free when you realize you belong no place — you belong every place. — Maya Angelou",
        "We're born alone, we live alone, we die alone. Only through love and friendship can we create the illusion that we're not. — Orson Welles",
     ]},
    {"task": "B", "theme": "Friendship",
     "quotes": [
        "A friend is one that knows you as you are and still, gently, allows you to grow. — attributed to William Shakespeare",
        "Friendship is unnecessary, like philosophy, like art. It has no survival value; rather it gives value to survival. — C. S. Lewis",
        "It is not a lack of love, but a lack of friendship that makes unhappy marriages. — Friedrich Nietzsche",
        "The worst solitude is to be destitute of sincere friendship. — Francis Bacon",
     ]},
]


def writing_stimulus(task: str = "A") -> dict:
    """Generate a fresh GAMSAT-style stimulus: theme + 3–5 quotes + instruction."""
    task = "B" if str(task).upper() == "B" else "A"
    kind = ("a socio-political theme suited to argumentative/analytical writing "
            "(e.g. justice, power, censorship, education, science and society)"
            if task == "A" else
            "an interpersonal/personal theme suited to reflective writing "
            "(e.g. grief, courage, family, ambition, solitude)")
    prompt = (
        "Create one GAMSAT Section II writing stimulus.\n"
        f"Pick {kind} — choose something fresh, not one of these already-used themes: "
        + ", ".join(sorted({s["theme"] for s in WRITING_SEED_STIMULI})) + ".\n"
        "Provide 4 short quotes on that theme from real, attributable sources where possible "
        "(mark uncertain attributions with 'attributed to'). The quotes should TENSION against "
        "each other — at least two should pull in opposing directions.\n"
        'Return ONLY JSON: {"theme": "...", "quotes": ["<quote> — <source>", ...]}'
    )
    raw = _write_call(None, prompt, effort="medium", max_tokens=1200)
    data = parse_json_response(raw)
    if isinstance(data, list):
        data = data[0] if data and isinstance(data[0], dict) else {}
    return {
        "task": task,
        "theme": (data.get("theme") or "").strip(),
        "quotes": [q for q in (data.get("quotes") or []) if isinstance(q, str)][:5],
        "instruction": WRITING_TASK_INSTRUCTIONS[task],
    }


def writing_drills(categories: list, examples_by_cat: dict) -> list:
    """ONE call producing one CONCEPT drill per due category, seeded from that
    student's own past mistakes. Each drill is a short conceptual question about
    the rule the student keeps breaking; they write an answer in their own words,
    then reveal a model answer and self-mark (SM-2). Teaches the foundation, not
    just pattern-matching."""
    payload = json.dumps(
        [{"category": c, "past_mistakes": examples_by_cat.get(c, [])} for c in categories],
        indent=1)
    prompt = (
        "You are a writing tutor building CONCEPT drills for one student, targeting the "
        "grammar/expression rules they personally keep breaking. For EACH category below, "
        "write ONE short conceptual question that makes the student explain or apply the "
        "underlying rule in their OWN words — not just spot an error. Ground the question in "
        "the KIND of mistake shown in their past examples, but ask about the principle "
        "(e.g. for a 'then/than' slip: \"In your own words, when do you use 'than' vs "
        "'then'? Give an example of each.\").\n\n"
        f"CATEGORIES AND THE STUDENT'S PAST MISTAKES:\n{payload}\n\n"
        "Return ONLY a JSON array, one element per category, in the same order:\n"
        '[{"category": "<category>", '
        '"question": "<a concept question the student answers in their own words>", '
        '"ideal_answer": "<a clear model answer (2-4 sentences) with a concrete example>", '
        '"explanation": "<one sentence naming the rule and the trap to watch for>"}]'
    )
    raw = _write_call(None, prompt, effort="medium", max_tokens=2500)
    data = parse_json_response(raw)
    if isinstance(data, dict):
        data = data.get("drills") or []
    out = []
    for d in data or []:
        if not isinstance(d, dict):
            continue
        cat = (d.get("category") or "").strip().lower().replace(" ", "-").replace("_", "-")
        if cat in categories and d.get("question") and d.get("ideal_answer"):
            out.append({"category": cat, "question": d["question"].strip(),
                        "ideal_answer": d["ideal_answer"].strip(),
                        "explanation": (d.get("explanation") or "").strip()})
    return out


# ── Consistency self-test (standardisation, measured) ─────────────────────────

WRITING_BENCHMARK_ESSAY = """Comfort is often praised as the reward of a life well lived, but I would argue it is better understood as a test. The quotes suggest that comfort can soothe us or soften us; both are true, and the difference lies in what we do next. When a person becomes comfortable, they stop asking questions about their situation. For example, an employee with a secure job rarely challenge the decisions of their manager, even when those decisions are wrong, because the comfort of the salary matters more then the discomfort of speaking up. In this way comfort operates like a quiet contract: we trade our voice for our ease. However, it would be too simple to say comfort is bad. Without some comfort, people cannot take risks at all — a person worried about food does not write novels. The real question is wether we treat comfort as a base camp or as a destination. Societies that treat it as a destination stop building, stop arguing, and slowly stop meaning anything. I believe the task for each person is to accept comfort gratefully and then deliberately leave it, again and again."""

def writing_calibrate(n: int = 3) -> dict:
    """Grade one fixed benchmark essay N times against the frozen rubric and report
    per-criterion band variance. Standardisation target: spread ≤ 1 per criterion.
    Each run costs credit — keep N small (default 3)."""
    n = max(1, min(3, int(n)))
    stim = dict(WRITING_SEED_STIMULI[0])  # fixed stimulus; benchmark is on 'comfort'
    stim = {"task": "A", "theme": "Comfort",
            "quotes": ["Comfort is the enemy of achievement. — attributed to Farrah Gray",
                       "Civilisation is the history of making people more comfortable. — anon"],
            "instruction": WRITING_TASK_INSTRUCTIONS["A"]}
    errors = writing_analyze_language(WRITING_BENCHMARK_ESSAY)  # once — audit is shared
    wc = len(WRITING_BENCHMARK_ESSAY.split())
    stats = {
        "word_count": wc, "error_count": len(errors),
        "per_100_words": round(len(errors) * 100.0 / wc, 2),
        "by_category": {c: sum(1 for e in errors if e["category"] == c)
                        for c in {e["category"] for e in errors}},
        "language_band_cap": _language_band_cap(wc, len(errors)),
    }
    runs = []
    for _ in range(n):
        a = writing_assess(stim, WRITING_BENCHMARK_ESSAY, stats)
        runs.append({
            "overall": a["overall_band"],
            "bands": {c["name"]: c["band"] for c in a["criteria"]},
        })
    names = sorted({name for r in runs for name in r["bands"]})
    spread = {}
    for name in names:
        vals = [r["bands"][name] for r in runs if name in r["bands"]]
        spread[name] = {"min": min(vals), "max": max(vals), "spread": max(vals) - min(vals)}
    overalls = [r["overall"] for r in runs]
    return {
        "n": n, "runs": runs, "per_criterion": spread,
        "overall": {"min": min(overalls), "max": max(overalls),
                    "spread": max(overalls) - min(overalls)},
        "error_audit": stats,
    }


# ── Auth / identity (password-based with session tokens) ───────────────────────

SESSION_TTL_DAYS = 30  # sessions last 30 days

def _create_session(db, user_id: int) -> str:
    """Create a new session token for a user, clean up expired sessions."""
    token = secrets.token_hex(32)
    expires = (datetime.utcnow() + timedelta(days=SESSION_TTL_DAYS)).strftime("%Y-%m-%d %H:%M:%S")
    db.execute("INSERT INTO sessions (token, user_id, expires_at) VALUES (?,?,?)", (token, user_id, expires))
    # Prune old sessions while we're here
    db.execute("DELETE FROM sessions WHERE expires_at < ?", (datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S"),))
    db.commit()
    return token

def get_current_user(request: Request) -> int:
    """Resolve the active user from either:
    1. Authorization: Bearer <token>  (new session-token auth)
    2. X-User-Id header               (legacy — kept for bookmarklet/console clipper)
    """
    # Try session token first
    auth = request.headers.get("Authorization", "")
    if auth.startswith("Bearer "):
        token = auth[7:].strip()
        if token:
            db = get_db()
            row = db.execute(
                "SELECT user_id FROM sessions WHERE token = ? AND expires_at > ?",
                (token, datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S"))
            ).fetchone()
            db.close()
            if row:
                return row["user_id"]
            raise HTTPException(401, "Session expired — please log in again")

    # Fallback: legacy X-User-Id header (bookmarklet, old clients)
    x_user_id = request.headers.get("X-User-Id")
    if x_user_id:
        try:
            uid = int(x_user_id)
        except (TypeError, ValueError):
            raise HTTPException(401, "Invalid profile")
        db = get_db()
        row = db.execute("SELECT id FROM users WHERE id = ?", (uid,)).fetchone()
        db.close()
        if not row:
            raise HTTPException(401, "Profile not found")
        return uid

    raise HTTPException(401, "Not logged in")


def user_can_access(db, mid: int, user_id: int) -> bool:
    """True if the material is owned by the user, in their library, or public."""
    row = db.execute("SELECT user_id, visibility FROM materials WHERE id = ?", (mid,)).fetchone()
    if not row:
        return False
    if row["user_id"] == user_id or row["visibility"] == "public":
        return True
    lib = db.execute(
        "SELECT 1 FROM user_materials WHERE user_id = ? AND material_id = ?", (user_id, mid)
    ).fetchone()
    return lib is not None


# ── Routes ───────────────────────────────────────────────────────────────────

@app.get("/")
def root():
    # Stamp asset URLs with the newest mtime of app.js/style.css so browsers
    # always fetch fresh JS/CSS after a deploy (no stale-cache surprises).
    try:
        ver = int(max(
            os.path.getmtime("static/app.js"),
            os.path.getmtime("static/style.css"),
        ))
    except OSError:
        ver = int(datetime.now().timestamp())
    html = Path("static/index.html").read_text(encoding="utf-8").replace("__ASSET_VER__", str(ver))
    return HTMLResponse(html, headers={"Cache-Control": "no-cache"})

@app.get("/bookmarklet")
def bookmarklet_page():
    return FileResponse("static/bookmarklet.html")


@app.get("/guides")
def guides_gallery():
    # Public, shareable study-guides gallery (bypasses the access gate — see access_guard).
    # Guide files + guides.json live in static/guides/ and are served via /static/.
    return FileResponse("static/guides/index.html", headers={"Cache-Control": "no-cache"})


# ── Gallery access logging (self-hosted: who opened the shared link, when & roughly where) ──
def _geo_lookup(ip: str) -> dict:
    """Best-effort city-level geolocation for an IP (free ip-api.com, no key).
    Returns {} for private/local IPs or on any failure — never raises."""
    if not ip or ip.startswith(("127.", "10.", "192.168.", "172.", "169.254.")) or ip in ("::1", "localhost"):
        return {}
    try:
        import urllib.request
        url = f"http://ip-api.com/json/{ip}?fields=status,country,regionName,city,org"
        with urllib.request.urlopen(url, timeout=3) as r:
            d = json.loads(r.read().decode())
        if d.get("status") == "success":
            return {"country": d.get("country"), "region": d.get("regionName"),
                    "city": d.get("city"), "org": d.get("org")}
    except Exception:
        pass
    return {}

# 1×1 transparent GIF
_BEACON_PIXEL = base64.b64decode("R0lGODlhAQABAIAAAAAAAP///yH5BAEAAAAALAAAAAABAAEAAAIBRAA7")

@app.get("/api/beacon")
async def access_beacon(request: Request, g: str = "", r: str = ""):
    """Public tracking pixel for the study-guide gallery. Records one access row
    (time, best-effort city, guide, referrer, user-agent) then returns a 1×1 GIF.
    Whitelisted in access_guard so remote visitors log without needing the code."""
    ip = (request.headers.get("CF-Connecting-IP")
          or request.headers.get("X-Forwarded-For", "").split(",")[0].strip()
          or (request.client.host if request.client else ""))
    geo = _geo_lookup(ip)
    try:
        conn = get_db()
        conn.execute(
            "INSERT INTO access_log (ip,country,region,city,org,guide,referer,user_agent,source) "
            "VALUES (?,?,?,?,?,?,?,?,?)",
            (ip, geo.get("country"), geo.get("region"), geo.get("city"), geo.get("org"),
             (g or "")[:200], (r or request.headers.get("Referer", ""))[:300],
             request.headers.get("User-Agent", "")[:300], "gallery"))
        conn.commit(); conn.close()
    except Exception:
        pass
    return Response(content=_BEACON_PIXEL, media_type="image/gif",
                    headers={"Cache-Control": "no-store, no-cache, must-revalidate, max-age=0"})

@app.get("/api/access-log")
def access_log_view(limit: int = 200):
    """Owner-only recent gallery accesses (newest first). Protected by the normal
    access-code middleware — not whitelisted, so it needs X-Access-Code or ?ac=."""
    conn = get_db()
    rows = conn.execute(
        "SELECT ts, city, region, country, org, guide, referer, ip "
        "FROM access_log ORDER BY id DESC LIMIT ?", (min(max(limit, 1), 1000),)).fetchall()
    conn.close()
    return {"count": len(rows), "hits": [dict(x) for x in rows]}



@app.get("/api/profiles")
def list_profiles():
    """Public list of profiles (names + avatars only). No passwords exposed."""
    db = get_db()
    rows = db.execute("SELECT id, username, created_at FROM users ORDER BY created_at").fetchall()
    db.close()
    return [{"id": r["id"], "username": r["username"], "has_password": True} for r in rows]


@app.post("/api/register")
async def register(request: Request):
    """Create a new account with username + password. Returns a session token."""
    body = await request.json()
    username = (body.get("username") or "").strip()[:40]
    password = body.get("password", "")
    if not username:
        raise HTTPException(400, "Username required")
    if len(password) < 4:
        raise HTTPException(400, "Password must be at least 4 characters")
    db = get_db()
    existing = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()
    if existing:
        db.close()
        raise HTTPException(409, "That name is taken — pick another")
    pw_hash = hash_password(password)
    cur = db.execute("INSERT INTO users (username, password_hash) VALUES (?,?)", (username, pw_hash))
    uid = cur.lastrowid
    db.commit()
    token = _create_session(db, uid)
    db.close()
    return {"id": uid, "username": username, "token": token}


@app.post("/api/login")
async def login(request: Request):
    """Log in with username + password. Returns a session token."""
    body = await request.json()
    username = (body.get("username") or "").strip()
    password = body.get("password", "")
    if not username or not password:
        raise HTTPException(400, "Username and password required")
    db = get_db()
    row = db.execute("SELECT id, username, password_hash FROM users WHERE username = ?", (username,)).fetchone()
    if not row:
        db.close()
        raise HTTPException(401, "Wrong username or password")
    # Legacy profile with no password — let them set one
    if not row["password_hash"]:
        pw_hash = hash_password(password)
        db.execute("UPDATE users SET password_hash = ? WHERE id = ?", (pw_hash, row["id"]))
        db.commit()
        token = _create_session(db, row["id"])
        db.close()
        return {"id": row["id"], "username": row["username"], "token": token, "password_set": True}
    if not verify_password(password, row["password_hash"]):
        db.close()
        raise HTTPException(401, "Wrong username or password")
    token = _create_session(db, row["id"])
    db.close()
    return {"id": row["id"], "username": row["username"], "token": token}


@app.post("/api/logout")
async def logout(request: Request):
    """Invalidate the current session token."""
    auth = request.headers.get("Authorization", "")
    if auth.startswith("Bearer "):
        token = auth[7:].strip()
        if token:
            db = get_db()
            db.execute("DELETE FROM sessions WHERE token = ?", (token,))
            db.commit()
            db.close()
    return {"ok": True}


@app.post("/api/profiles")
async def create_profile_legacy(request: Request):
    """Legacy endpoint — creates profile with a default password for old clients."""
    body = await request.json()
    username = (body.get("username") or "").strip()[:40]
    password = body.get("password") or "1234"
    if not username:
        raise HTTPException(400, "Username required")
    if len(password) < 4:
        password = "1234"
    db = get_db()
    existing = db.execute("SELECT id FROM users WHERE username = ?", (username,)).fetchone()
    if existing:
        db.close()
        raise HTTPException(409, "That name is taken — pick another")
    pw_hash = hash_password(password)
    cur = db.execute("INSERT INTO users (username, password_hash) VALUES (?,?)", (username, pw_hash))
    uid = cur.lastrowid
    db.commit()
    token = _create_session(db, uid)
    db.close()
    return {"id": uid, "username": username, "token": token}


@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...), subject: str = Form(default="Medicine"),
                      user_id: int = Depends(get_current_user)):
    ext = Path(file.filename).suffix.lower()
    allowed = {".pdf", ".pptx", ".ppt", ".png", ".jpg", ".jpeg", ".webp"}
    if ext not in allowed:
        raise HTTPException(400, f"Unsupported file type: {ext}")

    safe = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{file.filename}"
    path = UPLOAD_DIR / safe
    content = await file.read()
    path.write_bytes(content)

    if ext == ".pdf":
        text, ftype = extract_pdf(str(path)), "pdf"
    elif ext in {".pptx", ".ppt"}:
        text, ftype = extract_pptx(str(path)), "pptx"
    else:
        text, ftype = extract_image(str(path)), "image"

    if len(text) < 50:
        text = "[Minimal text — AI will process the content directly]"

    db = get_db()
    cur = db.execute(
        "INSERT INTO materials (user_id, filename, original_name, subject, content, file_type, images) VALUES (?,?,?,?,?,?,?)",
        (user_id, safe, file.filename, subject, text, ftype, '[]')
    )
    mid = cur.lastrowid
    db.execute("INSERT OR IGNORE INTO user_materials (user_id, material_id) VALUES (?,?)", (user_id, mid))
    db.commit()

    # Extract images now that we have the material ID
    image_urls: list = []
    if ext == ".pdf":
        image_urls = extract_pdf_images(str(path), mid)
    elif ext in {".pptx", ".ppt"}:
        image_urls = extract_pptx_images(str(path), mid)
    else:
        # Direct image upload — copy to IMAGES_DIR so it's URL-accessible
        img_fname = f"mat{mid}_0{ext}"
        IMAGES_DIR.mkdir(parents=True, exist_ok=True)
        (IMAGES_DIR / img_fname).write_bytes(content)
        image_urls = [f"/images/{img_fname}"]

    if image_urls:
        db.execute("UPDATE materials SET images = ? WHERE id = ?", (json.dumps(image_urls), mid))
        db.commit()

    db.close()
    return {"id": mid, "name": file.filename, "subject": subject, "type": ftype,
            "chars": len(text), "images": len(image_urls)}


@app.get("/api/materials")
def list_materials(user_id: int = Depends(get_current_user)):
    """Materials in the current profile's library — their own uploads plus public
    materials they've added from other people."""
    db = get_db()
    _ensure_all_guide_quizzes(db, user_id)   # so the quiz dropdown lists every gallery guide
    rows = db.execute(
        """SELECT m.id, m.original_name, m.subject, m.file_type, m.uploaded_at,
                  LENGTH(m.content) as chars, m.images, m.visibility,
                  COALESCE(m.sort_order, 0) as sort_order,
                  m.user_id as owner_id, u.username as owner_name
           FROM materials m
           JOIN user_materials um ON um.material_id = m.id
           LEFT JOIN users u ON u.id = m.user_id
           WHERE um.user_id = ?
           ORDER BY COALESCE(m.sort_order, 0) ASC, m.uploaded_at DESC""",
        (user_id,)
    ).fetchall()
    db.close()
    result = []
    for r in rows:
        d = dict(r)
        d['is_owner'] = (d.get('owner_id') == user_id)
        try:
            d['images'] = json.loads(d.get('images') or '[]')
        except Exception:
            d['images'] = []
        result.append(d)
    return result


@app.get("/api/discover")
def discover_materials(user_id: int = Depends(get_current_user)):
    """Public materials shared by other profiles that aren't already in my library."""
    db = get_db()
    rows = db.execute(
        """SELECT m.id, m.original_name, m.subject, m.file_type, m.uploaded_at,
                  LENGTH(m.content) as chars, u.username as owner_name
           FROM materials m
           LEFT JOIN users u ON u.id = m.user_id
           WHERE m.visibility = 'public' AND m.user_id != ?
             AND m.id NOT IN (SELECT material_id FROM user_materials WHERE user_id = ?)
           ORDER BY m.uploaded_at DESC""",
        (user_id, user_id)
    ).fetchall()
    db.close()
    return [dict(r) for r in rows]


@app.post("/api/materials/{mid}/visibility")
async def set_visibility(mid: int, request: Request, user_id: int = Depends(get_current_user)):
    body = await request.json()
    vis = body.get("visibility")
    if vis not in ("public", "private"):
        raise HTTPException(400, "visibility must be 'public' or 'private'")
    db = get_db()
    mat = db.execute("SELECT user_id FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        db.close()
        raise HTTPException(404, "Material not found")
    if mat["user_id"] != user_id:
        db.close()
        raise HTTPException(403, "Only the owner can change visibility")
    db.execute("UPDATE materials SET visibility = ? WHERE id = ?", (vis, mid))
    db.commit()
    db.close()
    return {"ok": True, "visibility": vis}


@app.post("/api/materials/{mid}/add")
def add_material(mid: int, user_id: int = Depends(get_current_user)):
    """Add a public material to my library and COPY the owner's generated content
    (slides, flashcards, quiz, mind map) as my own starting point — SRS reset, no AI used."""
    db = get_db()
    mat = db.execute("SELECT user_id, visibility FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        db.close()
        raise HTTPException(404, "Material not found")
    if mat["user_id"] == user_id:
        db.close()
        raise HTTPException(400, "That's already your material")
    if mat["visibility"] != "public":
        db.close()
        raise HTTPException(403, "This material isn't shared")

    already = db.execute(
        "SELECT 1 FROM user_materials WHERE user_id = ? AND material_id = ?", (user_id, mid)
    ).fetchone()
    if already:
        db.close()
        return {"ok": True, "already": True}

    owner = mat["user_id"]
    db.execute("INSERT OR IGNORE INTO user_materials (user_id, material_id) VALUES (?,?)", (user_id, mid))

    # Copy the owner's slides
    for r in db.execute(
        "SELECT title, content, slide_order FROM revision_slides WHERE material_id = ? AND user_id = ? ORDER BY slide_order",
        (mid, owner)
    ).fetchall():
        db.execute(
            "INSERT INTO revision_slides (material_id, user_id, title, content, slide_order) VALUES (?,?,?,?,?)",
            (mid, user_id, r["title"], r["content"], r["slide_order"])
        )

    # Copy the owner's flashcards — reset SRS so the borrower starts fresh
    for r in db.execute(
        "SELECT topic, question, answer FROM flashcards WHERE material_id = ? AND user_id = ?",
        (mid, owner)
    ).fetchall():
        db.execute(
            "INSERT INTO flashcards (material_id, user_id, topic, question, answer) VALUES (?,?,?,?,?)",
            (mid, user_id, r["topic"], r["question"], r["answer"])
        )

    # Copy the owner's quiz questions
    for r in db.execute(
        "SELECT topic, question, options, correct_answer, explanation, difficulty FROM quiz_questions WHERE material_id = ? AND user_id = ?",
        (mid, owner)
    ).fetchall():
        db.execute(
            "INSERT INTO quiz_questions (material_id, user_id, topic, question, options, correct_answer, explanation, difficulty) VALUES (?,?,?,?,?,?,?,?)",
            (mid, user_id, r["topic"], r["question"], r["options"], r["correct_answer"], r["explanation"], r["difficulty"])
        )

    # Copy the owner's mind map
    mm = db.execute(
        "SELECT title, data FROM mind_maps WHERE material_id = ? AND user_id = ?", (mid, owner)
    ).fetchone()
    if mm:
        db.execute(
            "INSERT INTO mind_maps (material_id, user_id, title, data) VALUES (?,?,?,?)",
            (mid, user_id, mm["title"], mm["data"])
        )

    db.commit()
    db.close()
    return {"ok": True}


@app.delete("/api/materials/{mid}")
def delete_material(mid: int, user_id: int = Depends(get_current_user)):
    db = get_db()
    mat = db.execute("SELECT user_id FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        db.close()
        raise HTTPException(404, "Material not found")
    # Remove the current profile's own derived content + their library link
    for tbl in ["flashcard_log", "flashcards", "quiz_questions", "quiz_attempts", "revision_slides", "mind_maps"]:
        db.execute(f"DELETE FROM {tbl} WHERE material_id = ? AND user_id = ?", (mid, user_id))
    db.execute("DELETE FROM user_materials WHERE material_id = ? AND user_id = ?", (mid, user_id))
    # If the owner removes it and nobody else still has it, delete the material itself
    if mat["user_id"] == user_id:
        others = db.execute("SELECT COUNT(*) as c FROM user_materials WHERE material_id = ?", (mid,)).fetchone()["c"]
        if others == 0:
            db.execute("DELETE FROM materials WHERE id = ?", (mid,))
    db.commit()
    db.close()
    return {"ok": True}


@app.patch("/api/materials/{mid}")
async def update_material(mid: int, request: Request, user_id: int = Depends(get_current_user)):
    """Rename a material or change its subject. Owner-only."""
    db = get_db()
    mat = db.execute("SELECT user_id FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        db.close(); raise HTTPException(404, "Material not found")
    if mat["user_id"] != user_id:
        db.close(); raise HTTPException(403, "Only the owner can rename")
    body = await request.json()
    fields, params = [], []
    if "name" in body:
        fields.append("original_name = ?"); params.append(body["name"].strip()[:200])
    if "subject" in body:
        fields.append("subject = ?");        params.append(body["subject"].strip()[:100])
    if not fields:
        db.close(); return {"ok": True}
    params.append(mid)
    db.execute(f"UPDATE materials SET {', '.join(fields)} WHERE id = ?", params)
    db.commit(); db.close()
    return {"ok": True}


@app.post("/api/materials/reorder")
async def reorder_materials(request: Request, user_id: int = Depends(get_current_user)):
    """Save a new sort_order for each material. Body: [{id, sort_order}, ...]"""
    body = await request.json()  # list of {id: int, sort_order: int}
    db = get_db()
    for item in body:
        try:
            mid = int(item["id"]); order = int(item["sort_order"])
        except (KeyError, ValueError, TypeError):
            continue
        # Only update materials the user owns
        db.execute(
            "UPDATE materials SET sort_order = ? WHERE id = ? AND user_id = ?",
            (order, mid, user_id)
        )
    db.commit(); db.close()
    return {"ok": True}


@app.get("/api/materials/{mid}/content")
def get_material_content(mid: int, user_id: int = Depends(get_current_user)):
    """Raw extracted text for a material — pure DB read, NO Claude tokens.
    Used by the 'Copy for Claude' widget to bundle context for the web chat."""
    db = get_db()
    if not user_can_access(db, mid, user_id):
        db.close()
        raise HTTPException(403, "No access to this material")
    mat = db.execute(
        "SELECT original_name, subject, content FROM materials WHERE id = ?", (mid,)
    ).fetchone()
    db.close()
    if not mat:
        raise HTTPException(404, "Material not found")
    return {"name": mat["original_name"], "subject": mat["subject"], "content": mat["content"] or ""}


@app.get("/api/context/export")
def export_context(material_id: Optional[int] = None, user_id: int = Depends(get_current_user)):
    """Bundle a student's study context (material text + weak topics + missed questions +
    stats) so it can be copied into the Claude.ai web chat. Pure DB reads — NO API tokens."""
    db = get_db()
    bundle: dict = {}

    if material_id and user_can_access(db, material_id, user_id):
        mat = db.execute(
            "SELECT original_name, subject, content FROM materials WHERE id = ?", (material_id,)
        ).fetchone()
        if mat:
            bundle["material"] = {
                "name": mat["original_name"], "subject": mat["subject"],
                "content": mat["content"] or "",
            }

    # Overall quiz accuracy
    qs = dict(db.execute("SELECT COUNT(*) as t, COALESCE(SUM(is_correct),0) as c FROM (SELECT is_correct, MAX(id) FROM quiz_attempts WHERE user_id = ? GROUP BY question_id)", (user_id,)).fetchone())
    bundle["total_questions"] = qs["t"] or 0
    bundle["accuracy"] = round(((qs["c"] or 0) / max(qs["t"] or 1, 1)) * 100, 1)

    # Weakest topics by accuracy
    bundle["weak_topics"] = [dict(r) for r in db.execute(
        """SELECT topic, COUNT(*) as attempts, SUM(is_correct) as correct,
                  ROUND(CAST(SUM(is_correct) AS REAL)/COUNT(*)*100) as pct
           FROM quiz_attempts WHERE user_id = ? GROUP BY topic ORDER BY pct ASC, attempts DESC LIMIT 8""",
        (user_id,)
    ).fetchall()]

    # Recently missed questions (for targeted help)
    bundle["missed"] = [dict(r) for r in db.execute(
        """SELECT q.question, q.correct_answer, q.explanation, a.topic
           FROM quiz_attempts a JOIN quiz_questions q ON a.question_id = q.id
           WHERE a.is_correct = 0 AND a.user_id = ? ORDER BY a.attempted_at DESC LIMIT 10""",
        (user_id,)
    ).fetchall()]

    # Flashcards due today
    today = date.today().isoformat()
    bundle["due_today"] = db.execute(
        "SELECT COUNT(*) as c FROM flashcards WHERE user_id = ? AND (next_review IS NULL OR next_review <= ?)",
        (user_id, today)
    ).fetchone()["c"]

    db.close()
    return bundle


@app.post("/api/import-web")
async def import_web(request: Request):
    # Sent as text/plain to avoid a CORS preflight from the uni portal page.
    # The bookmarklet includes the active profile id in the payload (no header available cross-origin).
    raw = await request.body()
    try:
        data = json.loads(raw.decode("utf-8"))
    except Exception:
        raise HTTPException(400, "Invalid payload")

    title = (data.get("title") or "Web import").strip()[:200]
    text = (data.get("text") or "").strip()
    subject = (data.get("subject") or "Medicine").strip()
    url = (data.get("url") or "").strip()

    if len(text) < 20:
        raise HTTPException(400, "No usable text found on the page")

    db = get_db()
    # Resolve owner: explicit user_id in payload, else the first profile (single-user fallback)
    uid = data.get("user_id")
    try:
        uid = int(uid) if uid is not None else None
    except (TypeError, ValueError):
        uid = None
    if uid is None or not db.execute("SELECT 1 FROM users WHERE id = ?", (uid,)).fetchone():
        first = db.execute("SELECT id FROM users ORDER BY id LIMIT 1").fetchone()
        uid = first["id"] if first else None
    if uid is None:
        db.close()
        raise HTTPException(400, "No profile exists yet — create one in MedVault first")

    content = f"[Imported from: {url}]\n\n{text}"
    cur = db.execute(
        "INSERT INTO materials (user_id, filename, original_name, subject, content, file_type) VALUES (?,?,?,?,?,?)",
        (uid, f"web_{datetime.now().strftime('%Y%m%d_%H%M%S')}", title, subject, content, "web")
    )
    mid = cur.lastrowid
    db.execute("INSERT OR IGNORE INTO user_materials (user_id, material_id) VALUES (?,?)", (uid, mid))
    db.commit()
    db.close()
    return {"id": mid, "name": title, "chars": len(content)}


@app.post("/api/import-web-form")
async def import_web_form(request: Request):
    """Form-based import — fallback for portals whose CSP blocks fetch/XHR.
    The clipper submits a hidden <form> with the payload in a textarea field.
    Returns an HTML confirmation page (opens in a new tab)."""
    form = await request.form()
    raw = form.get("payload", "")
    try:
        data = json.loads(raw)
    except Exception:
        return HTMLResponse("<h2>Import failed</h2><p>Invalid payload.</p>", status_code=400)

    title = (data.get("title") or "Web import").strip()[:200]
    text = (data.get("text") or "").strip()
    subject = (data.get("subject") or "Medicine").strip()
    url = (data.get("url") or "").strip()

    if len(text) < 20:
        return HTMLResponse("<h2>Import failed</h2><p>No usable text found on the page.</p>", status_code=400)

    db = get_db()
    uid = data.get("user_id")
    try:
        uid = int(uid) if uid is not None else None
    except (TypeError, ValueError):
        uid = None
    if uid is None or not db.execute("SELECT 1 FROM users WHERE id = ?", (uid,)).fetchone():
        first = db.execute("SELECT id FROM users ORDER BY id LIMIT 1").fetchone()
        uid = first["id"] if first else None
    if uid is None:
        db.close()
        return HTMLResponse("<h2>Import failed</h2><p>No profile exists yet.</p>", status_code=400)

    content = f"[Imported from: {url}]\n\n{text}"
    cur = db.execute(
        "INSERT INTO materials (user_id, filename, original_name, subject, content, file_type) VALUES (?,?,?,?,?,?)",
        (uid, f"web_{datetime.now().strftime('%Y%m%d_%H%M%S')}", title, subject, content, "web")
    )
    mid = cur.lastrowid
    db.execute("INSERT OR IGNORE INTO user_materials (user_id, material_id) VALUES (?,?)", (uid, mid))
    db.commit()
    db.close()

    html = f"""<!DOCTYPE html>
<html><head><title>MedVault Import</title>
<style>body{{font-family:system-ui;display:flex;align-items:center;justify-content:center;min-height:100vh;margin:0;background:#f0fdf4}}
.card{{background:#fff;border-radius:12px;padding:2rem;box-shadow:0 4px 12px rgba(0,0,0,.1);text-align:center;max-width:400px}}
h2{{color:#16a34a;margin:0 0 .5rem}}p{{color:#555;margin:.3rem 0}}.close{{margin-top:1rem;color:#999;font-size:.9rem}}</style></head>
<body><div class="card"><h2>✅ Imported!</h2>
<p><strong>{title}</strong></p>
<p>{len(content):,} characters from {data.get("pages", 1)} slide(s)</p>
<p class="close">You can close this tab and go back to your portal.</p></div></body></html>"""
    return HTMLResponse(html)


# ─────────────────────────────────────────────────────────────────────────────
# Guide → quiz. The study-guides gallery's "Quiz" button seeds a guide's
# HAND-AUTHORED MCQs into the normal quiz engine. We find-or-create one material
# per guide and idempotently load its questions into quiz_questions, so the usual
# flow (play → answer → SRS → dashboard → mistakes) runs with NO AI call — the
# bank is pre-filled, so generate_quiz's QUIZ_MIN_UNSEEN path serves it directly.
# Questions live in guide_quizzes.json (tracked, not served), keyed by guide file.
# ─────────────────────────────────────────────────────────────────────────────
GUIDE_QUIZ_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "guide_quizzes.json")
# Copyright-noticed guides' banks live in a gitignored local overlay so they never
# reach the public repo / deploy — merged in only when present (mirrors guides.local.json).
GUIDE_QUIZ_LOCAL_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "guide_quizzes.local.json")

def _load_guide_quizzes() -> dict:
    banks: dict = {}
    for _p in (GUIDE_QUIZ_PATH, GUIDE_QUIZ_LOCAL_PATH):
        try:
            with open(_p, encoding="utf-8") as f:
                banks.update(json.load(f))
        except Exception:
            pass
    return banks

def _ensure_guide_quiz(db, user_id: int, key: str):
    """Find-or-create the material for a guide and idempotently seed its authored
    MCQs into quiz_questions. Returns (material_id, inserted_count), or (None, 0)
    when the guide has no authored quiz. Option positions are shuffled on insert
    (mirrors generate_quiz) so the correct answer isn't biased to one letter. No AI."""
    bank = _load_guide_quizzes().get(key)
    if not bank or not bank.get("questions"):
        return None, 0
    fname = f"guide:{key}"
    row = db.execute(
        "SELECT id FROM materials WHERE user_id = ? AND filename = ?", (user_id, fname)
    ).fetchone()
    if row:
        mid = row["id"]
    else:
        cur = db.execute(
            "INSERT INTO materials (user_id, filename, original_name, subject, content, file_type) VALUES (?,?,?,?,?,?)",
            (user_id, fname, (bank.get("title") or key)[:200], bank.get("subject") or "Medicine",
             bank.get("content") or "", "guide")
        )
        mid = cur.lastrowid
        db.execute("INSERT OR IGNORE INTO user_materials (user_id, material_id) VALUES (?,?)", (user_id, mid))
    # Re-seed when the authored bank changes. The insert below is idempotent-by-stem, which
    # alone would leave a user's OLD questions in place when guide_quizzes.json is updated
    # (quality/guideline fixes would never reach them). So: if the questions currently stored
    # for this guide differ from the bank's current stem-set, the bank was rewritten — wipe
    # this material's questions and re-seed fresh. Guide materials are owned by their authored
    # bank; any ad-hoc AI "More questions" added onto a guide are transient and cleared here.
    bank_stems = {(q.get("question") or "").strip() for q in bank["questions"] if (q.get("question") or "").strip()}
    db_stems = {r["question"] for r in db.execute(
        "SELECT question FROM quiz_questions WHERE material_id = ? AND user_id = ?", (mid, user_id)
    ).fetchall()}
    if db_stems and db_stems != bank_stems:
        db.execute("DELETE FROM quiz_questions WHERE material_id = ? AND user_id = ?", (mid, user_id))
        existing = set()
    else:
        existing = db_stems
    _LETTERS = ["A", "B", "C", "D"]
    inserted = 0
    for q in bank["questions"]:
        stem = (q.get("question") or "").strip()
        if not stem or stem in existing:
            continue
        opts = q.get("options") or []
        ca = (q.get("correct_answer") or "A").strip().upper()[:1]
        if len(opts) == 4 and ca in _LETTERS:
            bodies = []
            for o in opts:
                m = re.match(r'^\s*[A-Da-d]\s*[\.\)\:\-]\s*(.*)$', str(o))
                bodies.append((m.group(1) if m else str(o)).strip())
            correct_body = bodies[_LETTERS.index(ca)]
            random.shuffle(bodies)
            ca = _LETTERS[bodies.index(correct_body)]
            opts = [f"{_LETTERS[i]}. {bodies[i]}" for i in range(4)]
        db.execute(
            "INSERT INTO quiz_questions (material_id, user_id, topic, difficulty, question, options, correct_answer, explanation, related_topics) VALUES (?,?,?,?,?,?,?,?,?)",
            (mid, user_id, q.get("topic") or "General", q.get("difficulty") or "medium",
             stem, json.dumps(opts), ca, q.get("explanation") or "", json.dumps(q.get("related") or []))
        )
        existing.add(stem)
        inserted += 1
    db.commit()
    return mid, inserted

def _ensure_all_guide_quizzes(db, user_id: int) -> None:
    """Seed EVERY gallery guide's quiz for this user (idempotent) so the quiz
    dropdown lists all guides, not just ones they've opened. Cheap guard skips
    the work once fully seeded. Zero AI."""
    banks = _load_guide_quizzes()
    if not banks:
        return
    have = db.execute(
        "SELECT COUNT(*) AS c FROM materials m JOIN user_materials um ON um.material_id = m.id "
        "WHERE um.user_id = ? AND m.filename LIKE 'guide:%'", (user_id,)
    ).fetchone()["c"]
    if have >= len(banks):
        return
    for key in banks:
        try:
            _ensure_guide_quiz(db, user_id, key)
        except Exception:
            pass

@app.post("/api/guide-quiz/{key}")
def guide_quiz(key: str, user_id: int = Depends(get_current_user)):
    """Seed (once) and return the material id for a guide's hand-authored quiz.
    Called by the in-app gallery's Quiz button; zero AI, fully offline of credits."""
    db = get_db()
    try:
        mid, inserted = _ensure_guide_quiz(db, user_id, key)
        if mid is None:
            raise HTTPException(404, "No quiz is available for this guide yet.")
        total = db.execute(
            "SELECT COUNT(*) AS c FROM quiz_questions WHERE material_id = ? AND user_id = ?",
            (mid, user_id)
        ).fetchone()["c"]
        return {"material_id": mid, "inserted": inserted, "count": total}
    finally:
        db.close()

@app.get("/api/guide-quiz/stats")
def guide_quiz_stats(user_id: int = Depends(get_current_user)):
    """Per-guide quiz progress for the gallery: which guides have a quiz at all, and
    (once attempted) how many questions the user has seen plus their accuracy on the
    LATEST attempt of each. Drives the gallery's progress chip + Quiz gating. No AI."""
    banks = _load_guide_quizzes()
    db = get_db()
    try:
        out = {}
        for key, bank in banks.items():
            total = len(bank.get("questions") or [])
            if not total:
                continue
            row = db.execute(
                "SELECT id FROM materials WHERE user_id = ? AND filename = ?",
                (user_id, f"guide:{key}")
            ).fetchone()
            seen = correct = 0
            if row:
                r = db.execute(
                    """SELECT COUNT(*) AS seen, COALESCE(SUM(is_correct), 0) AS correct FROM (
                           SELECT is_correct, MAX(id) FROM quiz_attempts
                           WHERE material_id = ? AND user_id = ? GROUP BY question_id
                       )""",
                    (row["id"], user_id)
                ).fetchone()
                seen = r["seen"]; correct = r["correct"]
            out[key] = {"total": total, "seen": seen,
                        "accuracy": round(100 * correct / seen) if seen else None}
        return out
    finally:
        db.close()


# ── Slides ───────────────────────────────────────────────────────────────────

@app.post("/api/generate/slides/{mid}")
def generate_slides(mid: int, force: bool = False, user_id: int = Depends(get_current_user)):
    db = get_db()
    if not user_can_access(db, mid, user_id):
        db.close()
        raise HTTPException(403, "No access to this material")
    mat = db.execute("SELECT * FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        raise HTTPException(404, "Material not found")

    # Memory: reuse existing slides without an AI call unless regeneration is forced
    existing = db.execute("SELECT COUNT(*) as c FROM revision_slides WHERE material_id = ? AND user_id = ?", (mid, user_id)).fetchone()["c"]
    if existing > 0 and not force:
        rows = db.execute(
            "SELECT content FROM revision_slides WHERE material_id = ? AND user_id = ? ORDER BY slide_order", (mid, user_id)
        ).fetchall()
        slides = []
        for r in rows:
            try:
                slides.append(json.loads(r["content"]))
            except Exception:
                pass
        db.close()
        return {"count": existing, "existing": True, "slides": slides}

    # Load any images extracted from this material
    image_urls: list = []
    try:
        image_urls = json.loads(mat['images'] or '[]')
    except Exception:
        pass

    # Build image context block for the prompt
    image_block = ""
    image_note = ""
    if image_urls:
        img_list = "\n".join(f"  {i+1}. {u}" for i, u in enumerate(image_urls))
        image_block = f"\n\nIMAGES EXTRACTED FROM THIS MATERIAL (you may reference these in diagram slides):\n{img_list}\n"
        image_note = (
            '\nFor "diagram" slides: if one of the images above directly illustrates the topic, '
            'add "image_url": "<url>" to the object and set "nodes": [], "connections": []. '
            'The image will be displayed in the diagram panel instead of a generated SVG.'
        )

    instructions = f"""You are a medical education designer creating revision slides.

STRICT RULE: Base EVERY slide exclusively on information explicitly stated in the SOURCE STUDY MATERIAL above. Do NOT introduce assumed knowledge, background context, or facts not present in the source text. If information is not in the source, omit it.
{image_block}
Generate 16-22 slides. You have 10 TEMPLATE TYPES. For EACH section of the material, CHOOSE THE TEMPLATE THAT BEST FITS THE SHAPE OF THAT CONTENT — do not default to "concept" for everything. A great deck feels intentional: the template matches what the content actually is.

🎯 TEMPLATE SELECTION GUIDE — pick by what the content IS:
- A single core idea with a definition + supporting points → "concept"
- Numeric values, normal ranges, pressures, rates, percentages, measurements → "stat"
- A cluster of vocabulary/terms the student must know → "keyterms"
- A spatial/relational structure, anatomy layout, or flow between parts → "diagram"
- TWO things that contrast (A vs B, left vs right, before vs after) → "comparison"
- An ordered sequence of steps/events that happen in order → "process"
- A genuine memory aid present or strongly implied in the material → "mnemonic"
- A real-world/clinical application of the concept → "clinical"
- A crisp high-yield recap of a section's must-know points → "takeaway"
- The module's opening framing → "overview" (use ONCE, first slide)

VARIETY REQUIREMENT (this matters — repetitive decks are low quality):
- Use AT LEAST 6 of the 10 types across the deck.
- NEVER use the same type more than twice in a row.
- "concept" must be at most ~40% of the slides — if you're writing a 3rd concept slide in a row, a different template almost certainly fits better.
- Aim to include at least one "stat", one "keyterms" or "takeaway", and one "diagram" or "comparison" wherever the material supports them.

TYPE SCHEMAS (exact field names required):

"overview" — module intro:
{{"type":"overview","topic":"string","title":"string","subtitle":"string","key_themes":["theme1","theme2","theme3","theme4"]}}

"concept" — core definition:
{{"type":"concept","topic":"string","title":"string","definition":"one clear sentence from the material","key_points":["point from material","point from material","point from material"],"icon":"relevant emoji","color":"teal|blue|purple|green|amber|red"}}

"diagram" — spatial or relational structure:{image_note}
{{"type":"diagram","topic":"string","title":"string","diagram_type":"flow|cycle|hierarchy","nodes":[{{"id":"1","label":"short label","description":"brief"}}],"connections":[{{"from":"1","to":"2"}}],"caption":"describe what the diagram shows","clinical_pearl":"string or null"}}

"comparison" — two contrasting concepts:
{{"type":"comparison","topic":"string","title":"string","left":{{"label":"Term A","color":"blue","points":["...","...","..."]}},"right":{{"label":"Term B","color":"red","points":["...","...","..."]}},"key_difference":"one sentence"}}

"process" — step-by-step sequence:
{{"type":"process","topic":"string","title":"string","steps":[{{"number":1,"title":"step name","description":"what happens — from the source"}}],"clinical_pearl":"string or null"}}

"mnemonic" — memory aid:
{{"type":"mnemonic","topic":"string","title":"string","mnemonic":"WORD","expansion":[{{"letter":"W","meaning":"Word — full explanation"}}],"context":"Full explanation of WHAT this topic is and WHY this mnemonic applies — not just 'use this to remember'"}}

"clinical" — real-world application:
{{"type":"clinical","topic":"string","title":"string","scenario":"clinical sentence drawn from the material","question":"question about this scenario","answer":"answer based on the material","teaching_point":"key takeaway from the material"}}

"stat" — key figures / values (use when the material gives numbers, ranges, or measurements):
{{"type":"stat","topic":"string","title":"string","stats":[{{"value":"80–120","unit":"mmHg","label":"what this number is","note":"optional brief context or null"}}],"clinical_pearl":"string or null"}}
(2 to 4 stats. "value" is the number/range, "unit" is the unit or empty string.)

"keyterms" — glossary of must-know terms (use for vocabulary-dense sections):
{{"type":"keyterms","topic":"string","title":"string","terms":[{{"term":"Systole","definition":"short, clear definition from the material"}}]}}
(3 to 6 terms. Definitions one concise sentence each.)

"takeaway" — high-yield recap of a section (use as a section closer or when content is a flat list of must-knows):
{{"type":"takeaway","topic":"string","title":"string","headline":"one-sentence summary of the section","points":["crisp must-know point","crisp must-know point","crisp must-know point"]}}
(3 to 5 points, each a single sharp sentence.)

RULES:
- Every fact must be sourced from the content above — no assumed knowledge
- Use comparison for contrasting pairs found in the source
- Use mnemonic only when the source contains or strongly implies a memory aid
- Mnemonic context MUST explain the full topic concept, not just say 'use this to remember'
- Vary concept slide colors
- "stat" values, "keyterms" definitions, and "takeaway" points must all come from the material — invent nothing
- Return ONLY the JSON array, no markdown, no extra text"""

    # Sonnet, not Haiku: choosing the best-fit template per section and keeping
    # variety needs stronger reasoning than Haiku reliably provides.
    #
    # Option B — auto-sectioned generation. A girthy module (> GEN_CONTENT_CHARS)
    # would otherwise be silently truncated by gen_source_block and lose its tail,
    # "shortening everything out". Instead we split it into sections and generate
    # each at full fidelity, then stitch — still ONE click for the user. Normal-size
    # modules take the unchanged single call (no extra cost).
    full = mat['content'] or ''
    if len(full) <= GEN_CONTENT_CHARS:
        try:
            slides = parse_json_response(generate_json(mat, instructions, model=MODEL, max_tokens=8000))
            if not isinstance(slides, list):
                slides = []
        except Exception:
            slides = []
    else:
        sections = _split_material_sections(full)
        n = len(sections)
        slides, seen = [], set()
        for i, sec in enumerate(sections):
            note = (
                f"\n\n━━━ SECTIONED GENERATION (PART {i+1} OF {n}) ━━━\n"
                "The SOURCE STUDY MATERIAL above is ONLY this part of a large module. "
                "Generate slides for THIS PART ONLY and cover it COMPLETELY — every subtopic, "
                "definition, value, figure/caption and niche detail in this part. Omit nothing; do NOT summarise. "
                "Ignore the '16-22 slides' total above (that was for a whole module) — make as many slides as "
                "this part genuinely needs. Use the \"overview\" type ONLY if this is Part 1; otherwise never use it. "
                "Keep the variety, best-fit-template and source-only rules. Return ONLY the JSON array for this part."
            )
            try:
                part = parse_json_response(generate_json(mat, instructions + note, model=MODEL, max_tokens=8000, content=sec))
            except Exception:
                continue
            if not isinstance(part, list):
                continue
            for s in part:                                   # de-dup across the small section overlaps
                key = (s.get("type"), (s.get("title") or "").strip().lower())
                if key in seen:
                    continue
                seen.add(key)
                slides.append(s)
    if not slides:
        slides = [{"type":"concept","title":"Generation Error","topic":"Error","definition":"Could not generate slides.","key_points":["Please try again"],"icon":"⚠️","color":"red"}]

    db.execute("DELETE FROM revision_slides WHERE material_id = ? AND user_id = ?", (mid, user_id))
    for i, s in enumerate(slides):
        db.execute(
            "INSERT INTO revision_slides (material_id, user_id, title, content, slide_order) VALUES (?,?,?,?,?)",
            (mid, user_id, s.get("title", f"Slide {i+1}"), json.dumps(s), i)
        )
    db.commit()
    db.close()
    return {"count": len(slides), "slides": slides}


@app.get("/api/slides")
def get_slides(material_id: Optional[int] = None, user_id: int = Depends(get_current_user)):
    db = get_db()
    q = "SELECT rs.*, m.original_name FROM revision_slides rs JOIN materials m ON rs.material_id = m.id"
    rows = db.execute(q + " WHERE rs.material_id = ? AND rs.user_id = ? ORDER BY slide_order", (material_id, user_id)).fetchall() \
        if material_id else db.execute(q + " WHERE rs.user_id = ? ORDER BY rs.material_id, slide_order", (user_id,)).fetchall()
    db.close()
    result = []
    for r in rows:
        d = dict(r)
        try:
            d["content"] = json.loads(d["content"])
        except Exception:
            pass
        result.append(d)
    return result


# ── Flashcards ───────────────────────────────────────────────────────────────

@app.post("/api/generate/flashcards/{mid}")
def generate_flashcards(mid: int, force: bool = False, user_id: int = Depends(get_current_user)):
    db = get_db()
    if not user_can_access(db, mid, user_id):
        db.close()
        raise HTTPException(403, "No access to this material")
    mat = db.execute("SELECT * FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        raise HTTPException(404, "Material not found")
    # Return existing cards without an AI call unless forced
    existing = db.execute("SELECT COUNT(*) as c FROM flashcards WHERE material_id = ? AND user_id = ?", (mid, user_id)).fetchone()["c"]
    if existing > 0 and not force:
        db.close()
        return {"count": existing, "existing": True}

    # Auto-linking: gather topics from user's OTHER materials for cross-references
    other_topics = db.execute("""
        SELECT DISTINCT q.topic FROM quiz_questions q
        JOIN user_materials um ON um.material_id = q.material_id
        WHERE um.user_id = ? AND q.material_id != ? AND q.topic IS NOT NULL
        LIMIT 20
    """, (user_id, mid)).fetchall()
    related_str = ", ".join(_normalize_topic(r["topic"]) for r in other_topics) if other_topics else ""
    cross_link = ""
    if related_str:
        cross_link = f"""
CROSS-LINKING: The student is also studying these topics: {related_str}.
Where relevant, connect flashcard content to these other topics the student is learning.
For example, reference how a concept applies across subjects, or note clinical/practical links.
Add a "related" field with 1-2 related topic names when there's a genuine connection."""

    # Detect subject to tailor flashcard style
    subject = (mat["subject"] or "").lower()
    content_sample = (mat["content"] or "")[:500].lower()
    is_quantitative = any(kw in subject or kw in content_sample for kw in
        ["chemistry", "chem", "physics", "pharmacology", "biochem", "calcul", "equation", "molar", "reaction"])

    fc_chem_block = ""
    if is_quantitative:
        fc_chem_block = """
QUANTITATIVE CONTENT — include where the material supports it:
- Calculation cards: "Calculate the pH of a 0.1M solution of acetic acid (Ka = 1.8 × 10⁻⁵)"
- When a card involves a specific molecule, include its SMILES notation in a "smiles" field
- Example SMILES: ethanol = "CCO", acetic acid = "CC(=O)O", benzene = "c1ccccc1"
- Only include "smiles" when the structure is directly relevant"""
    else:
        fc_chem_block = """
SUBJECT-APPROPRIATE — match the card style to the subject, but only using concepts in the material:
- Anatomy: structure/relationship/spatial reasoning ("What structure sits between X and Y?")
- Physiology: mechanism questions, "What happens when X increases?"
- Only use a clinical or pathology scenario if the material itself introduces that condition. Do NOT invent named diseases, syndromes, or diagnoses the material doesn't cover.
- Do NOT include calculations or chemical structures unless the material explicitly covers them"""

    instructions = f"""You are a university-level educator writing CONCISE flashcards for EXAM PREPARATION.

⚠️ STAY GROUNDED IN THE SOURCE MATERIAL:
- Base EVERY card on concepts EXPLICITLY present in the SOURCE STUDY MATERIAL above.
- Do NOT introduce named conditions, syndromes, drugs, enzymes, channels, or pathways that the material does not mention. If a term does not appear in the source, do not build a card around it.
- Stay at the level of THIS course as reflected in the material — not a specialist/board-exam level. Match the vocabulary the student has actually been taught.
- A card may ask the student to APPLY or CONNECT what's in the material, but the underlying facts must come from the material, not from outside knowledge.

Create flashcards that test UNDERSTANDING and APPLICATION. Good question styles:
- Application: "Why does X cause Y?" or "What happens when..." (where X and Y are both in the material)
- Compare/contrast: "How does X differ from Y?"
- Reasoning: brief scenario requiring the student to reason from a concept in the material
{fc_chem_block}

AVOID simple definition cards like "What is X?" → "X is...". Every card should require THINKING — but only about material that was actually taught.

⚠️ BREVITY RULE — ANSWERS MUST BE SHORT:
- Maximum 2-3 sentences OR up to 4 bullet points (use • character)
- State the KEY FACT first, then brief mechanism if needed
- NO long paragraphs, NO preamble, NO padding
- Target 25-40 words per answer

Return ONLY a JSON array of 15-20 flashcards:
[{{
  "topic": "Broad topic (2-3 words max, e.g. 'Organic Chemistry', 'Cell Biology', 'Pharmacology')",
  "question": "University exam-style question (1-2 sentences max)",
  "answer": "Short answer — 2-3 sentences or ≤4 bullet points. Key fact + brief mechanism only.",
  "related": ["Related Topic 1"],
  "smiles": "SMILES string ONLY if a chemical structure is relevant, otherwise omit entirely"
}}]

IMPORTANT: The "topic" must be a BROAD subject category (2-3 words), NOT a specific question description.
{cross_link}"""

    cards = None
    if STRUCTURED_FLASHCARDS:
        # #13 structured-outputs path — schema-guaranteed JSON from Haiku.
        # Any error (API, schema, empty) falls through to the legacy path.
        try:
            cards = generate_flashcards_structured(mat, instructions, max_tokens=6000)
        except Exception as e:
            print(f"[flashcards] structured path failed ({e!r}); falling back to legacy parse")
            cards = None
    if cards is None:
        text = generate_json(mat, instructions, model=HAIKU, max_tokens=6000)
        try:
            cards = parse_json_response(text)
        except Exception:
            cards = [{"topic": "Error", "question": "Could not generate flashcards", "answer": "Please try again"}]

    db.execute("DELETE FROM flashcards WHERE material_id = ? AND user_id = ?", (mid, user_id))
    for c in cards:
        # legacy prompt emits "related"; the structured schema emits "related_topics"
        related = json.dumps(c.get("related") or c.get("related_topics") or [])
        smiles = c.get("smiles") or None
        db.execute(
            "INSERT INTO flashcards (material_id, user_id, topic, question, answer, related_topics, smiles) VALUES (?,?,?,?,?,?,?)",
            (mid, user_id, c.get("topic", "General"), c.get("question", ""), c.get("answer", ""), related, smiles)
        )
    db.commit()
    db.close()
    return {"count": len(cards)}


@app.get("/api/flashcards")
def get_flashcards(material_id: Optional[int] = None, adaptive: bool = False, due_only: bool = False,
                   user_id: int = Depends(get_current_user)):
    db = get_db()
    today = date.today().isoformat()
    conditions, params = ["f.user_id = ?"], [user_id]
    if material_id:
        conditions.append("f.material_id = ?")
        params.append(material_id)
    if due_only:
        conditions.append("(f.next_review IS NULL OR f.next_review <= ?)")
        params.append(today)
    where = ("WHERE " + " AND ".join(conditions)) if conditions else ""
    if due_only:
        order = "COALESCE(f.next_review, '1970-01-01') ASC"  # most overdue first
    elif adaptive:
        order = "CASE WHEN f.times_seen = 0 THEN 0 ELSE CAST(f.times_correct AS REAL)/f.times_seen END ASC"
    else:
        order = "f.id"
    base = "SELECT f.*, m.original_name FROM flashcards f JOIN materials m ON f.material_id = m.id"
    rows = db.execute(f"{base} {where} ORDER BY {order}", params).fetchall()
    db.close()
    return [dict(r) for r in rows]


@app.post("/api/flashcards/{cid}/result")
async def flashcard_result(cid: int, request: Request, user_id: int = Depends(get_current_user)):
    body = await request.json()
    correct = body.get("correct", False)
    db = get_db()
    card = db.execute("SELECT * FROM flashcards WHERE id = ? AND user_id = ?", (cid, user_id)).fetchone()
    if not card:
        db.close()
        raise HTTPException(404, "Card not found")
    # SM-2 spaced repetition (shared scheduler)
    new_interval, new_ease, new_count, next_review = sm2_schedule(
        card["ease_factor"], card["srs_interval"], card["review_count"], correct)
    db.execute(
        """UPDATE flashcards SET
           times_seen=times_seen+1, times_correct=times_correct+?,
           last_seen=?, srs_interval=?, ease_factor=?, review_count=?, next_review=?
           WHERE id=?""",
        (1 if correct else 0, datetime.now().isoformat(),
         new_interval, round(new_ease, 4), new_count, next_review, cid)
    )
    # Log every review so the dashboard activity chart includes flashcard study
    db.execute(
        "INSERT INTO flashcard_log (flashcard_id, material_id, user_id, topic, correct) VALUES (?,?,?,?,?)",
        (cid, card["material_id"], user_id, card["topic"], 1 if correct else 0)
    )
    db.commit()
    db.close()
    return {"ok": True, "next_review": next_review, "interval_days": new_interval}


@app.get("/api/srs/stats")
def get_srs_stats(user_id: int = Depends(get_current_user)):
    db = get_db()
    today    = date.today().isoformat()
    week_end = (date.today() + timedelta(days=7)).isoformat()
    due      = db.execute("SELECT COUNT(*) as c FROM flashcards WHERE user_id = ? AND (next_review IS NULL OR next_review <= ?)", (user_id, today)).fetchone()["c"]
    new_c    = db.execute("SELECT COUNT(*) as c FROM flashcards WHERE user_id = ? AND review_count = 0", (user_id,)).fetchone()["c"]
    mature   = db.execute("SELECT COUNT(*) as c FROM flashcards WHERE user_id = ? AND srs_interval > 21", (user_id,)).fetchone()["c"]
    # "Coming up" = scheduled strictly after today, within 7 days (not today's cards)
    upcoming = db.execute(
        "SELECT COUNT(*) as c FROM flashcards WHERE user_id = ? AND next_review > ? AND next_review <= ?",
        (user_id, today, week_end)
    ).fetchone()["c"]
    db.close()
    return {"due_today": due, "new_cards": new_c, "mature": mature, "upcoming_week": upcoming}


# ── GAMSAT Writing coach ──────────────────────────────────────────────────────

@app.get("/api/writing/stimulus")
def writing_get_stimulus(task: str = "A", generate: bool = False,
                         user_id: int = Depends(get_current_user)):
    """A GAMSAT-style stimulus. By default serves from the hand-written seed bank
    (zero API cost); pass ?generate=true for a fresh AI-generated one."""
    task = "B" if str(task).upper() == "B" else "A"
    if generate:
        stim = writing_stimulus(task)
        if stim.get("theme") and stim.get("quotes"):
            return stim
        # fall through to seed bank on a malformed generation
    pool = [s for s in WRITING_SEED_STIMULI if s["task"] == task]
    stim = dict(random.choice(pool))
    stim["instruction"] = WRITING_TASK_INSTRUCTIONS[task]
    return stim


@app.post("/api/writing/assess")
async def writing_assess_route(request: Request, user_id: int = Depends(get_current_user)):
    """Grade an essay: objective grammar audit → standardised rubric assessment →
    persist essay + errors → lazily upsert one SM-2 drill card per error category."""
    body = await request.json()
    essay = (body.get("essay") or "").strip()
    words = essay.split()
    if len(words) < 30:
        raise HTTPException(400, "Essay too short to assess — write at least ~30 words")
    task = "B" if str(body.get("task", "A")).upper() == "B" else "A"
    stimulus = {
        "task": task,
        "theme": (body.get("theme") or "").strip(),
        "quotes": [q for q in (body.get("quotes") or []) if isinstance(q, str)],
        "instruction": body.get("instruction") or WRITING_TASK_INSTRUCTIONS[task],
    }

    # 1. Objective grammar pass (the anchor)
    errors = writing_analyze_language(essay)
    by_category = defaultdict(int)
    for e in errors:
        by_category[e["category"]] += 1
    error_stats = {
        "word_count": len(words),
        "error_count": len(errors),
        "per_100_words": round(len(errors) * 100.0 / len(words), 2),
        "by_category": dict(by_category),
        "language_band_cap": _language_band_cap(len(words), len(errors)),
    }

    # 2. Standardised rubric assessment (cached rubric prefix, evidence-grounded)
    assessment = writing_assess(stimulus, essay, error_stats)

    # 3. Persist essay + errors; upsert one SM-2 card per category
    db = get_db()
    cur = db.execute(
        "INSERT INTO writing_essays (user_id, task, theme, essay_text, overall_band, assessment_json) VALUES (?,?,?,?,?,?)",
        (user_id, task, stimulus["theme"], essay, assessment["overall_band"],
         json.dumps({"assessment": assessment, "errors": errors, "error_stats": error_stats})))
    essay_id = cur.lastrowid
    for e in errors:
        db.execute(
            "INSERT INTO writing_errors (user_id, essay_id, original, corrected, category, explanation) VALUES (?,?,?,?,?,?)",
            (user_id, essay_id, e["original"], e["corrected"], e["category"], e["explanation"]))
    today = date.today().isoformat()
    for cat, cnt in by_category.items():
        db.execute(
            """INSERT INTO writing_cards (user_id, category, error_count, next_review)
               VALUES (?,?,?,?)
               ON CONFLICT(user_id, category) DO UPDATE SET error_count = error_count + ?""",
            (user_id, cat, cnt, today, cnt))
    db.commit()
    db.close()

    return {"essay_id": essay_id, "assessment": assessment,
            "errors": errors, "error_stats": error_stats}


@app.get("/api/writing/drills")
def writing_get_drills(user_id: int = Depends(get_current_user)):
    """Due SM-2 category cards → one AI-generated drill per category, seeded from
    the user's own past mistakes in that category (single API call for the batch)."""
    db = get_db()
    today = date.today().isoformat()
    cards = db.execute(
        """SELECT * FROM writing_cards WHERE user_id = ?
           AND (next_review IS NULL OR next_review <= ?)
           ORDER BY COALESCE(next_review, '1970-01-01') ASC LIMIT 6""",
        (user_id, today)).fetchall()
    if not cards:
        db.close()
        return {"drills": []}
    card_by_cat, examples = {}, {}
    for c in cards:
        card_by_cat[c["category"]] = c["id"]
        rows = db.execute(
            "SELECT original, corrected FROM writing_errors WHERE user_id = ? AND category = ? ORDER BY id DESC LIMIT 3",
            (user_id, c["category"])).fetchall()
        examples[c["category"]] = [dict(r) for r in rows]
    db.close()
    drills = writing_drills(list(card_by_cat.keys()), examples)
    for d in drills:
        d["card_id"] = card_by_cat.get(d["category"])
    return {"drills": [d for d in drills if d.get("card_id")]}


@app.post("/api/writing/drills/{cid}/result")
async def writing_drill_result(cid: int, request: Request,
                               user_id: int = Depends(get_current_user)):
    """SM-2 step on a writing category card — mirrors /api/flashcards/{cid}/result."""
    body = await request.json()
    correct = body.get("correct", False)
    db = get_db()
    card = db.execute("SELECT * FROM writing_cards WHERE id = ? AND user_id = ?",
                      (cid, user_id)).fetchone()
    if not card:
        db.close()
        raise HTTPException(404, "Drill card not found")
    new_interval, new_ease, new_count, next_review = sm2_schedule(
        card["ease_factor"], card["srs_interval"], card["review_count"], correct)
    db.execute(
        """UPDATE writing_cards SET last_seen = ?, srs_interval = ?, ease_factor = ?,
           review_count = ?, next_review = ? WHERE id = ?""",
        (datetime.now().isoformat(), new_interval, round(new_ease, 4),
         new_count, next_review, cid))
    db.execute(
        "INSERT INTO writing_reviews (user_id, card_id, quality, was_correct) VALUES (?,?,?,?)",
        (user_id, cid, 4 if correct else 1, 1 if correct else 0))
    db.commit()
    db.close()
    return {"ok": True, "next_review": next_review, "interval_days": new_interval}


@app.get("/api/writing/stats")
def writing_get_stats(user_id: int = Depends(get_current_user)):
    """Band trend (now with per-criterion bands per essay), per-criterion
    averages, an overall-band improvement estimate, error counts by category,
    and due-drill count. Purely read-side — computed from stored assessments."""
    db = get_db()
    today = date.today().isoformat()
    rows = db.execute(
        "SELECT id, created_at, task, theme, overall_band, assessment_json FROM writing_essays WHERE user_id = ? ORDER BY created_at ASC",
        (user_id,)).fetchall()[-30:]
    trend = []
    crit_bands = defaultdict(list)  # criterion name → [band, …] in essay order
    for r in rows:
        item = {k: r[k] for k in ("id", "created_at", "task", "theme", "overall_band")}
        item["criteria"] = {}
        try:  # assessment_json may be missing/malformed on old rows — never fail stats
            crits = (json.loads(r["assessment_json"] or "{}")
                     .get("assessment") or {}).get("criteria") or []
            for c in crits:
                name, band = c.get("name"), c.get("band")
                if name and isinstance(band, (int, float)):
                    item["criteria"][name] = band
                    crit_bands[name].append(band)
        except (json.JSONDecodeError, AttributeError, TypeError):
            pass
        trend.append(item)

    criteria_averages = [
        {"name": name, "avg": round(sum(bands) / len(bands), 2), "count": len(bands)}
        for name, bands in crit_bands.items()
    ]

    # Improvement over overall_band: mean of the last third of essays minus the
    # first third (chronological), plus a least-squares slope (band per essay).
    bands = [t["overall_band"] for t in trend if t["overall_band"] is not None]
    n = len(bands)
    improvement = {"sample_size": n, "delta": None, "slope": None}
    if n >= 2:
        third = max(1, n // 3)
        improvement["delta"] = round(
            sum(bands[-third:]) / third - sum(bands[:third]) / third, 2)
        xm, ym = (n - 1) / 2.0, sum(bands) / n
        denom = sum((i - xm) ** 2 for i in range(n))
        if denom:
            improvement["slope"] = round(
                sum((i - xm) * (b - ym) for i, b in enumerate(bands)) / denom, 3)

    errors = {r["category"]: r["c"] for r in db.execute(
        "SELECT category, COUNT(*) AS c FROM writing_errors WHERE user_id = ? GROUP BY category ORDER BY c DESC",
        (user_id,)).fetchall()}
    due = db.execute(
        "SELECT COUNT(*) AS c FROM writing_cards WHERE user_id = ? AND (next_review IS NULL OR next_review <= ?)",
        (user_id, today)).fetchone()["c"]
    essays = db.execute("SELECT COUNT(*) AS c FROM writing_essays WHERE user_id = ?",
                        (user_id,)).fetchone()["c"]
    db.close()
    return {"band_trend": trend, "errors_by_category": errors,
            "due_drills": due, "essay_count": essays,
            "criteria_averages": criteria_averages, "improvement": improvement}


@app.post("/api/writing/coach")
def writing_coach(user_id: int = Depends(get_current_user)):
    """OPT-IN AI deep-dive (the only credit-spending writing route). Reads the
    student's recent essays + recurring errors + criterion averages and returns
    bespoke, prioritised coaching. Everything else in the Writing coach is zero-AI."""
    db = get_db()
    essays = db.execute(
        "SELECT task, theme, overall_band, essay_text FROM writing_essays "
        "WHERE user_id = ? ORDER BY created_at DESC LIMIT 4", (user_id,)).fetchall()
    if not essays:
        db.close()
        raise HTTPException(400, "Write at least one essay first")
    err_rows = db.execute(
        "SELECT category, COUNT(*) AS c FROM writing_errors WHERE user_id = ? "
        "GROUP BY category ORDER BY c DESC", (user_id,)).fetchall()
    crit_rows = db.execute(
        "SELECT assessment_json FROM writing_essays WHERE user_id = ? "
        "ORDER BY created_at DESC LIMIT 8", (user_id,)).fetchall()
    db.close()

    crit_bands = defaultdict(list)
    for r in crit_rows:
        try:
            for c in (json.loads(r["assessment_json"] or "{}").get("assessment") or {}).get("criteria") or []:
                if c.get("name") and isinstance(c.get("band"), (int, float)):
                    crit_bands[c["name"]].append(c["band"])
        except (json.JSONDecodeError, AttributeError, TypeError):
            pass
    crit_summary = {k: round(sum(v) / len(v), 2) for k, v in crit_bands.items() if v}
    errors_summary = {r["category"]: r["c"] for r in err_rows}
    # Cap essay text sent to keep the single call cheap.
    excerpts = "\n\n".join(
        f"[Task {e['task']} · {e['theme']} · band {e['overall_band']}]\n{(e['essay_text'] or '')[:1200]}"
        for e in essays)

    prompt = (
        "You are a GAMSAT writing coach. Below is one student's recent work and their "
        "measured patterns. Give SPECIFIC, prioritised coaching that only makes sense for "
        "THIS student — reference what they actually do, not generic advice. Focus on the "
        "highest-leverage changes.\n\n"
        f"CRITERION AVERAGES (band 1-6): {json.dumps(crit_summary, sort_keys=True)}\n"
        f"RECURRING LANGUAGE ERRORS (by count): {json.dumps(errors_summary, sort_keys=True)}\n\n"
        f"RECENT ESSAY EXCERPTS:\n{excerpts}\n\n"
        "Return ONLY this JSON object:\n"
        '{"summary": "<2-3 sentences naming the single most important thing to work on and why>", '
        '"tips": ["<specific, actionable coaching point tied to their writing>", "... 3-5 total"]}'
    )
    raw = _write_call(None, prompt, effort="high", max_tokens=2000)
    data = parse_json_response(raw)
    if isinstance(data, list):
        data = data[0] if data and isinstance(data[0], dict) else {}
    if not isinstance(data, dict):
        data = {}
    tips = [str(t).strip() for t in (data.get("tips") or []) if str(t).strip()]
    return {"summary": (data.get("summary") or "").strip(), "tips": tips[:6]}


@app.get("/api/writing/essays/{essay_id}")
def writing_get_essay(essay_id: int, user_id: int = Depends(get_current_user)):
    """One stored essay with its parsed assessment — powers the Progress archive."""
    db = get_db()
    row = db.execute("SELECT * FROM writing_essays WHERE id = ? AND user_id = ?",
                     (essay_id, user_id)).fetchone()
    db.close()
    if not row:
        raise HTTPException(404, "Essay not found")
    try:
        parsed = json.loads(row["assessment_json"] or "{}")
        if not isinstance(parsed, dict):
            parsed = {}
    except json.JSONDecodeError:
        parsed = {}
    return {"id": row["id"], "created_at": row["created_at"], "task": row["task"],
            "theme": row["theme"], "essay_text": row["essay_text"],
            "assessment": parsed.get("assessment") or {},
            "errors": parsed.get("errors") or [],
            "error_stats": parsed.get("error_stats") or {}}


@app.post("/api/writing/calibrate")
def writing_calibrate_route(n: int = 3, user_id: int = Depends(get_current_user)):
    """Dev route: consistency self-test. Grades the fixed benchmark essay N (≤3)
    times and reports per-criterion band spread. Costs credit — use sparingly."""
    return writing_calibrate(n)


# ── Word Bank ────────────────────────────────────────────────────────────────
# User-added vocabulary, merged client-side with the built-in seed list.
# No AI involved — words are supplied by the history-scan import flow.

WORD_BANK_CATEGORIES = {"vocab", "philosophy", "grammar", "synonyms", "idioms", "archaic"}


@app.get("/api/wordbank")
def wordbank_list(user_id: int = Depends(get_current_user)):
    """All of this user's saved words, newest first."""
    db = get_db()
    rows = db.execute(
        "SELECT id, category, word, pos, definition, example, source, created_at "
        "FROM word_bank WHERE user_id = ? ORDER BY id DESC",
        (user_id,)).fetchall()
    db.close()
    return {"words": [dict(r) for r in rows]}


@app.post("/api/wordbank")
async def wordbank_add(request: Request, user_id: int = Depends(get_current_user)):
    """Bulk-add words. Body: {words: [{category, word, pos, definition, example}]}.
    INSERT OR IGNORE on UNIQUE(user_id, word) makes re-imports idempotent."""
    body = await request.json()
    items = body.get("words")
    if not isinstance(items, list) or not items:
        raise HTTPException(400, "Provide a non-empty 'words' array")
    db = get_db()
    added = 0
    for it in items:
        if not isinstance(it, dict):
            continue
        word = (it.get("word") or "").strip()
        if not word:
            continue
        cat = (it.get("category") or "vocab").strip().lower()
        if cat not in WORD_BANK_CATEGORIES:
            cat = "vocab"
        cur = db.execute(
            "INSERT OR IGNORE INTO word_bank (user_id, category, word, pos, definition, example, source) "
            "VALUES (?, ?, ?, ?, ?, ?, ?)",
            (user_id, cat, word, (it.get("pos") or "").strip(),
             (it.get("definition") or "").strip(), (it.get("example") or "").strip(),
             (it.get("source") or "import").strip()))
        added += cur.rowcount
    db.commit()
    db.close()
    return {"added": added, "requested": len(items)}


@app.delete("/api/wordbank/{word_id}")
def wordbank_delete(word_id: int, user_id: int = Depends(get_current_user)):
    """Remove one saved word (scoped to the requesting user)."""
    db = get_db()
    cur = db.execute("DELETE FROM word_bank WHERE id = ? AND user_id = ?", (word_id, user_id))
    db.commit()
    db.close()
    if cur.rowcount == 0:
        raise HTTPException(404, "Word not found")
    return {"deleted": word_id}


# ── Russian learning module (zero-API: seed deck + SM-2 drilling + browser TTS) ──
# The 5-phase curriculum prose/resources live client-side (RU_CURRICULUM in app.js).
# The backend only stores vocab, schedules reviews with the shared sm2_schedule(),
# and tracks per-phase progress. There are NO AI calls anywhere in this module.
RUSSIAN_CATEGORIES = {
    "letters", "greetings", "intro", "politeness", "numbers",
    "food", "shopping", "directions", "lodging", "grammar",
    "questions", "conversation", "verbs", "vocab",
    # A1 expansion
    "pronouns", "family", "colors", "time", "adjectives",
    "places", "weather", "body", "home",
}

# (phase, category, cyrillic, translit, english, example, note)
# phase 0 = Cyrillic alphabet · 1 = survival · 2 = travel · 3 = conversation.
RUSSIAN_SEED = [
    # ── Phase 0 — Cyrillic alphabet (all 33 letters, grouped by difficulty) ──
    (0, "letters", "А а", "a",    "Vowel — 'a' as in 'father'",            "мама",      "look-alike"),
    (0, "letters", "Е е", "ye",   "'ye' as in 'yes' (stressed)",           "нет",       "new"),
    (0, "letters", "К к", "k",    "'k' as in 'kite'",                      "кофе",      "look-alike"),
    (0, "letters", "М м", "m",    "'m' as in 'map'",                       "метро",     "look-alike"),
    (0, "letters", "О о", "o",    "'o' as in 'more' when stressed",        "окно",      "look-alike"),
    (0, "letters", "Т т", "t",    "'t' as in 'top'",                       "такси",     "look-alike"),
    (0, "letters", "В в", "v",    "'v' as in 'van' — NOT 'b'",             "вода",      "false-friend"),
    (0, "letters", "Н н", "n",    "'n' as in 'net' — NOT 'h'",             "нос",       "false-friend"),
    (0, "letters", "Р р", "r",    "rolled 'r' — NOT 'p'",                  "ресторан",  "false-friend"),
    (0, "letters", "С с", "s",    "'s' as in 'sun' — NOT 'c'",             "спасибо",   "false-friend"),
    (0, "letters", "У у", "oo",   "'oo' as in 'boot' — NOT 'y'",           "утро",      "false-friend"),
    (0, "letters", "Х х", "kh",   "'ch' as in Scottish 'loch'",           "хлеб",      "false-friend"),
    (0, "letters", "Б б", "b",    "'b' as in 'bat'",                       "банк",      "new"),
    (0, "letters", "Г г", "g",    "'g' as in 'go'",                        "город",     "new"),
    (0, "letters", "Д д", "d",    "'d' as in 'dog'",                       "да",        "new"),
    (0, "letters", "Ё ё", "yo",   "'yo' as in 'yonder' (always stressed)", "ёлка",      "new"),
    (0, "letters", "Ж ж", "zh",   "'s' as in 'measure'",                   "жена",      "new"),
    (0, "letters", "З з", "z",    "'z' as in 'zoo'",                       "зонт",      "new"),
    (0, "letters", "И и", "ee",   "'ee' as in 'see'",                      "икра",      "new"),
    (0, "letters", "Й й", "y",    "short 'y' as in 'boy'",                 "чай",       "new"),
    (0, "letters", "Л л", "l",    "'l' as in 'lamp'",                      "лампа",     "new"),
    (0, "letters", "П п", "p",    "'p' as in 'pen'",                       "паспорт",   "new"),
    (0, "letters", "Ф ф", "f",    "'f' as in 'fun'",                       "кофе",      "new"),
    (0, "letters", "Ц ц", "ts",   "'ts' as in 'cats'",                     "центр",     "new"),
    (0, "letters", "Ч ч", "ch",   "'ch' as in 'chair'",                    "чай",       "new"),
    (0, "letters", "Ш ш", "sh",   "hard 'sh' as in 'shut'",                "школа",     "new"),
    (0, "letters", "Щ щ", "shch", "soft 'sh' — 'fresh sheet'",             "борщ",      "new"),
    (0, "letters", "Э э", "e",    "'e' as in 'met'",                       "это",       "new"),
    (0, "letters", "Ю ю", "yu",   "'yu' as in 'universe'",                 "юг",        "new"),
    (0, "letters", "Я я", "ya",   "'ya' as in 'yard'",                     "я",         "new"),
    (0, "letters", "Ъ ъ", "—",    "Hard sign — silent; separates sounds",  "объект",    "sign"),
    (0, "letters", "Ы ы", "y",    "hard 'i' — no English equivalent",      "сыр",       "sign"),
    (0, "letters", "Ь ь", "—",    "Soft sign — silent; softens the letter","соль",      "sign"),

    # ── Phase 1 — Survival basics ──
    (1, "greetings",  "Привет",         "privét",            "Hi (informal)",            "", "Informal — friends & peers"),
    (1, "greetings",  "Здравствуйте",   "zdrávstvuyte",      "Hello (formal)",           "", "The first 'в' is silent: 'zdrastvuyte'"),
    (1, "greetings",  "Пока",           "poká",              "Bye (informal)",           "", ""),
    (1, "greetings",  "До свидания",    "do svidániya",      "Goodbye (formal)",         "", ""),
    (1, "greetings",  "Доброе утро",    "dóbroye útro",      "Good morning",             "", ""),
    (1, "greetings",  "Добрый день",    "dóbryy den'",       "Good afternoon",           "", ""),
    (1, "intro",      "Меня зовут…",    "menyá zovút…",      "My name is…",              "", "lit. 'me they-call'"),
    (1, "intro",      "Как вас зовут?", "kak vas zovút?",    "What's your name? (formal)","", ""),
    (1, "intro",      "Я из…",          "ya iz…",            "I'm from…",                "", "Я из Австралии = I'm from Australia"),
    (1, "intro",      "Очень приятно",  "óchen' priyátno",   "Nice to meet you",         "", ""),
    (1, "intro",      "Я не понимаю",   "ya ne ponimáyu",    "I don't understand",       "", ""),
    (1, "intro",      "Вы говорите по-английски?", "vy govoríte po-anglíyski?", "Do you speak English?", "", ""),
    (1, "politeness", "Спасибо",        "spasíbo",           "Thank you",                "", ""),
    (1, "politeness", "Большое спасибо","bol'shóye spasíbo", "Thank you very much",      "", ""),
    (1, "politeness", "Пожалуйста",     "pozháluysta",       "Please / You're welcome",  "", ""),
    (1, "politeness", "Извините",       "izviníte",          "Excuse me / Sorry",        "", ""),
    (1, "politeness", "Да",             "da",                "Yes",                      "", ""),
    (1, "politeness", "Нет",            "net",               "No",                       "", ""),
    (1, "numbers",    "ноль",   "nol'",         "0",  "", ""),
    (1, "numbers",    "один",   "odín",         "1",  "", ""),
    (1, "numbers",    "два",    "dva",          "2",  "", ""),
    (1, "numbers",    "три",    "tri",          "3",  "", ""),
    (1, "numbers",    "четыре", "chetýre",      "4",  "", ""),
    (1, "numbers",    "пять",   "pyat'",        "5",  "", ""),
    (1, "numbers",    "шесть",  "shest'",       "6",  "", ""),
    (1, "numbers",    "семь",   "sem'",         "7",  "", ""),
    (1, "numbers",    "восемь", "vósem'",       "8",  "", ""),
    (1, "numbers",    "девять", "dévyat'",      "9",  "", ""),
    (1, "numbers",    "десять", "désyat'",      "10", "", ""),
    (1, "numbers",    "двадцать","dvádtsat'",   "20", "", ""),
    (1, "numbers",    "сто",    "sto",          "100","", "vowel reduction: unstressed 'о' → 'a'"),

    # ── Phase 2 — Travel & everyday ──
    (2, "food",       "Я хочу…",           "ya khochú…",         "I want…",                "Я хочу кофе — I want a coffee", ""),
    (2, "food",       "Меню, пожалуйста",  "menyú, pozháluysta", "The menu, please",       "", ""),
    (2, "food",       "Счёт, пожалуйста",  "schyot, pozháluysta","The check, please",      "", ""),
    (2, "food",       "вода",              "vodá",               "water",                  "", ""),
    (2, "food",       "кофе",              "kófe",               "coffee",                 "", ""),
    (2, "food",       "хлеб",              "khleb",              "bread",                  "", ""),
    (2, "food",       "Вкусно!",           "vkúsno",             "Delicious!",             "", ""),
    (2, "shopping",   "Сколько стоит?",    "skól'ko stóit?",     "How much is it?",        "", ""),
    (2, "shopping",   "У вас есть…?",      "u vas yest'…?",      "Do you have…?",          "", ""),
    (2, "shopping",   "Это дорого",        "éto dórogo",         "That's expensive",       "", ""),
    (2, "shopping",   "карта",             "kárta",              "card (payment)",         "", ""),
    (2, "directions", "Где…?",             "gde…?",              "Where is…?",             "Где метро? — Where's the metro?", ""),
    (2, "directions", "налево",            "nalévo",             "(to the) left",          "", ""),
    (2, "directions", "направо",           "naprávo",            "(to the) right",         "", ""),
    (2, "directions", "прямо",             "pryámo",             "straight ahead",         "", ""),
    (2, "directions", "Где туалет?",       "gde tualét?",        "Where is the toilet?",   "", ""),
    (2, "directions", "метро",             "metró",              "metro / subway",         "", ""),
    (2, "directions", "такси",             "taksí",              "taxi",                   "", ""),
    (2, "lodging",    "У меня бронь",      "u menyá bron'",      "I have a reservation",   "", ""),
    (2, "lodging",    "ключ",              "klyuch",             "key",                    "", ""),
    (2, "lodging",    "паспорт",           "pásport",            "passport",               "", ""),
    (2, "grammar",    "стол · книга · окно","stol · kniga · okno","Noun gender by ending",  "", "-consonant = masc, -а/-я = fem, -о/-е = neut"),
    (2, "grammar",    "я хочу · ты хочешь","ya khochú · ty khóchesh'","Present tense (хотеть, to want)","", "я хочу, ты хочешь, он хочет, мы хотим"),
    (2, "grammar",    "Я хочу воду",       "ya khochú vódu",     "Accusative = the object","", "fem -а → -у: вода → воду"),
    (2, "grammar",    "в отеле",           "v otéle",            "Prepositional = location","", "after в/на, ending often -е: 'in the hotel'"),

    # ── Phase 3 — Conversation ──
    (3, "questions",    "Что",          "shto",           "What",              "", "'ч' is pronounced 'sh' here"),
    (3, "questions",    "Где",          "gde",            "Where",             "", ""),
    (3, "questions",    "Когда",        "kogdá",          "When",              "", ""),
    (3, "questions",    "Почему",       "pochemú",        "Why",               "", ""),
    (3, "questions",    "Как",          "kak",            "How",               "", ""),
    (3, "questions",    "Кто",          "kto",            "Who",               "", ""),
    (3, "conversation", "Как дела?",    "kak delá?",      "How are you?",      "", ""),
    (3, "conversation", "Хорошо, спасибо","khoroshó, spasíbo","Good, thanks",   "", ""),
    (3, "conversation", "Мне нравится", "mne nrávitsya",  "I like it",         "", "lit. 'to-me it-is-pleasing'"),
    (3, "conversation", "Я буду…",      "ya búdu…",       "I will…",           "Я буду чай — I'll have tea", ""),
    (3, "grammar",      "Я был · Я была","ya byl · ya bylá","Past tense",        "", "agrees with GENDER/number, not person: был/была/было/были"),
    (3, "grammar",      "нет времени",  "net vrémeni",    "Genitive = 'of' / after нет","У меня нет времени — I have no time", "possession & negation"),
    (3, "grammar",      "мне",          "mne",            "Dative = to/for someone","Мне нравится — I like it", ""),
    (3, "grammar",      "кофе с молоком","kófe s molokóm","Instrumental = 'with'","", "с + instrumental: 'coffee with milk'"),

    # ══ A1 vocabulary expansion ══════════════════════════════════════════════
    # ── Phase 1 — pronouns ──
    (1, "pronouns", "я",   "ya",    "I",                "", ""),
    (1, "pronouns", "ты",  "ty",    "you (informal)",   "", "singular, to a friend"),
    (1, "pronouns", "он",  "on",    "he",               "", ""),
    (1, "pronouns", "она", "oná",   "she",              "", ""),
    (1, "pronouns", "оно", "onó",   "it",               "", "neuter"),
    (1, "pronouns", "мы",  "my",    "we",               "", ""),
    (1, "pronouns", "вы",  "vy",    "you (formal/plural)","", ""),
    (1, "pronouns", "они", "oní",   "they",             "", ""),
    (1, "pronouns", "это", "éto",   "this / it is",     "Это книга — This is a book", ""),
    (1, "pronouns", "мой", "moy",   "my (masc.)",       "мой друг — my friend", "моя (fem.), моё (neut.)"),
    (1, "pronouns", "твой","tvoy",  "your (informal)",  "", ""),

    # ── Phase 1 — family ──
    (1, "family", "семья",   "sem'yá",    "family",       "", ""),
    (1, "family", "мама",    "máma",      "mum",          "", ""),
    (1, "family", "папа",    "pápa",      "dad",          "", ""),
    (1, "family", "мать",    "mat'",      "mother",       "", ""),
    (1, "family", "отец",    "otéts",     "father",       "", ""),
    (1, "family", "брат",    "brat",      "brother",      "", ""),
    (1, "family", "сестра",  "sestrá",    "sister",       "", ""),
    (1, "family", "сын",     "syn",       "son",          "", ""),
    (1, "family", "дочь",    "doch'",     "daughter",     "", ""),
    (1, "family", "бабушка", "bábushka",  "grandmother",  "", ""),
    (1, "family", "дедушка", "dédushka",  "grandfather",  "", ""),
    (1, "family", "муж",     "muzh",      "husband",      "", ""),
    (1, "family", "жена",    "zhená",     "wife",         "", ""),
    (1, "family", "друг",    "drug",      "friend (male)","", "подруга = female friend"),
    (1, "family", "ребёнок", "rebyónok",  "child",        "", "дети = children"),

    # ── Phase 1 — numbers (11–19, tens, thousand) ──
    (1, "numbers", "одиннадцать",  "odínnadtsat'",    "11", "", ""),
    (1, "numbers", "двенадцать",   "dvenádtsat'",     "12", "", ""),
    (1, "numbers", "тринадцать",   "trinádtsat'",     "13", "", ""),
    (1, "numbers", "четырнадцать", "chetýrnadtsat'",  "14", "", ""),
    (1, "numbers", "пятнадцать",   "pyatnádtsat'",    "15", "", ""),
    (1, "numbers", "шестнадцать",  "shestnádtsat'",   "16", "", ""),
    (1, "numbers", "семнадцать",   "semnádtsat'",     "17", "", ""),
    (1, "numbers", "восемнадцать", "vosemnádtsat'",   "18", "", ""),
    (1, "numbers", "девятнадцать", "devyatnádtsat'",  "19", "", ""),
    (1, "numbers", "тридцать",     "trídtsat'",       "30", "", ""),
    (1, "numbers", "сорок",        "sórok",           "40", "", ""),
    (1, "numbers", "пятьдесят",    "pyat'desyát",     "50", "", ""),
    (1, "numbers", "шестьдесят",   "shest'desyát",    "60", "", ""),
    (1, "numbers", "семьдесят",    "sém'desyat",      "70", "", ""),
    (1, "numbers", "восемьдесят",  "vósem'desyat",    "80", "", ""),
    (1, "numbers", "девяносто",    "devyanósto",      "90", "", ""),
    (1, "numbers", "тысяча",       "týsyacha",        "1000","", ""),

    # ── Phase 1 — colours ──
    (1, "colors", "красный",    "krásnyy",     "red",    "", ""),
    (1, "colors", "синий",      "síniy",       "blue",   "", ""),
    (1, "colors", "зелёный",    "zelyónyy",    "green",  "", ""),
    (1, "colors", "жёлтый",     "zhyóltyy",    "yellow", "", ""),
    (1, "colors", "белый",      "bélyy",       "white",  "", ""),
    (1, "colors", "чёрный",     "chyórnyy",    "black",  "", ""),
    (1, "colors", "серый",      "séryy",       "grey",   "", ""),
    (1, "colors", "оранжевый",  "oránzhevyy",  "orange", "", ""),
    (1, "colors", "розовый",    "rózovyy",     "pink",   "", ""),
    (1, "colors", "коричневый", "koríchnevyy", "brown",  "", ""),

    # ── Phase 1 — time: days, months, time words ──
    (1, "time", "понедельник",  "ponedél'nik",   "Monday",    "", ""),
    (1, "time", "вторник",      "vtórnik",       "Tuesday",   "", ""),
    (1, "time", "среда",        "sredá",         "Wednesday", "", ""),
    (1, "time", "четверг",      "chetvérg",      "Thursday",  "", ""),
    (1, "time", "пятница",      "pyátnitsa",     "Friday",    "", ""),
    (1, "time", "суббота",      "subbóta",       "Saturday",  "", ""),
    (1, "time", "воскресенье",  "voskresén'ye",  "Sunday",    "", ""),
    (1, "time", "сегодня",      "sevódnya",      "today",     "", "'г' pronounced 'v'"),
    (1, "time", "завтра",       "záftra",        "tomorrow",  "", ""),
    (1, "time", "вчера",        "vcherá",        "yesterday", "", ""),
    (1, "time", "сейчас",       "seychás",       "now",       "", ""),
    (1, "time", "утро",         "útro",          "morning",   "", ""),
    (1, "time", "день",         "den'",          "day",       "", ""),
    (1, "time", "вечер",        "vécher",        "evening",   "", ""),
    (1, "time", "ночь",         "noch'",         "night",     "", ""),
    (1, "time", "неделя",       "nedélya",       "week",      "", ""),
    (1, "time", "месяц",        "mésyats",       "month",     "", ""),
    (1, "time", "год",          "god",           "year",      "", ""),
    (1, "time", "время",        "vrémya",        "time",      "", ""),
    (1, "time", "час",          "chas",          "hour / o'clock","", ""),
    (1, "time", "минута",       "minúta",        "minute",    "", ""),
    (1, "time", "январь",       "yanvár'",       "January",   "", ""),
    (1, "time", "февраль",      "fevrál'",       "February",  "", ""),
    (1, "time", "март",         "mart",          "March",     "", ""),
    (1, "time", "апрель",       "aprél'",        "April",     "", ""),
    (1, "time", "май",          "may",           "May",       "", ""),
    (1, "time", "июнь",         "iyún'",         "June",      "", ""),
    (1, "time", "июль",         "iyúl'",         "July",      "", ""),
    (1, "time", "август",       "ávgust",        "August",    "", ""),
    (1, "time", "сентябрь",     "sentyábr'",     "September", "", ""),
    (1, "time", "октябрь",      "oktyábr'",      "October",   "", ""),
    (1, "time", "ноябрь",       "noyábr'",       "November",  "", ""),
    (1, "time", "декабрь",      "dekábr'",       "December",  "", ""),

    # ── Phase 1 — core verbs (infinitives) ──
    (1, "verbs", "быть",      "byt'",        "to be",              "", ""),
    (1, "verbs", "идти",      "idtí",        "to go (on foot)",    "", ""),
    (1, "verbs", "ходить",    "khodít'",     "to go / walk",       "", "habitual"),
    (1, "verbs", "ехать",     "yékhat'",     "to go (by transport)","", ""),
    (1, "verbs", "есть",      "yest'",       "to eat",             "", ""),
    (1, "verbs", "пить",      "pit'",        "to drink",           "", ""),
    (1, "verbs", "знать",     "znat'",       "to know",            "", ""),
    (1, "verbs", "думать",    "dúmat'",      "to think",           "", ""),
    (1, "verbs", "хотеть",    "khotét'",     "to want",            "", ""),
    (1, "verbs", "говорить",  "govorít'",    "to speak / say",     "", ""),
    (1, "verbs", "понимать",  "ponimát'",    "to understand",      "", ""),
    (1, "verbs", "читать",    "chitát'",     "to read",            "", ""),
    (1, "verbs", "писать",    "pisát'",      "to write",           "", ""),
    (1, "verbs", "работать",  "rabótat'",    "to work",            "", ""),
    (1, "verbs", "жить",      "zhit'",       "to live",            "", ""),
    (1, "verbs", "любить",    "lyubít'",     "to love",            "", ""),
    (1, "verbs", "делать",    "délat'",      "to do / make",       "", ""),
    (1, "verbs", "видеть",    "vídet'",      "to see",             "", ""),
    (1, "verbs", "слушать",   "slúshat'",    "to listen",          "", ""),
    (1, "verbs", "смотреть",  "smotrét'",    "to watch / look",    "", ""),
    (1, "verbs", "купить",    "kupít'",      "to buy",             "", ""),
    (1, "verbs", "дать",      "dat'",        "to give",            "", ""),
    (1, "verbs", "спать",     "spat'",       "to sleep",           "", ""),
    (1, "verbs", "играть",    "igrát'",      "to play",            "", ""),

    # ── Phase 1 — core adjectives ──
    (1, "adjectives", "большой",    "bol'shóy",   "big",         "", ""),
    (1, "adjectives", "маленький",  "málen'kiy",  "small",       "", ""),
    (1, "adjectives", "хороший",    "khoróshiy",  "good",        "", ""),
    (1, "adjectives", "плохой",     "plokhóy",    "bad",         "", ""),
    (1, "adjectives", "новый",      "nóvyy",      "new",         "", ""),
    (1, "adjectives", "старый",     "stáryy",     "old",         "", ""),
    (1, "adjectives", "красивый",   "krasívyy",   "beautiful",   "", ""),
    (1, "adjectives", "дорогой",    "dorogóy",    "expensive / dear","", ""),
    (1, "adjectives", "дешёвый",    "deshyóvyy",  "cheap",       "", ""),
    (1, "adjectives", "горячий",    "goryáchiy",  "hot (object)","", ""),
    (1, "adjectives", "холодный",   "kholódnyy",  "cold",        "", ""),
    (1, "adjectives", "быстрый",    "býstryy",    "fast",        "", ""),
    (1, "adjectives", "медленный",  "médlennyy",  "slow",        "", ""),
    (1, "adjectives", "лёгкий",     "lyókhkiy",   "easy / light","", ""),
    (1, "adjectives", "трудный",    "trúdnyy",    "difficult",   "", ""),
    (1, "adjectives", "вкусный",    "vkúsnyy",    "tasty",       "", ""),
    (1, "adjectives", "счастливый", "schastlívyy","happy",       "", "'т' is silent"),

    # ── Phase 2 — food & drink ──
    (2, "food", "чай",      "chay",      "tea",        "", ""),
    (2, "food", "молоко",   "molokó",    "milk",       "", "ma-la-KO — unstressed о→a"),
    (2, "food", "сок",      "sok",       "juice",      "", ""),
    (2, "food", "пиво",     "pívo",      "beer",       "", ""),
    (2, "food", "вино",     "vinó",      "wine",       "", ""),
    (2, "food", "суп",      "sup",       "soup",       "", ""),
    (2, "food", "мясо",     "myáso",     "meat",       "", ""),
    (2, "food", "рыба",     "rýba",      "fish",       "", ""),
    (2, "food", "курица",   "kúritsa",   "chicken",    "", ""),
    (2, "food", "овощи",    "óvoshchi",  "vegetables", "", ""),
    (2, "food", "фрукты",   "frúkty",    "fruit",      "", ""),
    (2, "food", "яблоко",   "yábloko",   "apple",      "", ""),
    (2, "food", "сыр",      "syr",       "cheese",     "", ""),
    (2, "food", "масло",    "máslo",     "butter / oil","", ""),
    (2, "food", "сахар",    "sákhar",    "sugar",      "", ""),
    (2, "food", "соль",     "sol'",      "salt",       "", ""),
    (2, "food", "завтрак",  "záftrak",   "breakfast",  "", ""),
    (2, "food", "обед",     "obéd",      "lunch",      "", ""),
    (2, "food", "ужин",     "úzhin",     "dinner",     "", ""),
    (2, "food", "ресторан", "restorán",  "restaurant", "", ""),
    (2, "food", "кафе",     "kafé",      "café",       "", ""),

    # ── Phase 2 — shopping & money ──
    (2, "shopping", "деньги",     "dén'gi",       "money",      "", ""),
    (2, "shopping", "рубль",      "rubl'",        "rouble",     "", ""),
    (2, "shopping", "чек",        "chek",         "receipt",    "", ""),
    (2, "shopping", "Я беру это", "ya berú éto",  "I'll take it","", ""),
    (2, "shopping", "открыто",    "otkrýto",      "open",       "", ""),
    (2, "shopping", "закрыто",    "zakrýto",      "closed",     "", ""),

    # ── Phase 2 — directions & transport ──
    (2, "directions", "автобус",  "avtóbus",   "bus",       "", ""),
    (2, "directions", "поезд",    "póezd",     "train",     "", ""),
    (2, "directions", "самолёт",  "samolyót",  "plane",     "", ""),
    (2, "directions", "машина",   "mashína",   "car",       "", ""),
    (2, "directions", "билет",    "bilét",     "ticket",    "", ""),
    (2, "directions", "здесь",    "zdes'",     "here",      "", ""),
    (2, "directions", "там",      "tam",       "there",     "", ""),
    (2, "directions", "близко",   "blízko",    "near",      "", ""),
    (2, "directions", "далеко",   "dalekó",    "far",       "", ""),
    (2, "directions", "остановка","ostanóvka", "(bus) stop","", ""),

    # ── Phase 2 — places ──
    (2, "places", "город",        "górod",         "city",        "", ""),
    (2, "places", "улица",        "úlitsa",        "street",      "", ""),
    (2, "places", "дом",          "dom",           "house / home","", ""),
    (2, "places", "магазин",      "magazín",       "shop",        "", ""),
    (2, "places", "рынок",        "rýnok",         "market",      "", ""),
    (2, "places", "аптека",       "aptéka",        "pharmacy",    "", ""),
    (2, "places", "банк",         "bank",          "bank",        "", ""),
    (2, "places", "больница",     "bol'nítsa",     "hospital",    "", ""),
    (2, "places", "школа",        "shkóla",        "school",      "", ""),
    (2, "places", "университет",  "universitét",   "university",  "", ""),
    (2, "places", "гостиница",    "gostínitsa",    "hotel",       "", ""),
    (2, "places", "вокзал",       "vokzál",        "train station","", ""),
    (2, "places", "аэропорт",     "aeropórt",      "airport",     "", ""),

    # ── Phase 2 — lodging ──
    (2, "lodging", "номер",   "nómer",   "room (hotel)", "", ""),
    (2, "lodging", "комната", "kómnata", "room",         "", ""),

    # ── Phase 3 — question words ──
    (3, "questions", "Сколько?", "skól'ko?", "How much / many?", "", ""),
    (3, "questions", "Какой?",   "kakóy?",   "Which / what kind?","", ""),
    (3, "questions", "Чей?",     "chey?",    "Whose?",           "", ""),

    # ── Phase 3 — conversation ──
    (3, "conversation", "Конечно",         "konéchno",       "Of course",        "", "'ч' → 'sh'"),
    (3, "conversation", "Может быть",       "mózhet byt'",    "Maybe",            "", ""),
    (3, "conversation", "Я не знаю",        "ya ne znáyu",    "I don't know",     "", ""),
    (3, "conversation", "Я согласен",       "ya soglásen",    "I agree",          "", "m; согласна (f)"),
    (3, "conversation", "Всё хорошо",       "vsyo khoroshó",  "All good",         "", ""),
    (3, "conversation", "Помогите!",        "pomogíte!",      "Help!",            "", ""),
    (3, "conversation", "Я люблю тебя",     "ya lyublyú tebyá","I love you",      "", ""),
    (3, "conversation", "Хорошего дня",     "khoróshevo dnya","Have a good day",  "", ""),
    (3, "conversation", "Как по-русски…?",  "kak po-rússki…?","How do you say … in Russian?","", ""),
    (3, "conversation", "Я устал",          "ya ustál",       "I'm tired",        "", "m; устала (f)"),
    (3, "conversation", "Я голоден",        "ya góloden",     "I'm hungry",       "", "m; голодна (f)"),

    # ── Phase 3 — connectors & common words ──
    (3, "grammar", "и",           "i",            "and",     "", ""),
    (3, "grammar", "но",          "no",           "but",     "", ""),
    (3, "grammar", "или",         "íli",          "or",      "", ""),
    (3, "grammar", "потому что",  "potomú chto",  "because", "", ""),
    (3, "grammar", "тоже",        "tózhe",        "also / too","", ""),
    (3, "grammar", "очень",       "óchen'",       "very",    "", ""),
    (3, "grammar", "немного",     "nemnógo",      "a little","", ""),

    # ── Phase 3 — weather ──
    (3, "weather", "погода",           "pogóda",           "weather",       "", ""),
    (3, "weather", "дождь",            "dozhd'",           "rain",          "", ""),
    (3, "weather", "снег",             "sneg",             "snow",          "", ""),
    (3, "weather", "солнце",           "sóntse",           "sun",           "", "'л' is silent"),
    (3, "weather", "ветер",            "véter",            "wind",          "", ""),
    (3, "weather", "Сегодня жарко",    "sevódnya zhárko",  "It's hot today","", ""),
    (3, "weather", "Сегодня холодно",  "sevódnya khólodno","It's cold today","", ""),

    # ── Phase 3 — body & health ──
    (3, "body", "голова",         "golová",        "head",           "", ""),
    (3, "body", "рука",           "ruká",          "hand / arm",     "", ""),
    (3, "body", "нога",           "nogá",          "leg / foot",     "", ""),
    (3, "body", "глаз",           "glaz",          "eye",            "", ""),
    (3, "body", "рот",            "rot",           "mouth",          "", ""),
    (3, "body", "ухо",            "úkho",          "ear",            "", ""),
    (3, "body", "живот",          "zhivót",        "stomach",        "", ""),
    (3, "body", "Мне плохо",      "mne plókho",    "I feel unwell",  "", ""),
    (3, "body", "Мне нужен врач", "mne núzhen vrach","I need a doctor","", ""),

    # ── Phase 3 — home & objects ──
    (3, "home", "стол",     "stol",     "table",  "", ""),
    (3, "home", "стул",     "stul",     "chair",  "", ""),
    (3, "home", "дверь",    "dver'",    "door",   "", ""),
    (3, "home", "окно",     "oknó",     "window", "", ""),
    (3, "home", "кровать",  "krovát'",  "bed",    "", ""),
    (3, "home", "книга",    "kníga",    "book",   "", ""),
    (3, "home", "телефон",  "telefón",  "phone",  "", ""),

    # ── A1 content top-up (more verbs, adjectives, time, phrases, places) ──
    (1, "verbs", "гулять",     "gulyát",     "to walk / stroll",   "", ""),
    (1, "verbs", "встречать",  "vstrechát",  "to meet",            "", ""),
    (1, "verbs", "продавать",  "prodavát",   "to sell",            "", ""),
    (1, "verbs", "помогать",   "pomogát",    "to help",            "", ""),
    (1, "verbs", "открывать",  "otkryvát",   "to open",            "", ""),
    (1, "verbs", "закрывать",  "zakryvát",   "to close",           "", ""),
    (1, "verbs", "ждать",      "zhdat",      "to wait",            "", ""),
    (1, "verbs", "начинать",   "nachinát",   "to begin",           "", ""),
    (1, "verbs", "спрашивать", "spráshivat", "to ask",             "", ""),
    (1, "verbs", "отвечать",   "otvechát",   "to answer",          "", ""),
    (1, "verbs", "звонить",    "zvonít",     "to call (phone)",    "", ""),
    (1, "verbs", "готовить",   "gotóvit",    "to cook",            "", ""),
    (1, "verbs", "учить",      "uchít",      "to learn / study",   "", ""),
    (1, "adjectives", "молодой",     "molodóy",    "young",       "", ""),
    (1, "adjectives", "интересный",  "interésnyy", "interesting", "", ""),
    (1, "adjectives", "скучный",     "skúchnyy",   "boring",      "", ""),
    (1, "adjectives", "тёплый",      "tyóplyy",    "warm",        "", ""),
    (1, "adjectives", "сильный",     "sílnyy",     "strong",      "", ""),
    (1, "adjectives", "важный",      "vázhnyy",    "important",   "", ""),
    (1, "adjectives", "правильный",  "právilnyy",  "correct",     "", ""),
    (1, "adjectives", "грустный",    "grústnyy",   "sad",         "", ""),
    (1, "time", "рано",     "ráno",     "early",      "", ""),
    (1, "time", "поздно",   "pózdno",   "late",       "", ""),
    (1, "time", "всегда",   "vsegdá",   "always",     "", ""),
    (1, "time", "иногда",   "inogdá",   "sometimes",  "", ""),
    (1, "time", "никогда",  "nikogdá",  "never",      "", ""),
    (1, "time", "обычно",   "obýchno",  "usually",    "", ""),
    (2, "places", "кухня",    "kúkhnya",  "kitchen",    "", ""),
    (2, "places", "спальня",  "spálnya",  "bedroom",    "", ""),
    (2, "places", "квартира", "kvartíra", "apartment",  "", ""),
    (2, "places", "парк",     "park",     "park",       "", ""),
    (3, "home", "диван",   "diván",   "sofa",   "", ""),
    (3, "conversation", "Ничего",              "nichevó",               "Never mind / it's nothing", "", ""),
    (3, "conversation", "Правда?",             "právda?",               "Really?",                   "", ""),
    (3, "conversation", "Я не уверен",         "ya ne uvéren",          "I'm not sure",              "", "m; не уверена (f)"),
    (3, "conversation", "Повторите, пожалуйста","povtoríte, pozháluysta","Please repeat",            "", ""),
    (3, "conversation", "Помедленнее, пожалуйста","pomédlennee, pozháluysta","Slower, please",       "", ""),
    (3, "conversation", "С удовольствием",     "s udovólstviyem",       "With pleasure / gladly",    "", ""),
]


def _russian_seed_user(db, user_id: int):
    """Idempotently load the starter deck into a user's russian_vocab. INSERT OR
    IGNORE dedupes on UNIQUE(user_id, cyrillic, english), so it's safe to re-run."""
    db.executemany(
        "INSERT OR IGNORE INTO russian_vocab "
        "(user_id, phase, category, cyrillic, translit, english, example, note, source) "
        "VALUES (?,?,?,?,?,?,?,?, 'seed')",
        [(user_id, p, cat, cy, tr, en, ex, nt) for (p, cat, cy, tr, en, ex, nt) in RUSSIAN_SEED])
    db.commit()


def _russian_ensure_seed(db, user_id: int):
    """Seed on first use, and top-up existing users when the seed deck grows.
    Runs the (idempotent) INSERT OR IGNORE only when the user is missing seed rows."""
    seeded = db.execute(
        "SELECT COUNT(*) AS c FROM russian_vocab WHERE user_id = ? AND source = 'seed'",
        (user_id,)).fetchone()["c"]
    if seeded < len(RUSSIAN_SEED):
        _russian_seed_user(db, user_id)


@app.get("/api/russian/vocab")
def russian_vocab_list(phase: int = -1, user_id: int = Depends(get_current_user)):
    """The user's Russian cards (optionally one phase). Lazy-seeds the starter
    deck on first call, so a new user has content with zero setup."""
    db = get_db()
    _russian_ensure_seed(db, user_id)
    if phase >= 0:
        rows = db.execute(
            "SELECT * FROM russian_vocab WHERE user_id = ? AND phase = ? ORDER BY category, id",
            (user_id, phase)).fetchall()
    else:
        rows = db.execute(
            "SELECT * FROM russian_vocab WHERE user_id = ? ORDER BY phase, category, id",
            (user_id,)).fetchall()
    db.close()
    return {"words": [dict(r) for r in rows]}


@app.post("/api/russian/vocab")
async def russian_vocab_add(request: Request, user_id: int = Depends(get_current_user)):
    """Bulk-add custom cards. Body: {words:[{phase, category, cyrillic, translit,
    english, example, note}]}. INSERT OR IGNORE keeps re-imports idempotent."""
    body = await request.json()
    items = body.get("words")
    if not isinstance(items, list) or not items:
        raise HTTPException(400, "Provide a non-empty 'words' array")
    db = get_db()
    added = 0
    for it in items:
        if not isinstance(it, dict):
            continue
        cy = (it.get("cyrillic") or "").strip()
        en = (it.get("english") or "").strip()
        if not cy or not en:
            continue
        cat = (it.get("category") or "vocab").strip().lower()
        if cat not in RUSSIAN_CATEGORIES:
            cat = "vocab"
        try:
            phase = int(it.get("phase", 1))
        except (TypeError, ValueError):
            phase = 1
        cur = db.execute(
            "INSERT OR IGNORE INTO russian_vocab "
            "(user_id, phase, category, cyrillic, translit, english, example, note, source) "
            "VALUES (?,?,?,?,?,?,?,?, 'user')",
            (user_id, phase, cat, cy, (it.get("translit") or "").strip(), en,
             (it.get("example") or "").strip(), (it.get("note") or "").strip()))
        added += cur.rowcount
    db.commit()
    db.close()
    return {"added": added, "requested": len(items)}


@app.delete("/api/russian/vocab/{vid}")
def russian_vocab_delete(vid: int, user_id: int = Depends(get_current_user)):
    """Remove one card (scoped to the requesting user)."""
    db = get_db()
    cur = db.execute("DELETE FROM russian_vocab WHERE id = ? AND user_id = ?", (vid, user_id))
    db.commit()
    db.close()
    if cur.rowcount == 0:
        raise HTTPException(404, "Card not found")
    return {"deleted": vid}


@app.get("/api/russian/drills")
def russian_drills(phase: int = -1, user_id: int = Depends(get_current_user)):
    """Due cards for an SM-2 review session (optionally within one phase):
    next_review NULL or <= today, due-first then newest. Pure DB — no AI."""
    db = get_db()
    _russian_ensure_seed(db, user_id)
    today = date.today().isoformat()
    params = [user_id, today]
    if phase >= 0:
        phase_clause = " AND phase = ?"
        params.append(phase)
    else:
        # "All due" excludes the alphabet — phase 0 is practiced by voice, not flip-drilled
        phase_clause = " AND phase > 0"
    rows = db.execute(
        f"""SELECT * FROM russian_vocab WHERE user_id = ?
            AND (next_review IS NULL OR next_review <= ?){phase_clause}
            ORDER BY COALESCE(next_review, '1970-01-01') ASC, review_count ASC LIMIT 30""",
        params).fetchall()
    db.close()
    return {"cards": [dict(r) for r in rows]}


@app.post("/api/russian/vocab/{vid}/result")
async def russian_vocab_result(vid: int, request: Request, user_id: int = Depends(get_current_user)):
    """One SM-2 step on a Russian card — mirrors /api/flashcards/{cid}/result."""
    body = await request.json()
    correct = body.get("correct", False)
    db = get_db()
    card = db.execute("SELECT * FROM russian_vocab WHERE id = ? AND user_id = ?", (vid, user_id)).fetchone()
    if not card:
        db.close()
        raise HTTPException(404, "Card not found")
    new_interval, new_ease, new_count, next_review = sm2_schedule(
        card["ease_factor"], card["srs_interval"], card["review_count"], correct)
    db.execute(
        """UPDATE russian_vocab SET
           times_seen = times_seen + 1, times_correct = times_correct + ?,
           last_seen = ?, srs_interval = ?, ease_factor = ?, review_count = ?, next_review = ?
           WHERE id = ?""",
        (1 if correct else 0, datetime.now().isoformat(),
         new_interval, round(new_ease, 4), new_count, next_review, vid))
    db.execute(
        "INSERT INTO russian_review_log (vocab_id, user_id, phase, correct) VALUES (?,?,?,?)",
        (vid, user_id, card["phase"], 1 if correct else 0))
    db.commit()
    db.close()
    return {"ok": True, "next_review": next_review, "interval_days": new_interval}


@app.get("/api/russian/stats")
def russian_stats(user_id: int = Depends(get_current_user)):
    """Due / new / learned / mature counts + per-phase totals and saved progress."""
    db = get_db()
    today = date.today().isoformat()
    def _count(where, *a):
        return db.execute(f"SELECT COUNT(*) AS c FROM russian_vocab WHERE user_id = ?{where}",
                          (user_id, *a)).fetchone()["c"]
    total   = _count("")
    # due_today drives the nav + Drill badges → count only drillable vocab (phase > 0);
    # the alphabet (phase 0) is voice-practiced, tracked client-side.
    due     = _count(" AND phase > 0 AND (next_review IS NULL OR next_review <= ?)", today)
    new_c   = _count(" AND review_count = 0")
    learned = _count(" AND review_count > 0")
    mature  = _count(" AND srs_interval > 21")
    by_phase = {}
    for r in db.execute(
        """SELECT phase, COUNT(*) AS total,
                  SUM(CASE WHEN review_count > 0 THEN 1 ELSE 0 END) AS learned,
                  SUM(CASE WHEN (next_review IS NULL OR next_review <= ?) THEN 1 ELSE 0 END) AS due
           FROM russian_vocab WHERE user_id = ? GROUP BY phase""",
        (today, user_id)).fetchall():
        by_phase[r["phase"]] = {"total": r["total"], "learned": r["learned"] or 0, "due": r["due"] or 0}
    progress = {r["phase"]: r["status"] for r in db.execute(
        "SELECT phase, status FROM russian_progress WHERE user_id = ?", (user_id,)).fetchall()}
    db.close()
    return {"due_today": due, "new_cards": new_c, "learned": learned, "mature": mature,
            "total": total, "by_phase": by_phase, "progress": progress}


@app.post("/api/russian/progress")
async def russian_progress_set(request: Request, user_id: int = Depends(get_current_user)):
    """Set a phase's status. Body: {phase:int 0–4, status:'not_started'|'in_progress'|'done'}."""
    body = await request.json()
    try:
        phase = int(body.get("phase"))
    except (TypeError, ValueError):
        raise HTTPException(400, "phase must be an integer 0–4")
    status = (body.get("status") or "").strip()
    if status not in ("not_started", "in_progress", "done"):
        raise HTTPException(400, "status must be not_started, in_progress, or done")
    db = get_db()
    db.execute(
        """INSERT INTO russian_progress (user_id, phase, status, updated_at)
           VALUES (?,?,?,CURRENT_TIMESTAMP)
           ON CONFLICT(user_id, phase) DO UPDATE SET status = ?, updated_at = CURRENT_TIMESTAMP""",
        (user_id, phase, status, status))
    db.commit()
    db.close()
    return {"ok": True, "phase": phase, "status": status}


# ── Quiz ─────────────────────────────────────────────────────────────────────

DIFF_INSTRUCTIONS = {
    'easy': """Generate EASY questions only — direct factual recall.
- Single concept per question; answer is clearly stated in the material
- No multi-step reasoning or clinical vignettes required
- Set "difficulty": "easy" on every question""",
    'medium': """Generate MEDIUM difficulty questions — application and mechanism understanding.
- Require understanding relationships and mechanisms, not just bare recall
- Short clinical scenarios (1–2 sentences) where helpful; two-step reasoning
- Set "difficulty": "medium" on every question""",
    'hard': """Generate HARD questions — built on NICHE, easily-missed details from the material.
- MINE THE FINE PRINT: go past the headline facts every student already knows. Build questions on the SPECIFIC details actually stated in the source — exact named structures, precise spatial/sequential relationships, specific values, qualifying conditions, exceptions, "only/except/unless" caveats, and the particular wording the material uses. The kind of detail a student skims past on first read.
- AVOID textbook headlines. A fact so famous a student could answer it without ever opening this module (e.g. "the left ventricle is thicker because of systemic resistance") is too easy — find the subtler point instead.
- A great hard question hinges on ONE subtle distinction or detail from the text; the student must have actually read carefully to get it.
- Every fact and term must be drawn from THIS module — never outside knowledge. The difficulty is in the subtlety of the detail, not in reaching beyond the material.
- Distractors must be wrong only because of a subtle, specific detail — each is what a student would pick if they half-remembered the material.
- Set "difficulty": "hard" on every question""",
    'mixed': """Generate a MIX of difficulties — label every question:
- ~30% easy  ("difficulty":"easy")  — direct recall
- ~40% medium ("difficulty":"medium") — application/mechanism, short vignettes
- ~30% hard  ("difficulty":"hard")  — deep reasoning that APPLIES taught concepts to novel scenarios; difficulty from reasoning depth, never from facts the module didn't teach""",
}

@app.post("/api/generate/quiz/{mid}")
def generate_quiz(mid: int, difficulty: str = "mixed", force: bool = False,
                  user_id: int = Depends(get_current_user)):
    db = get_db()
    if not user_can_access(db, mid, user_id):
        db.close()
        raise HTTPException(403, "No access to this material")
    mat = db.execute("SELECT * FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        raise HTTPException(404, "Material not found")

    # Difficulty filter shared by the bank queries below. "adaptive" and "mixed" both
    # span all levels, so they don't constrain the bank by a single difficulty.
    if difficulty in ("mixed", "adaptive"):
        diff_sql, diff_args = "", []
    else:
        diff_sql, diff_args = " AND difficulty = ?", [difficulty]

    # Adaptive bank: serve UNSEEN questions (ones the student hasn't answered yet)
    # straight from the saved pool with NO AI call. We only spend on the model when
    # the fresh pool runs low — or when the student forces a brand-new batch.
    if not force:
        unseen = db.execute(
            f"""SELECT COUNT(*) AS c FROM quiz_questions
                WHERE material_id = ? AND user_id = ?{diff_sql}
                  AND id NOT IN (SELECT question_id FROM quiz_attempts WHERE user_id = ?)""",
            tuple([mid, user_id] + diff_args + [user_id])
        ).fetchone()["c"]
        # Serve the saved bank with NO AI call when either: enough unseen remain, OR we're
        # out of API credits (no key) but the bank already has questions — so a fully-seeded
        # guide stays re-quizzable offline instead of erroring once its unseen run low.
        bank_total = db.execute(
            f"SELECT COUNT(*) AS c FROM quiz_questions WHERE material_id = ? AND user_id = ?{diff_sql}",
            tuple([mid, user_id] + diff_args)
        ).fetchone()["c"]
        if unseen >= QUIZ_MIN_UNSEEN or (bank_total > 0 and not os.getenv("ANTHROPIC_API_KEY")):
            db.close()
            return {"count": unseen if unseen >= QUIZ_MIN_UNSEEN else bank_total, "existing": True}

    # Feed the model the recent stems already in the bank so it writes genuinely NEW
    # questions instead of reworded duplicates (the main source of "stale" quizzes).
    prior_stems = [r["question"] for r in db.execute(
        f"""SELECT question FROM quiz_questions
            WHERE material_id = ? AND user_id = ?{diff_sql}
            ORDER BY id DESC LIMIT 40""",
        tuple([mid, user_id] + diff_args)
    ).fetchall()]
    avoid_block = ""
    if prior_stems:
        listed = "\n".join(f"- {s}" for s in prior_stems)
        avoid_block = f"""
⛔ ALREADY IN THE STUDENT'S BANK — do NOT repeat or lightly reword any of these. Write questions on DIFFERENT facts, sub-topics, and angles, with different sentence structure. If your idea is a near-duplicate of one below, pick a less-covered point from the material instead:
{listed}
"""

    # Concept-graph context: the links the student has already built between concepts,
    # plus where they're weak/at-risk — so regeneration writes INTEGRATION questions
    # that test the CONNECTIONS, not more isolated recall. (This is the payoff of the
    # node-linking system; it only spends API here, when the student asks for more.)
    graph_block = ""
    adj = concept_links(db, user_id)
    if adj:
        wscore = propagated_weakness(db, user_id)
        pairs = sorted(
            {tuple(sorted((a, b))): w for a, nb in adj.items() for b, w in nb.items()}.items(),
            key=lambda kv: -kv[1]
        )[:8]
        link_lines = "; ".join(f"{a} ↔ {b}" for (a, b), _ in pairs)
        weak_list = ", ".join(t for t, _ in sorted(wscore.items(), key=lambda kv: -kv[1])[:5])
        if link_lines:
            graph_block = f"""
🕸 CONCEPT MAP — the student's concepts are linked like this: {link_lines}.
Weakest / at-risk concepts right now (target these): {weak_list}.
Write 2-3 INTEGRATION questions that test the CONNECTION between two linked concepts above (how one drives, depends on, or contrasts with the other) — not isolated recall. Weight the set toward the weak/at-risk concepts."""

    if difficulty == "adaptive":
        # Centre new questions on the student's current level for this material, with a
        # consolidation rung below and a stretch rung above so the bank can keep adapting.
        tg = adaptive_targets(db, user_id, mid)
        center = min(2, round(sum(tg.values()) / len(tg))) if tg else 1
        c_name, easier, harder = RANK_NAME[center], RANK_NAME[max(0, center - 1)], RANK_NAME[min(2, center + 1)]
        diff_prompt = (
            f"Generate questions ADAPTED to the student's current performance "
            f"(they are working at roughly '{c_name}' level on this material):\n"
            f"- ~50% at {c_name} level\n"
            f"- ~25% at {easier} level (consolidate)\n"
            f"- ~25% at {harder} level (stretch)\n"
            f"Set each question's \"difficulty\" accurately to easy/medium/hard so the app can track mastery.\n\n"
            f"Level guidance:\n{DIFF_INSTRUCTIONS.get(c_name, DIFF_INSTRUCTIONS['medium'])}"
        )
    else:
        diff_prompt = DIFF_INSTRUCTIONS.get(difficulty, DIFF_INSTRUCTIONS['mixed'])

    # Auto-linking: gather topics from other materials
    other_topics = db.execute("""
        SELECT DISTINCT q.topic FROM quiz_questions q
        JOIN user_materials um ON um.material_id = q.material_id
        WHERE um.user_id = ? AND q.material_id != ? AND q.topic IS NOT NULL
        LIMIT 20
    """, (user_id, mid)).fetchall()
    related_str = ", ".join(_normalize_topic(r["topic"]) for r in other_topics) if other_topics else ""
    cross_link = ""
    if related_str:
        cross_link = f"""
CROSS-LINKING: The student is also studying: {related_str}.
Where relevant, test how concepts from THIS material connect to those other topics.
Add a "related" field with 1-2 related topic names when there's a genuine connection."""

    # Detect subject to tailor question style
    subject = (mat["subject"] or "").lower()
    content_sample = (mat["content"] or "")[:500].lower()
    is_quantitative = any(kw in subject or kw in content_sample for kw in
        ["chemistry", "chem", "physics", "pharmacology", "biochem", "calcul", "equation", "molar", "reaction"])

    chem_block = ""
    if is_quantitative:
        chem_block = """
QUANTITATIVE CONTENT DETECTED — include these where the material supports it:
- CALCULATION questions: pH, molarity, dilution, reaction yield, equilibrium constants, molar mass, stoichiometry, dosage calculations
- For calculation questions, show values in the stem and have options be specific numerical answers (with units)
- When a question involves a specific molecule, include its SMILES notation in a "smiles" field for structure rendering
- Example SMILES: ethanol = "CCO", acetic acid = "CC(=O)O", benzene = "c1ccccc1"
- Only include "smiles" when the structure is directly relevant to the question"""
    else:
        chem_block = """
SUBJECT-APPROPRIATE QUESTIONS — match the question style to the subject:
- Anatomy: spatial relationships, clinical correlations ("damage to X nerve would cause..."), identify structures by description
- Physiology: mechanisms, feedback loops, "what happens when..."
- Pathology: disease presentations, histological findings, differential diagnosis
- Medicine: clinical vignettes with diagnosis/management decisions
- Do NOT include calculation questions or chemical structures unless the material explicitly covers quantitative content"""

    instructions = f"""You are a UNIVERSITY EXAM question writer. Use the source material above as the TOPIC LIST, but write real exam questions that test understanding of the subject — exactly like a lecturer setting a final exam.

📚 COVER THE WHOLE MODULE — DO NOT CLUSTER ON ONE SECTION:
- The material spans several distinct sections/lessons. Before writing, mentally list every major section it contains, then DISTRIBUTE questions across ALL of them — roughly even coverage, not 5 questions on one topic and none on the rest.
- Each question should target a DIFFERENT concept. Do NOT ask the same fact two ways or write near-duplicate questions.
- If the material has 8 sections and you write 14 questions, aim for ~1-2 per section. A section covered later in the material is just as important as one covered early — make sure the back half of the module is represented, not just the opening topics.

QUESTION STYLE — adapt to the subject matter:
- Test APPLICATION and REASONING, not just recall
- Include questions that require INTEGRATING multiple concepts
- Some applied / clinical framing is good ("A patient...", "A researcher observes...") for harder questions, but it must test concepts THAT ARE ACTUALLY IN THE SOURCE MATERIAL. Do NOT invent advanced scenarios that require specialist terms, named enzymes, drug names, or biochemical pathways the student has not studied. If a term doesn't appear in the source material, don't build a question around it.
- KEEP QUESTION STEMS SHORT: 1-2 sentences, ideally under 35 words. A clear focused question beats a long scenario. No padding, no unnecessary backstory.
- Stay at the level of THIS course as reflected in the material — not a specialist board exam. Match the vocabulary the student has actually been taught.
{chem_block}

🔀 VARY THE PHRASING — the student reviews many questions over time, so repetitive templates make it stale:
- Do NOT start most stems with "Which of the following". Rotate formats across the set: direct question, "A patient/researcher…" vignette, cause→effect ("What is the consequence of…"), identify-the-exception ("All of the following EXCEPT…"), ordering/sequence, "Why does…", fill-in-the-concept, compare-two-things.
- No two questions in this set should share the same opening template or sentence shape.
{graph_block}
{avoid_block}
🎯 OPTION-WRITING RULES — these prevent the quiz from being guessable:
- ALL FOUR options MUST be roughly the SAME LENGTH and the SAME LEVEL OF DETAIL/SPECIFICITY. CONCRETE RULE: the longest option may not exceed the shortest by more than ~6 words, and every option's word count should be within that band. The correct answer must NEVER be the longest, most specific, or most technical-sounding one — that is the #1 way these quizzes become guessable. If your correct answer names specific structures or a detailed mechanism (e.g. "sinuses of Valsalva", "suction effect"), then EITHER trim it to match the distractors OR give every distractor an equally specific, equally detailed (but wrong) mechanism. A test-wise student must be UNABLE to spot the answer just by length or specificity.
- NO GIVE-AWAY CLARIFIER ON THE CORRECT ANSWER: never append a parenthetical, textbook tag, or extra qualifier to the correct option that the distractors lack (e.g. ✗ "Anchor the sarcomere to the sarcolemma (lateral force transfer)" — the bracket flags it). Put that depth in the EXPLANATION, never in the option text.
- DISTRACTORS = NEIGHBOURING TRUE FACTS, not implausible extremes. Build each wrong option by swapping the location / protein / mechanism / direction / value for an ADJACENT correct one (e.g. dystrophin↔titin↔nebulin, or a fact true at the opposite end of the same spectrum), so ruling it out needs a SPECIFIC piece of knowledge. Do NOT use eliminable absolutes ("there is no ATP", "calcium is excessive", "it never happens") — those are free eliminations. Make the student weigh "which is MORE right", and include the occasional deliberate TRAP where a common misconception is the tempting wrong option.
- VARY WHICH LETTER IS CORRECT. Do not default to A or B. Spread correct answers evenly across A, B, C, and D. (Positions are also shuffled automatically after generation, so never assume order.)
- Every distractor MUST be a genuine, plausible misconception that a real student could believe — something that tests whether they actually understand. NO throwaway / filler options.
- BANNED lazy distractors (unless that statement is genuinely the correct answer, i.e. a deliberate trick question): "no change occurs", "X is unaffected", "X is independent of Y", "activity stays the same", "none of the above", "it makes no difference". These give away that they're wrong. Only include a "no effect / no change" option when it is actually the correct, counter-intuitive answer — then make it a real trick.
- Distractors should represent DIFFERENT reasoning errors, not just be obviously-wrong noise. Make the student work to eliminate each one.

⛔ ABSOLUTELY FORBIDDEN — these ruin the quiz:
- NEVER write a question whose correct answer is "the material doesn't explain/mention/state this", "not covered", "the text is descriptive not explanatory", or any variant. If a concept isn't fully explained in the source, either USE YOUR OWN SUBJECT KNOWLEDGE to write a proper question about it, or skip that concept entirely.
- NEVER use options like "The study material does not explain...". EVERY option must be a real, substantive answer to the question.
- NEVER test "reading comprehension" or "what is stated vs inferred". Test the actual SCIENCE/subject knowledge.
- In explanations, NEVER reference "the study material", "the source", "the text", "the notes", or "the material". Just explain the concept directly as a teacher would. The material is implied — don't mention it.

DIFFICULTY INSTRUCTIONS:
{diff_prompt}

Return ONLY a JSON array of 12-15 questions:
[{{
  "topic": "Broad topic (2-3 words max, e.g. 'Organic Chemistry', 'Cell Biology', 'Pharmacology')",
  "difficulty": "easy|medium|hard",
  "question": "University exam-style question appropriate to the subject",
  "options": ["A. Option", "B. Option", "C. Option", "D. Option"],
  "correct_answer": "A",
  "explanation": "Direct explanation of WHY the answer is correct. For calculations, show formula → substitution → answer with units. Do NOT mention the source/material/text.",
  "related": ["Related Topic"],
  "smiles": "SMILES string ONLY if a chemical structure is relevant, otherwise omit entirely"
}}]

IMPORTANT: The "topic" must be a BROAD subject category (2-3 words), NOT a specific question description.
{cross_link}

Never repeat the same question. Every question must require THINKING, not just memory."""

    # Quiz uses Sonnet, not Haiku: option-balancing and even-coverage rules
    # need stronger instruction-following than Haiku reliably gives.
    try:
        text = generate_json(mat, instructions, model=MODEL, max_tokens=8000, temperature=0.9)
    except Exception as e:
        # AI unavailable (out of credits, rate-limited, network). Don't 500 — keep
        # the app usable by serving whatever is already in the student's bank.
        have = db.execute(
            f"SELECT COUNT(*) AS c FROM quiz_questions WHERE material_id = ? AND user_id = ?{diff_sql}",
            tuple([mid, user_id] + diff_args)
        ).fetchone()["c"]
        db.close()
        if have:
            return {"count": have, "existing": True, "ai_unavailable": True}
        msg = str(e).lower()
        if "credit balance" in msg or "billing" in msg:
            detail = "Out of API credits — top up to generate new questions. Any saved quizzes still work."
        elif "rate" in msg and "limit" in msg:
            detail = "AI is rate-limited right now — try again in a minute. Saved quizzes still work."
        else:
            detail = "AI generation is temporarily unavailable. Saved quizzes still work."
        raise HTTPException(503, detail)
    try:
        qs = parse_json_response(text)
    except Exception:
        qs = [{"topic": "Error", "difficulty": "medium", "question": "Could not generate quiz", "options": ["A. Error", "B. Error", "C. Error", "D. Error"], "correct_answer": "A", "explanation": "Please try again"}]

    # Shuffle option positions so the correct answer isn't biased toward A/B.
    # The model tends to write the correct option first; we re-randomise it here.
    _LETTERS = ["A", "B", "C", "D"]
    for q in qs:
        opts = q.get("options", [])
        ca = (q.get("correct_answer") or "A").strip().upper()[:1]
        if len(opts) != 4 or ca not in _LETTERS:
            continue
        # Strip any leading "A. " / "B) " style prefix to get clean option text
        bodies = []
        for o in opts:
            txt = str(o)
            m = re.match(r'^\s*[A-Da-d]\s*[\.\)\:\-]\s*(.*)$', txt)
            bodies.append((m.group(1) if m else txt).strip())
        correct_body = bodies[_LETTERS.index(ca)]
        random.shuffle(bodies)
        new_correct = _LETTERS[bodies.index(correct_body)]
        q["options"] = [f"{_LETTERS[i]}. {bodies[i]}" for i in range(4)]
        q["correct_answer"] = new_correct

    # ── Grounding guard: drop questions about concepts not in the source ──────
    # The model sometimes interpolates textbook knowledge the student's material
    # never covered (e.g. respiratory physiology turning up in a cardio module).
    # For each question we extract technical "anchor" words from the stem/topic;
    # if it has several and NONE of them appear anywhere in the source text, the
    # question is off-syllabus and gets removed.
    source_lc = (mat["content"] or "").lower()
    _STOP = {
        "researcher","patient","research","observes","explain","explains","describe",
        "describes","statement","correctly","difference","differs","function","functions",
        "following","because","between","during","increase","increases","decrease",
        "decreases","process","structure","structures","mechanism","mechanisms","produce",
        "produces","require","requires","reduce","reduces","result","results","occurs",
        "compared","whereas","through","without","within","greater","smaller","higher",
        "lower","number","numerous","example","another","recovery","property","properties",
        "release","releases","secrete","secretes","divide","replace","correct","answer",
        "question","scenario","likely","change","changes","affect","affects","relationship",
    }
    def _anchors(text):
        seen = []
        for w in re.findall(r"[a-zA-Z]{6,}", (text or "").lower()):
            if w not in _STOP and w not in seen:
                seen.append(w)
        return seen
    def _in_source(term):
        return term in source_lc or (term.endswith("s") and term[:-1] in source_lc)

    if source_lc:
        kept, dropped = [], 0
        for q in qs:
            anchors = _anchors((q.get("question", "") + " " + q.get("topic", "")))
            if len(anchors) >= 2 and not any(_in_source(a) for a in anchors):
                dropped += 1
                continue
            kept.append(q)
        # Only apply if it doesn't gut the quiz (guards against a misfiring
        # heuristic or a very sparse source leaving almost no questions).
        if kept and len(kept) >= max(6, len(qs) // 2):
            qs = kept

    # Accumulate — APPEND this batch to the bank instead of wiping it, so the
    # student builds a growing pool of fresh questions to draw from over time.
    for q in qs:
        fallback_diff = difficulty if difficulty != 'mixed' else 'medium'
        related = json.dumps(q.get("related", []))
        smiles = q.get("smiles") or None
        db.execute(
            "INSERT INTO quiz_questions (material_id, user_id, topic, difficulty, question, options, correct_answer, explanation, related_topics, smiles) VALUES (?,?,?,?,?,?,?,?,?,?)",
            (mid, user_id, q.get("topic", "General"), q.get("difficulty", fallback_diff),
             q.get("question", ""), json.dumps(q.get("options", [])),
             q.get("correct_answer", "A"), q.get("explanation", ""), related, smiles)
        )

    # Bound the bank: prune the oldest questions the student has ALREADY seen
    # (never an unseen one) so storage stays capped without throwing away fresh content.
    total = db.execute(
        "SELECT COUNT(*) AS c FROM quiz_questions WHERE material_id = ? AND user_id = ?",
        (mid, user_id)
    ).fetchone()["c"]
    if total > QUIZ_BANK_CAP:
        stale = db.execute(
            """SELECT id FROM quiz_questions
               WHERE material_id = ? AND user_id = ?
                 AND id IN (SELECT question_id FROM quiz_attempts WHERE user_id = ?)
               ORDER BY id ASC LIMIT ?""",
            (mid, user_id, user_id, total - QUIZ_BANK_CAP)
        ).fetchall()
        for r in stale:
            db.execute("DELETE FROM quiz_questions WHERE id = ?", (r["id"],))
    db.commit()
    db.close()
    return {"count": len(qs)}


@app.get("/api/quiz")
def get_quiz(material_id: Optional[int] = None, difficulty: Optional[str] = None,
             user_id: int = Depends(get_current_user)):
    db = get_db()
    # Graph-aware weakness: prioritise topics you're failing AND their at-risk
    # neighbours via the concept graph (a topic linked to weak ones surfaces too).
    # Topics absent from the map (never attempted, no links) take a neutral 0.4 so
    # they sort between weak topics and already-mastered ones.
    wscore = propagated_weakness(db, user_id)

    # Performance-driven difficulty. "adaptive" (default) tunes the served level per
    # topic to the student's recent accuracy; an explicit level prioritises that level;
    # "mixed" applies no difficulty preference (legacy behaviour).
    mode = (difficulty or "adaptive").lower()
    targets = adaptive_targets(db, user_id, material_id) if mode == "adaptive" else {}

    base = """SELECT q.*, m.original_name,
                     (SELECT MAX(a.attempted_at) FROM quiz_attempts a
                       WHERE a.question_id = q.id AND a.user_id = ?) AS last_attempt
              FROM quiz_questions q JOIN materials m ON q.material_id = m.id
              WHERE q.user_id = ?"""
    args = [user_id, user_id]
    if material_id:
        base += " AND q.material_id = ?"
        args.append(material_id)
    rows = db.execute(base, tuple(args)).fetchall()
    db.close()

    result = []
    for r in rows:
        d = dict(r)
        try:
            d["options"] = json.loads(d["options"])
        except Exception:
            d["options"] = []
        result.append(d)

    # How well a question's difficulty matches what the student should be doing now.
    # 0 = perfect match; bigger = further off. Used as the lead sort so the served
    # level tracks performance (or an explicitly chosen level).
    def diff_distance(q):
        rank = DIFF_RANK.get(q.get("difficulty") or "medium", 1)
        if mode in DIFF_RANK:                       # explicit level requested
            return abs(rank - DIFF_RANK[mode])
        if mode == "adaptive":
            return abs(rank - targets.get(q.get("topic") or "General", 1))
        return 0                                    # "mixed" → no preference

    # Build the sitting from three pools:
    #   • due      — answered before and the spaced-repetition schedule says it's time
    #                (next_review ≤ today; legacy-answered-but-unscheduled count as due)
    #   • unseen   — never answered: fresh material, adaptive difficulty + weakness order
    #   • notdue   — answered recently, not yet due: lowest priority
    today = date.today().isoformat()
    def weak_pos(q):  # more weak/at-risk first; unknown topics neutral, mastered last
        return -wscore.get(_normalize_topic(q.get("topic") or ""), 0.4)
    unseen = [q for q in result if not q.get("last_attempt")]
    seen   = [q for q in result if q.get("last_attempt")]
    due    = [q for q in seen if (q.get("next_review") or "") <= today]
    notdue = [q for q in seen if (q.get("next_review") or "") >  today]
    random.shuffle(unseen)
    due.sort(key=lambda q: (q.get("next_review") or "", weak_pos(q)))   # most overdue first
    unseen.sort(key=lambda q: (diff_distance(q), weak_pos(q)))
    notdue.sort(key=lambda q: (diff_distance(q), q["last_attempt"]))
    for q in due:
        q["is_review"] = True

    # Due reviews lead (spaced repetition is time-sensitive), but reserve room for fresh
    # questions so a big review backlog never crowds out new material entirely.
    if unseen:
        cap = max(1, (QUIZ_SESSION * 2) // 3)
        ordered = due[:cap] + unseen + due[cap:] + notdue
    else:
        ordered = due + notdue
    for q in ordered:
        q.pop("last_attempt", None)
    return ordered[:QUIZ_SESSION]


@app.post("/api/quiz/{qid}/answer")
async def submit_answer(qid: int, request: Request, user_id: int = Depends(get_current_user)):
    body = await request.json()
    answer = body.get("answer", "").strip()
    db = get_db()
    q = db.execute("SELECT * FROM quiz_questions WHERE id = ? AND user_id = ?", (qid, user_id)).fetchone()
    if not q:
        raise HTTPException(404, "Question not found")
    correct = answer and answer[0].upper() == q["correct_answer"].strip()[0].upper()
    db.execute(
        "INSERT INTO quiz_attempts (question_id, material_id, user_id, topic, user_answer, is_correct) VALUES (?,?,?,?,?,?)",
        (qid, q["material_id"], user_id, q["topic"], answer, 1 if correct else 0)
    )
    # Spaced repetition: schedule when this question should resurface. Miss it and it
    # returns tomorrow; keep getting it right and the gap widens (1 → 6 → ×ease days).
    new_interval, new_ease, new_count, next_review = sm2_schedule(
        q["ease_factor"], q["srs_interval"], q["review_count"], correct)
    db.execute(
        """UPDATE quiz_questions SET
           srs_interval=?, ease_factor=?, review_count=?, next_review=?, last_seen=?
           WHERE id=?""",
        (new_interval, new_ease, new_count, next_review, datetime.now().isoformat(), qid)
    )
    db.commit()
    db.close()
    return {"correct": correct, "correct_answer": q["correct_answer"], "explanation": q["explanation"],
            "next_review": next_review, "interval_days": new_interval}


@app.get("/api/quiz/mistakes")
def get_mistakes(user_id: int = Depends(get_current_user)):
    """Questions the user has most recently answered INCORRECTLY, for review.
    Only shows mistakes that haven't since been answered correctly (still unmastered)."""
    db = get_db()
    rows = db.execute(
        """SELECT a.question_id, a.topic, a.user_answer, MAX(a.attempted_at) AS last_wrong,
                  q.question, q.options, q.correct_answer, q.explanation, q.difficulty
           FROM quiz_attempts a
           JOIN quiz_questions q ON a.question_id = q.id
           WHERE a.user_id = ? AND a.is_correct = 0
             AND a.question_id NOT IN (
                 SELECT question_id FROM quiz_attempts
                 WHERE user_id = ? AND is_correct = 1
             )
           GROUP BY a.question_id
           ORDER BY last_wrong DESC
           LIMIT 50""",
        (user_id, user_id)
    ).fetchall()
    out = []
    for r in rows:
        d = dict(r)
        try:
            d["options"] = json.loads(d["options"])
        except Exception:
            d["options"] = []
        out.append(d)
    db.close()
    return {"mistakes": out, "count": len(out)}


# ── Leaderboard & Quiz Battles ────────────────────────────────────────────────

@app.get("/api/leaderboard")
def get_leaderboard(user_id: int = Depends(get_current_user)):
    db = get_db()
    rows = db.execute("""
        SELECT u.id, u.username,
               COUNT(a.id) as total_attempts,
               COALESCE(SUM(a.is_correct), 0) as total_correct,
               CASE WHEN COUNT(a.id) > 0
                    THEN ROUND(CAST(SUM(a.is_correct) AS REAL) / COUNT(a.id) * 100, 1)
                    ELSE 0 END as accuracy,
               COUNT(DISTINCT DATE(a.attempted_at)) as active_days
        FROM users u
        LEFT JOIN quiz_attempts a ON a.user_id = u.id
        GROUP BY u.id
        ORDER BY total_correct DESC, accuracy DESC
    """).fetchall()

    # Calculate streaks for each user
    leaderboard = []
    for i, r in enumerate(rows):
        # Get streak: consecutive days with quiz activity ending today or yesterday
        dates = db.execute(
            "SELECT DISTINCT DATE(attempted_at) as d FROM quiz_attempts WHERE user_id = ? ORDER BY d DESC",
            (r["id"],)
        ).fetchall()
        streak = 0
        if dates:
            from datetime import date, timedelta
            today = date.today()
            expected = today
            for row in dates:
                d = date.fromisoformat(row["d"])
                if d == expected:
                    streak += 1
                    expected -= timedelta(days=1)
                elif d == expected - timedelta(days=1) and streak == 0:
                    # Allow streak to start from yesterday
                    streak = 1
                    expected = d - timedelta(days=1)
                else:
                    break

        leaderboard.append({
            "rank": i + 1,
            "user_id": r["id"],
            "username": r["username"],
            "total_correct": r["total_correct"],
            "total_attempts": r["total_attempts"],
            "accuracy": r["accuracy"],
            "streak": streak,
            "active_days": r["active_days"],
            "is_me": r["id"] == user_id,
        })
    db.close()
    return leaderboard


@app.post("/api/battles")
async def create_battle(request: Request, user_id: int = Depends(get_current_user)):
    """Create a quiz battle. Picks questions from a topic for everyone to answer."""
    body = await request.json()
    topic = (body.get("topic") or "").strip()
    num_questions = min(int(body.get("num_questions", 10)), 20)

    db = get_db()
    # Find questions matching the topic (from any user's materials — shared battle)
    questions = db.execute(
        "SELECT id FROM quiz_questions WHERE topic LIKE ? ORDER BY RANDOM() LIMIT ?",
        (f"%{topic}%", num_questions)
    ).fetchall()

    if len(questions) < 3:
        db.close()
        raise HTTPException(400, f"Not enough questions for topic '{topic}'. Need at least 3, found {len(questions)}. Generate more quizzes first!")

    qids = ",".join(str(q["id"]) for q in questions)
    cur = db.execute(
        "INSERT INTO quiz_battles (creator_id, topic, question_ids) VALUES (?,?,?)",
        (user_id, topic, qids)
    )
    battle_id = cur.lastrowid
    # Auto-join creator
    db.execute(
        "INSERT INTO battle_participants (battle_id, user_id, total) VALUES (?,?,?)",
        (battle_id, user_id, len(questions))
    )
    db.commit()

    creator = db.execute("SELECT username FROM users WHERE id = ?", (user_id,)).fetchone()
    db.close()
    return {"id": battle_id, "topic": topic, "num_questions": len(questions), "creator": creator["username"]}


@app.get("/api/battles")
def list_battles(user_id: int = Depends(get_current_user)):
    db = get_db()
    battles = db.execute("""
        SELECT b.id, b.topic, b.status, b.created_at, b.question_ids,
               u.username as creator_name, b.creator_id
        FROM quiz_battles b JOIN users u ON b.creator_id = u.id
        ORDER BY b.created_at DESC LIMIT 20
    """).fetchall()

    result = []
    for b in battles:
        participants = db.execute("""
            SELECT bp.user_id, u.username, bp.score, bp.total, bp.completed
            FROM battle_participants bp JOIN users u ON bp.user_id = u.id
            WHERE bp.battle_id = ?
            ORDER BY bp.score DESC
        """, (b["id"],)).fetchall()

        num_qs = len(b["question_ids"].split(","))
        result.append({
            "id": b["id"],
            "topic": b["topic"],
            "status": b["status"],
            "created_at": b["created_at"],
            "creator": b["creator_name"],
            "num_questions": num_qs,
            "participants": [
                {"user_id": p["user_id"], "username": p["username"],
                 "score": p["score"], "total": p["total"], "completed": bool(p["completed"]),
                 "is_me": p["user_id"] == user_id}
                for p in participants
            ],
            "i_joined": any(p["user_id"] == user_id for p in participants),
            "i_completed": any(p["user_id"] == user_id and p["completed"] for p in participants),
        })
    db.close()
    return result


@app.post("/api/battles/{bid}/join")
def join_battle(bid: int, user_id: int = Depends(get_current_user)):
    db = get_db()
    battle = db.execute("SELECT * FROM quiz_battles WHERE id = ?", (bid,)).fetchone()
    if not battle:
        db.close()
        raise HTTPException(404, "Battle not found")
    existing = db.execute(
        "SELECT 1 FROM battle_participants WHERE battle_id = ? AND user_id = ?", (bid, user_id)
    ).fetchone()
    if existing:
        db.close()
        return {"ok": True, "message": "Already joined"}
    num_qs = len(battle["question_ids"].split(","))
    db.execute(
        "INSERT INTO battle_participants (battle_id, user_id, total) VALUES (?,?,?)",
        (bid, user_id, num_qs)
    )
    db.commit()
    db.close()
    return {"ok": True}


@app.get("/api/battles/{bid}/questions")
def get_battle_questions(bid: int, user_id: int = Depends(get_current_user)):
    """Get the questions for a battle (must have joined)."""
    db = get_db()
    participant = db.execute(
        "SELECT * FROM battle_participants WHERE battle_id = ? AND user_id = ?", (bid, user_id)
    ).fetchone()
    if not participant:
        db.close()
        raise HTTPException(403, "Join the battle first")
    if participant["completed"]:
        db.close()
        raise HTTPException(400, "You already completed this battle")

    battle = db.execute("SELECT * FROM quiz_battles WHERE id = ?", (bid,)).fetchone()
    qids = battle["question_ids"].split(",")
    placeholders = ",".join("?" * len(qids))
    questions = db.execute(
        f"SELECT id, topic, question, options, difficulty FROM quiz_questions WHERE id IN ({placeholders})",
        qids
    ).fetchall()
    db.close()

    return [{
        "id": q["id"],
        "topic": q["topic"],
        "question": q["question"],
        "options": json.loads(q["options"]) if isinstance(q["options"], str) else q["options"],
        "difficulty": q["difficulty"],
    } for q in questions]


@app.post("/api/battles/{bid}/submit")
async def submit_battle(bid: int, request: Request, user_id: int = Depends(get_current_user)):
    """Submit all battle answers at once. Body: {answers: {question_id: chosen_answer, ...}}"""
    body = await request.json()
    answers = body.get("answers", {})

    db = get_db()
    participant = db.execute(
        "SELECT * FROM battle_participants WHERE battle_id = ? AND user_id = ?", (bid, user_id)
    ).fetchone()
    if not participant:
        db.close()
        raise HTTPException(403, "Join the battle first")
    if participant["completed"]:
        db.close()
        raise HTTPException(400, "Already submitted")

    battle = db.execute("SELECT * FROM quiz_battles WHERE id = ?", (bid,)).fetchone()
    qids = battle["question_ids"].split(",")
    placeholders = ",".join("?" * len(qids))
    questions = db.execute(
        f"SELECT id, correct_answer, topic, material_id, explanation FROM quiz_questions WHERE id IN ({placeholders})",
        qids
    ).fetchall()

    score = 0
    results = []
    for q in questions:
        given = answers.get(str(q["id"]), "")
        correct = given == q["correct_answer"]
        if correct:
            score += 1
        results.append({
            "question_id": q["id"],
            "your_answer": given,
            "correct_answer": q["correct_answer"],
            "correct": correct,
            "explanation": q["explanation"],
        })
        # Also record in quiz_attempts for leaderboard stats
        db.execute(
            "INSERT INTO quiz_attempts (question_id, material_id, user_id, topic, user_answer, is_correct) VALUES (?,?,?,?,?,?)",
            (q["id"], q["material_id"], user_id, q["topic"], given, 1 if correct else 0)
        )

    db.execute(
        "UPDATE battle_participants SET score = ?, completed = 1, finished_at = CURRENT_TIMESTAMP WHERE battle_id = ? AND user_id = ?",
        (score, bid, user_id)
    )
    db.commit()

    # Return updated standings
    standings = db.execute("""
        SELECT u.username, bp.score, bp.total, bp.completed
        FROM battle_participants bp JOIN users u ON bp.user_id = u.id
        WHERE bp.battle_id = ? ORDER BY bp.score DESC
    """, (bid,)).fetchall()

    db.close()
    return {
        "score": score,
        "total": len(questions),
        "results": results,
        "standings": [{"username": s["username"], "score": s["score"],
                       "total": s["total"], "completed": bool(s["completed"])} for s in standings],
    }


# ── Tutor (streaming) ────────────────────────────────────────────────────────

@app.post("/api/chat")
async def chat(request: Request, user_id: int = Depends(get_current_user)):
    body = await request.json()
    message = body.get("message", "")
    mode = body.get("mode", "explain")
    session_id = body.get("session_id", "default")
    mid = body.get("material_id")

    db = get_db()
    history = list(reversed(db.execute(
        "SELECT role, content FROM chat_messages WHERE session_id = ? AND user_id = ? ORDER BY timestamp DESC LIMIT 12",
        (session_id, user_id)
    ).fetchall()))

    mat_ctx = ""
    if mid and user_can_access(db, mid, user_id):
        mat = db.execute("SELECT original_name, content FROM materials WHERE id = ?", (mid,)).fetchone()
        if mat:
            mat_ctx = f"\n\nContext from '{mat['original_name']}':\n{mat['content'][:3000]}"

    db.execute("INSERT INTO chat_messages (user_id, session_id, role, content) VALUES (?,?,?,?)", (user_id, session_id, "user", message))
    db.commit()
    db.close()

    if mode == "socratic":
        system_text = f"""You are a Socratic medical tutor. Guide the student to discover answers through questions rather than giving them directly. Ask one focused question at a time. Be encouraging but intellectually challenging. Use clinical scenarios. Keep responses to 2-3 sentences + one guiding question.{mat_ctx}"""
    else:
        system_text = f"""You are an expert medical tutor for a health sciences student. Give clear, accurate explanations with clinical context. Use mnemonics and bullet points. Relate concepts to clinical practice. Be thorough but concise.{mat_ctx}"""

    # Cache the system block (contains the material context) so each follow-up turn in a
    # multi-turn conversation reuses it at ~10% cost instead of re-sending it in full.
    system = [{"type": "text", "text": system_text, "cache_control": {"type": "ephemeral"}}]

    messages = [{"role": r["role"], "content": r["content"]} for r in history]
    messages.append({"role": "user", "content": message})

    def stream():
        full = ""
        try:
            with get_client().messages.stream(model=MODEL, max_tokens=1500, system=system, messages=messages) as s:
                for chunk in s.text_stream:
                    full += chunk
                    yield f"data: {json.dumps({'text': chunk})}\n\n"
            conn = sqlite3.connect(DB_PATH)
            conn.execute("INSERT INTO chat_messages (user_id, session_id, role, content) VALUES (?,?,?,?)", (user_id, session_id, "assistant", full))
            conn.commit()
            conn.close()
            yield f"data: {json.dumps({'done': True})}\n\n"
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(stream(), media_type="text/event-stream",
                             headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})


@app.delete("/api/chat/{session_id}")
def clear_chat(session_id: str, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("DELETE FROM chat_messages WHERE session_id = ? AND user_id = ?", (session_id, user_id))
    db.commit()
    db.close()
    return {"ok": True}


# ── Progress ─────────────────────────────────────────────────────────────────

@app.get("/api/progress")
def get_progress(user_id: int = Depends(get_current_user)):
    db = get_db()

    # ── Quiz stats ─────────────────────────────────────────────────────────
    qs = dict(db.execute("SELECT COUNT(*) as t, COALESCE(SUM(is_correct),0) as c FROM (SELECT is_correct, MAX(id) FROM quiz_attempts WHERE user_id = ? GROUP BY question_id)", (user_id,)).fetchone())
    quiz_topics = [dict(r) for r in db.execute(
        """SELECT topic, COUNT(*) as attempts, SUM(is_correct) as correct,
                  CAST(SUM(is_correct) AS REAL)/COUNT(*) as accuracy
           FROM quiz_attempts WHERE user_id = ? GROUP BY topic ORDER BY accuracy ASC""",
        (user_id,)
    ).fetchall()]

    # ── Flashcard stats ────────────────────────────────────────────────────
    fc = dict(db.execute(
        "SELECT COUNT(*) as total, SUM(times_seen) as reviews, SUM(times_correct) as correct FROM flashcards WHERE user_id = ?",
        (user_id,)
    ).fetchone())
    fc_topics = [dict(r) for r in db.execute(
        """SELECT topic, COUNT(*) as attempts, SUM(correct) as correct
           FROM flashcard_log WHERE user_id = ? AND topic IS NOT NULL GROUP BY topic""",
        (user_id,)
    ).fetchall()]

    # ── Daily activity — last 7 days ───────────────────────────────────────
    daily_quiz = [dict(r) for r in db.execute(
        """SELECT DATE(attempted_at) as date, COUNT(*) as attempts, SUM(is_correct) as correct
           FROM quiz_attempts WHERE user_id = ? AND attempted_at >= DATE('now','-7 days')
           GROUP BY DATE(attempted_at) ORDER BY date""",
        (user_id,)
    ).fetchall()]
    daily_fc = [dict(r) for r in db.execute(
        """SELECT DATE(reviewed_at) as date, COUNT(*) as reviews, SUM(correct) as correct
           FROM flashcard_log WHERE user_id = ? AND reviewed_at >= DATE('now','-7 days')
           GROUP BY DATE(reviewed_at) ORDER BY date""",
        (user_id,)
    ).fetchall()]

    # ── Combined topic mastery (quiz + flashcard, no API needed) ───────────
    # Normalize long topic names so chart labels stay readable
    topic_agg = defaultdict(lambda: {"attempts": 0, "correct": 0})
    for r in quiz_topics:
        t = _normalize_topic(r["topic"])
        topic_agg[t]["attempts"] += r["attempts"]
        topic_agg[t]["correct"]  += int(r["correct"] or 0)
    for r in fc_topics:
        t = _normalize_topic(r["topic"])
        topic_agg[t]["attempts"] += r["attempts"]
        topic_agg[t]["correct"]  += int(r["correct"] or 0)
    combined_topics = sorted(
        [{"topic": t, "attempts": d["attempts"], "correct": d["correct"],
          "accuracy": d["correct"] / d["attempts"] if d["attempts"] else 0}
         for t, d in topic_agg.items()],
        key=lambda x: x["accuracy"]
    )

    # ── Reliability-aware weak topics (fixes the "0% after 1 wrong" glitch) ─
    # Delegated to compute_weak_topics (pure, unit-tested): Laplace smoothing
    # + WEAK_MIN_ATTEMPTS guard, sorted weakest first.
    weak_topics = compute_weak_topics(topic_agg)

    # ── Performance per MATERIAL (quiz accuracy grouped by the material) ────
    by_material = [dict(r) for r in db.execute(
        """SELECT m.original_name AS material, COUNT(*) AS attempts, SUM(a.is_correct) AS correct,
                  CAST(SUM(a.is_correct) AS REAL)/COUNT(*) AS accuracy
           FROM quiz_attempts a JOIN materials m ON a.material_id = m.id
           WHERE a.user_id = ? GROUP BY a.material_id ORDER BY accuracy ASC""",
        (user_id,)
    ).fetchall()]

    # ── Counts ─────────────────────────────────────────────────────────────
    mats   = db.execute("SELECT COUNT(*) as c FROM user_materials WHERE user_id = ?", (user_id,)).fetchone()["c"]
    fcs    = db.execute("SELECT COUNT(*) as c FROM flashcards WHERE user_id = ?", (user_id,)).fetchone()["c"]
    slides = db.execute("SELECT COUNT(*) as c FROM revision_slides WHERE user_id = ?", (user_id,)).fetchone()["c"]
    db.close()

    return {
        "quiz": {
            "total":    qs["t"] or 0,
            "correct":  qs["c"] or 0,
            "accuracy": round(((qs["c"] or 0) / max(qs["t"] or 1, 1)) * 100, 1),
            "by_topic": [
                {**t, "topic": _normalize_topic(t["topic"])} for t in quiz_topics
            ],
        },
        "flashcards":      fc,
        "daily":           daily_quiz,
        "daily_fc":        daily_fc,
        "combined_topics": combined_topics,
        "weak_topics":     weak_topics,
        "by_material":     by_material,
        "counts":          {"materials": mats, "flashcards": fcs, "slides": slides},
    }


# ── Knowledge Graph ──────────────────────────────────────────────────────────

@app.get("/api/knowledge-graph")
def knowledge_graph(user_id: int = Depends(get_current_user)):
    """Build a graph of materials ↔ topics with performance data."""
    db = get_db()

    # Get user's materials
    mats = db.execute("""
        SELECT m.id, m.original_name, m.subject, m.file_type, LENGTH(m.content) as chars
        FROM materials m JOIN user_materials um ON um.material_id = m.id
        WHERE um.user_id = ?
    """, (user_id,)).fetchall()

    # Get topics from quiz questions linked to each material
    mat_topics = db.execute("""
        SELECT DISTINCT q.material_id, q.topic
        FROM quiz_questions q
        WHERE q.user_id = ? AND q.topic IS NOT NULL
    """, (user_id,)).fetchall()

    # Get topic performance
    topic_perf = db.execute("""
        SELECT topic, COUNT(*) as attempts, SUM(is_correct) as correct,
               CAST(SUM(is_correct) AS REAL)/COUNT(*) as accuracy
        FROM quiz_attempts WHERE user_id = ? GROUP BY topic
    """, (user_id,)).fetchall()
    perf_map = {_normalize_topic(r["topic"]): {
        "attempts": r["attempts"], "correct": r["correct"] or 0,
        "accuracy": round((r["accuracy"] or 0) * 100, 1)
    } for r in topic_perf}

    # Get flashcard topics per material
    fc_topics = db.execute("""
        SELECT DISTINCT material_id, topic FROM flashcards
        WHERE user_id = ? AND topic IS NOT NULL
    """, (user_id,)).fetchall()

    # Get user's graph customizations
    hidden_nodes = set(r["node_id"] for r in db.execute(
        "SELECT node_id FROM graph_hidden_nodes WHERE user_id = ?", (user_id,)
    ).fetchall())
    hidden_edges = set((r["source"], r["target"]) for r in db.execute(
        "SELECT source, target FROM graph_hidden_edges WHERE user_id = ?", (user_id,)
    ).fetchall())
    custom_edges = [(r["source"], r["target"]) for r in db.execute(
        "SELECT source, target FROM graph_custom_edges WHERE user_id = ?", (user_id,)
    ).fetchall()]

    # Real concept links the AI asserted via related_topics (previously unused).
    concept_adj = concept_links(db, user_id)

    db.close()

    # Build nodes and edges
    nodes = []
    edges = []
    topic_set = {}  # normalized_topic → node_id
    mat_node_ids = {}

    # Material nodes
    for m in mats:
        nid = f"mat_{m['id']}"
        if nid in hidden_nodes:
            continue
        mat_node_ids[m["id"]] = nid
        nodes.append({
            "id": nid, "type": "material", "label": m["original_name"][:30],
            "subject": m["subject"], "file_type": m["file_type"],
            "size": max(8, min(20, (m["chars"] or 0) / 5000 + 8)),
        })

    # Collect UNIQUE topic→material links (deduplicate across quiz + flashcard)
    all_links = set()
    for r in mat_topics:
        t = _normalize_topic(r["topic"])
        if t and r["material_id"] in mat_node_ids:
            all_links.add((r["material_id"], t))
    for r in fc_topics:
        t = _normalize_topic(r["topic"])
        if t and r["material_id"] in mat_node_ids:
            all_links.add((r["material_id"], t))

    # Edge dedup set — prevents any duplicate edges
    edge_set = set()
    def add_edge(src, tgt, **extra):
        key = tuple(sorted([src, tgt]))
        if key in edge_set:
            return False
        if key in hidden_edges or (src, tgt) in hidden_edges or (tgt, src) in hidden_edges:
            return False
        edge_set.add(key)
        edges.append({"source": src, "target": tgt, **extra})
        return True

    # Topic nodes + edges to materials
    for mid, topic in all_links:
        if topic not in topic_set:
            tid = f"topic_{len(topic_set)}"
            if tid in hidden_nodes:
                topic_set[topic] = tid
                continue
            topic_set[topic] = tid
            perf = perf_map.get(topic, {"attempts": 0, "correct": 0, "accuracy": 0})
            nodes.append({
                "id": tid, "type": "topic", "label": topic,
                "accuracy": perf["accuracy"], "attempts": perf["attempts"],
                "size": max(6, min(16, perf["attempts"] / 2 + 6)),
            })
        tid = topic_set[topic]
        node_ids_set = set(n["id"] for n in nodes)
        src, tgt = mat_node_ids[mid], tid
        if src in node_ids_set and tgt in node_ids_set:
            add_edge(src, tgt)

    # Subject nodes — group materials by subject
    subjects = {}
    for m in mats:
        subj = m["subject"] or "General"
        mid = m["id"]
        if mid not in mat_node_ids:
            continue
        if subj not in subjects:
            sid = f"subj_{len(subjects)}"
            if sid in hidden_nodes:
                subjects[subj] = sid
                continue
            subjects[subj] = sid
            nodes.append({"id": sid, "type": "subject", "label": subj, "size": 22})
        add_edge(subjects[subj], mat_node_ids[mid])

    # Custom user-created edges
    node_ids_set = set(n["id"] for n in nodes)
    for src, tgt in custom_edges:
        if src in node_ids_set and tgt in node_ids_set:
            add_edge(src, tgt, custom=True)

    # Concept edges from the AI's related_topics — real semantic links between
    # concepts (e.g. "Cardiac Output" ↔ "Stroke Volume"), not just shared words.
    for a, nbrs in concept_adj.items():
        if a not in topic_set:
            continue
        for b, weight in nbrs.items():
            if b in topic_set and a < b:  # a<b dedups the undirected pair
                add_edge(topic_set[a], topic_set[b], concept=True, weight=weight)

    # Auto-link similar topics (word-overlap ≥ 0.3) — no self-links
    topic_nodes_list = [n for n in nodes if n["type"] == "topic"]
    for i in range(len(topic_nodes_list)):
        wi = set(topic_nodes_list[i]["label"].lower().split())
        for j in range(i + 1, len(topic_nodes_list)):
            wj = set(topic_nodes_list[j]["label"].lower().split())
            inter = len(wi & wj)
            union = len(wi | wj)
            if union > 0 and inter / union >= 0.3:
                add_edge(topic_nodes_list[i]["id"], topic_nodes_list[j]["id"], similarity=True)

    return {"nodes": nodes, "edges": edges}


# ── Knowledge Graph Modifications ─────────────────────────────────────────────

class GraphNodeAction(BaseModel):
    node_id: str

class GraphEdgeAction(BaseModel):
    source: str
    target: str

@app.post("/api/knowledge-graph/hide-node")
def graph_hide_node(body: GraphNodeAction, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("INSERT OR IGNORE INTO graph_hidden_nodes (user_id, node_id) VALUES (?, ?)",
               (user_id, body.node_id))
    # Also hide all edges to/from this node
    db.execute("DELETE FROM graph_custom_edges WHERE user_id = ? AND (source = ? OR target = ?)",
               (user_id, body.node_id, body.node_id))
    db.commit(); db.close()
    return {"ok": True}

@app.post("/api/knowledge-graph/restore-node")
def graph_restore_node(body: GraphNodeAction, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("DELETE FROM graph_hidden_nodes WHERE user_id = ? AND node_id = ?",
               (user_id, body.node_id))
    db.commit(); db.close()
    return {"ok": True}

@app.post("/api/knowledge-graph/hide-edge")
def graph_hide_edge(body: GraphEdgeAction, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("INSERT OR IGNORE INTO graph_hidden_edges (user_id, source, target) VALUES (?, ?, ?)",
               (user_id, body.source, body.target))
    # Also remove if it was a custom edge
    db.execute("DELETE FROM graph_custom_edges WHERE user_id = ? AND source = ? AND target = ?",
               (user_id, body.source, body.target))
    db.commit(); db.close()
    return {"ok": True}

@app.post("/api/knowledge-graph/add-edge")
def graph_add_edge(body: GraphEdgeAction, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("INSERT OR IGNORE INTO graph_custom_edges (user_id, source, target) VALUES (?, ?, ?)",
               (user_id, body.source, body.target))
    # Remove from hidden if it was hidden before
    db.execute("DELETE FROM graph_hidden_edges WHERE user_id = ? AND source = ? AND target = ?",
               (user_id, body.source, body.target))
    db.commit(); db.close()
    return {"ok": True}

@app.post("/api/knowledge-graph/reset")
def graph_reset(user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("DELETE FROM graph_hidden_nodes WHERE user_id = ?", (user_id,))
    db.execute("DELETE FROM graph_hidden_edges WHERE user_id = ?", (user_id,))
    db.execute("DELETE FROM graph_custom_edges WHERE user_id = ?", (user_id,))
    db.commit(); db.close()
    return {"ok": True}


# ── Exam dates & study plan ───────────────────────────────────────────────────

class ExamIn(BaseModel):
    subject: str
    exam_date: str
    notes: Optional[str] = ""


@app.post("/api/exam-dates")
def add_exam(e: ExamIn, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("INSERT INTO exam_dates (user_id, subject, exam_date, notes) VALUES (?,?,?,?)", (user_id, e.subject, e.exam_date, e.notes))
    db.commit()
    db.close()
    return {"ok": True}


@app.get("/api/exam-dates")
def list_exams(user_id: int = Depends(get_current_user)):
    db = get_db()
    rows = db.execute("SELECT * FROM exam_dates WHERE user_id = ? ORDER BY exam_date", (user_id,)).fetchall()
    db.close()
    return [dict(r) for r in rows]


@app.delete("/api/exam-dates/{eid}")
def delete_exam(eid: int, user_id: int = Depends(get_current_user)):
    db = get_db()
    db.execute("DELETE FROM exam_dates WHERE id = ? AND user_id = ?", (eid, user_id))
    db.commit()
    db.close()
    return {"ok": True}


@app.post("/api/generate/studyplan")
async def gen_study_plan(request: Request, user_id: int = Depends(get_current_user)):
    body = await request.json()
    exam_date = body.get("exam_date", "")
    subject = body.get("subject", "Medicine")

    db = get_db()
    weak = [dict(r) for r in db.execute(
        "SELECT topic, CAST(SUM(is_correct) AS REAL)/COUNT(*) as acc FROM quiz_attempts WHERE user_id = ? GROUP BY topic ORDER BY acc ASC LIMIT 5",
        (user_id,)
    ).fetchall()]
    mats = [r["original_name"] for r in db.execute(
        """SELECT m.original_name FROM materials m
           JOIN user_materials um ON um.material_id = m.id WHERE um.user_id = ?""",
        (user_id,)
    ).fetchall()]
    db.close()

    days_left = 30
    if exam_date:
        try:
            days_left = max(1, (date.fromisoformat(exam_date) - date.today()).days)
        except Exception:
            pass

    weak_str = ", ".join([f"{r['topic']} ({round(r['acc']*100)}%)" for r in weak]) or "No data yet — start quizzing!"
    mats_str = ", ".join(mats) or "No materials uploaded yet"

    resp = get_client().messages.create(
        model=MODEL, max_tokens=2500,
        messages=[{"role": "user", "content": f"""Create a personalised medical study plan.

Subject: {subject}
Days until exam: {days_left}
Materials: {mats_str}
Weak topics: {weak_str}

Generate a plan for {min(days_left, 14)} days. Return ONLY JSON:
{{
  "overview": "Brief 2-sentence strategy",
  "daily_hours": 3,
  "days": [{{
    "day": 1,
    "focus": "Main topic",
    "tasks": ["Task 1", "Task 2", "Task 3"],
    "priority": "high"
  }}]
}}

Prioritise weak topics. Mix flashcard review, active recall, and new content."""}]
    )

    try:
        return parse_json_response(resp.content[0].text)
    except Exception:
        return {"overview": "Could not generate plan", "daily_hours": 3, "days": []}


# ── Mind maps ────────────────────────────────────────────────────────────────

@app.post("/api/generate/mindmap/{mid}")
def generate_mindmap(mid: int, force: bool = False, user_id: int = Depends(get_current_user)):
    db = get_db()
    if not user_can_access(db, mid, user_id):
        db.close()
        raise HTTPException(403, "No access to this material")
    mat = db.execute("SELECT * FROM materials WHERE id = ?", (mid,)).fetchone()
    if not mat:
        raise HTTPException(404, "Material not found")

    # Memory: reuse the existing mind map without an AI call unless forced
    if not force:
        row = db.execute("SELECT data FROM mind_maps WHERE material_id = ? AND user_id = ?", (mid, user_id)).fetchone()
        if row:
            db.close()
            try:
                return json.loads(row["data"])
            except Exception:
                pass
            db = get_db()  # reopen if the stored data was unparseable

    instructions = """Create a mind map from the SOURCE STUDY MATERIAL above.

Return ONLY valid JSON matching this exact structure — no markdown, no extra keys:
{
  "id": "root",
  "label": "Short topic name",
  "children": [
    {
      "id": "b1",
      "label": "Branch name",
      "children": [
        {"id": "l1_1", "label": "Leaf label", "children": []},
        {"id": "l1_2", "label": "Leaf label", "children": []}
      ]
    }
  ]
}

Rules:
- 4–6 main branches
- 2–4 leaves per branch
- Root label: 1–3 words (the topic name)
- Branch labels: 2–4 words (a key concept or category)
- Leaf labels: 2–5 words (a specific fact, term, or sub-concept)
- Never use colons or slashes inside labels — split into two separate leaves instead
- All ids must be unique strings"""

    text = generate_json(mat, instructions, model=HAIKU, max_tokens=3000)
    try:
        data = parse_json_response(text)
    except Exception:
        data = {"id": "root", "label": mat["original_name"], "children": [{"id": "e1", "label": "Error — try again", "children": []}]}

    db.execute("DELETE FROM mind_maps WHERE material_id = ? AND user_id = ?", (mid, user_id))
    db.execute("INSERT INTO mind_maps (material_id, user_id, title, data) VALUES (?,?,?,?)", (mid, user_id, mat["original_name"], json.dumps(data)))
    db.commit()
    db.close()
    return data


@app.get("/api/mindmaps")
def get_mindmaps(material_id: Optional[int] = None, user_id: int = Depends(get_current_user)):
    db = get_db()
    base = "SELECT mm.*, m.original_name FROM mind_maps mm JOIN materials m ON mm.material_id = m.id"
    rows = db.execute(base + " WHERE mm.material_id = ? AND mm.user_id = ?", (material_id, user_id)).fetchall() \
        if material_id else db.execute(base + " WHERE mm.user_id = ?", (user_id,)).fetchall()
    db.close()
    result = []
    for r in rows:
        d = dict(r)
        try:
            d["data"] = json.loads(d["data"])
        except Exception:
            pass
        result.append(d)
    return result


# ── Network info ─────────────────────────────────────────────────────────────

@app.get("/api/network-info")
def network_info():
    """Returns LAN IPs and the current Cloudflare tunnel URL (read from the tunnel log)."""
    import socket, re
    ips: list[str] = []
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        ip = s.getsockname()[0]
        s.close()
        if ip and ip not in ("127.0.0.1", "0.0.0.0"):
            ips.append(ip)
    except Exception:
        pass
    try:
        hostname = socket.gethostname()
        for info in socket.getaddrinfo(hostname, None):
            ip = info[4][0]
            if ":" in ip:
                continue
            if ip.startswith(("192.168.", "10.", "172.")) and ip not in ips:
                ips.append(ip)
    except Exception:
        pass

    # Read the current Cloudflare tunnel URL from the log file
    tunnel_url: Optional[str] = None
    log_path = Path("data/tunnel.log")
    if log_path.exists():
        try:
            text = log_path.read_text(errors="ignore")
            # Find the last trycloudflare.com URL in the log
            matches = re.findall(r'https://[a-z0-9-]+\.trycloudflare\.com', text)
            if matches:
                tunnel_url = matches[-1]
        except Exception:
            pass

    return {"ips": ips, "port": 8000, "tunnel_url": tunnel_url}


# ── Serve material images ─────────────────────────────────────────────────────
# Images are stored in IMAGES_DIR (either static/material_images locally, or
# /data/material_images on Railway). All new image URLs use /images/{filename}.
# Old /static/material_images/ URLs still work locally via the StaticFiles mount.

@app.get("/images/{filename}")
def serve_image(filename: str):
    """Serve a material image from IMAGES_DIR (works locally and on Railway)."""
    # Reject any path-traversal attempts — only allow a bare filename inside IMAGES_DIR
    if "/" in filename or "\\" in filename or ".." in filename:
        raise HTTPException(404, "Image not found")
    p = (IMAGES_DIR / filename).resolve()
    if IMAGES_DIR.resolve() not in p.parents or not p.is_file():
        raise HTTPException(404, "Image not found")
    return FileResponse(str(p))


# ── Serve static ─────────────────────────────────────────────────────────────

app.mount("/static", StaticFiles(directory="static"), name="static")


def _print_startup():
    import socket
    print("\n  MedVault — Your Medical Study Vault")
    print("  " + "─" * 42)
    print("  Local   → http://localhost:8000")
    try:
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        s.connect(("8.8.8.8", 80))
        ip = s.getsockname()[0]
        s.close()
        if ip and ip not in ("127.0.0.1", "0.0.0.0"):
            print(f"  Network → http://{ip}:8000  (share with classmates)")
    except Exception:
        pass
    print("  Press Ctrl+C to stop\n")


if __name__ == "__main__":
    import uvicorn
    _print_startup()
    uvicorn.run(app, host="0.0.0.0", port=8000)
